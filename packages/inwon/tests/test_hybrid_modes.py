from __future__ import annotations

import platform
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from canva_ppt_mcp.com_editor import (
    ComUnavailableError, _compact_text, _design_audit, _font_size_change,
    _issue_signature, _validate_watch_delay, design_fingerprint, edit_template_with_com,
)
from canva_ppt_mcp.pipeline import create_presentation
from canva_ppt_mcp.routing import (
    confirm_presentation_mode,
    consume_mode_confirmation,
    infer_presentation_mode,
    list_user_templates,
    prepare_presentation_task,
    resolve_template,
)


def test_intent_recommendation_never_executes():
    assert infer_presentation_mode("IBM 소개 PPT 만들어줘") == "generate"
    assert infer_presentation_mode("이 디자인 그대로 내용만 IBM에 맞게 수정해줘") == "template_com"
    assert infer_presentation_mode("PPT 좀 부탁해") is None


def test_confirmation_token_is_required_and_one_time(tmp_path: Path):
    result = prepare_presentation_task("IBM 소개 PPT 만들어줘", str(tmp_path))
    assert result["requires_user_confirmation"] is True
    assert result["execution_blocked"] is True
    with pytest.raises(RuntimeError, match="MODE_CONFIRMATION_REQUIRED"):
        consume_mode_confirmation(requested_mode="generate", execution_token=None)
    confirmed = confirm_presentation_mode(result["confirmation_id"], "generate")
    consume_mode_confirmation(requested_mode="generate", execution_token=confirmed["execution_token"])
    with pytest.raises(RuntimeError, match="already-used"):
        consume_mode_confirmation(requested_mode="generate", execution_token=confirmed["execution_token"])


def test_wrong_mode_token_is_rejected(tmp_path: Path):
    result = prepare_presentation_task("템플릿 내용만 수정해줘", str(tmp_path))
    confirmed = confirm_presentation_mode(result["confirmation_id"], "template_com")
    with pytest.raises(RuntimeError, match="confirmed mode"):
        consume_mode_confirmation(requested_mode="generate", execution_token=confirmed["execution_token"])


def test_template_folder_is_the_only_source(tmp_path: Path):
    root = tmp_path / "template"; root.mkdir()
    prs = Presentation(); prs.save(root / "sample.pptx")
    assert list_user_templates(str(root))[0]["name"] == "sample.pptx"
    assert resolve_template("sample.pptx", str(root)) == (root / "sample.pptx").resolve()
    with pytest.raises(FileNotFoundError):
        resolve_template("../outside.pptx", str(root))


def test_design_fingerprint_ignores_text_content(tmp_path: Path):
    source = tmp_path / "a.pptx"
    prs = Presentation(); slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    run = box.text_frame.paragraphs[0].add_run(); run.text = "Old topic"; run.font.size = Pt(24)
    prs.save(source)
    before = design_fingerprint(source)
    prs = Presentation(source); prs.slides[0].shapes[0].text_frame.paragraphs[0].runs[0].text = "New topic"
    prs.save(source)
    after = design_fingerprint(source)
    assert before == after


def _minimal_plan():
    return {
        "communication_job": "test",
        "design_system": {
            "palette": {
                "primary": "#102A43", "secondary": ["#D9EAF7"], "accent": "#00B8A9",
                "background_light": "#F7FBFF", "background_dark": "#071A2B"
            },
            "typography": {},
            "visual_motif": "test",
            "style_preset": "orbital",
            "layout_rotation": ["title", "two_column", "icon_rows", "closing"]
        },
        "slides": [
            {"title": "A", "layout": "title"},
            {"title": "B", "layout": "two_column", "items": [{"heading": "H", "body": "B"}]},
            {"title": "C", "layout": "closing"}
        ],
        "grounded": False,
        "language": "ko"
    }


def test_python_pptx_pipeline_refuses_template_path(tmp_path: Path):
    template = tmp_path / "template.pptx"; Presentation().save(template)
    with pytest.raises(ValueError, match="no longer accepted"):
        create_presentation(
            topic="test", output_path=str(tmp_path / "out.pptx"), template_path=str(template),
            content_json=_minimal_plan(), research_required=False, max_qa_rounds=1,
        )


def test_com_mode_fails_explicitly_without_windows(tmp_path: Path):
    if platform.system() == "Windows":
        pytest.skip("non-Windows behavior only")
    root = tmp_path / "template"; root.mkdir()
    prs = Presentation()
    for _ in range(3): prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(root / "sample.pptx")
    with pytest.raises(ComUnavailableError, match="Windows"):
        edit_template_with_com(
            topic="test", template_name="sample.pptx", output_path=str(tmp_path / "out.pptx"),
            template_dir=str(root), research_required=False,
        )


def test_com_watch_delay_is_never_headless_fast():
    assert _validate_watch_delay(0.55) == 0.55
    assert _validate_watch_delay(5) == 5.0
    with pytest.raises(ValueError, match="0.20"):
        _validate_watch_delay(0.0)
    with pytest.raises(ValueError, match="5.0"):
        _validate_watch_delay(5.1)


def test_com_public_api_uses_watch_mode_not_visible_flag():
    import inspect
    signature = inspect.signature(edit_template_with_com)
    assert "visible" not in signature.parameters
    assert signature.parameters["step_delay"].default == 0.55



def test_design_audit_allows_small_font_shrink_as_warning(tmp_path: Path):
    source = tmp_path / "before.pptx"
    prs = Presentation(); slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    run = box.text_frame.paragraphs[0].add_run(); run.text = "Original"; run.font.size = Pt(24)
    prs.save(source)
    before = design_fingerprint(source)
    prs = Presentation(source)
    prs.slides[0].shapes[0].text_frame.paragraphs[0].runs[0].font.size = Pt(22)
    prs.save(source)
    audit = _design_audit(before, design_fingerprint(source))
    assert audit["errors"] == []
    assert any("font size reduced" in warning for warning in audit["warnings"])


def test_design_audit_rejects_large_font_shrink(tmp_path: Path):
    source = tmp_path / "before.pptx"
    prs = Presentation(); slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    run = box.text_frame.paragraphs[0].add_run(); run.text = "Original"; run.font.size = Pt(24)
    prs.save(source)
    before = design_fingerprint(source)
    prs = Presentation(source)
    prs.slides[0].shapes[0].text_frame.paragraphs[0].runs[0].font.size = Pt(16)
    prs.save(source)
    audit = _design_audit(before, design_fingerprint(source))
    assert any("font size changed too much" in error for error in audit["errors"])


def test_design_audit_rejects_geometry_change(tmp_path: Path):
    source = tmp_path / "before.pptx"
    prs = Presentation(); slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    box.text = "Original"
    prs.save(source)
    before = design_fingerprint(source)
    prs = Presentation(source)
    prs.slides[0].shapes[0].left += Inches(0.1)
    prs.save(source)
    audit = _design_audit(before, design_fingerprint(source))
    assert any("shape geometry changed" in error for error in audit["errors"])


def test_font_size_tolerance_contract():
    assert _font_size_change(24, 23.9) == (None, None)
    assert _font_size_change(24, 22)[0] == "warning"
    assert _font_size_change(24, 18)[0] == "error"


def test_local_compaction_never_expands_or_invents():
    original = "데이터 플랫폼을 기반으로 분석 준비 시간을 줄이고 재현성을 높일 수 있습니다"
    compact = _compact_text(original, 28)
    assert len(compact) <= 28
    assert len(compact) <= len(original)
    assert "재현성" in original  # source fact remains the only source of content


def test_com_editor_does_not_force_autosize_shrink_to_fit():
    import inspect
    import canva_ppt_mcp.com_editor as module
    source = inspect.getsource(module._set_text_preserve_style) + inspect.getsource(module._replace_text)
    assert '"AutoSize", 2' not in source
    assert "AutoSize=2" not in source



def test_design_audit_tolerates_theme_to_rgb_normalization():
    shape = {
        "path": "1:Box", "shape_type": 17, "geometry": [1, 2, 3, 4], "rotation": 0.0,
        "placeholder": None, "fill": ("theme", 4, 0.0), "line": None,
        "text_style": {"name": None, "size": None, "bold": None, "italic": None, "color": ("theme", 1, 0.0)},
    }
    changed = dict(shape)
    changed["fill"] = ("rgb", "#123456", 0.0)
    changed["text_style"] = {"name": "Aptos", "size": None, "bold": None, "italic": None, "color": ("rgb", "#000000", 0.0)}
    before = {"slide_count": 1, "slide_size": [1, 1], "slides": [{"slide": 1, "background": None, "shapes": [shape]}]}
    after = {"slide_count": 1, "slide_size": [1, 1], "slides": [{"slide": 1, "background": None, "shapes": [changed]}]}
    assert _design_audit(before, after)["errors"] == []


def test_design_audit_still_rejects_explicit_fill_color_change():
    base = {
        "path": "1:Box", "shape_type": 17, "geometry": [1, 2, 3, 4], "rotation": 0.0,
        "placeholder": None, "fill": ("rgb", "#111111", 0.0), "line": None,
        "text_style": {"name": None, "size": None, "bold": None, "italic": None, "color": None},
    }
    changed = dict(base); changed["fill"] = ("rgb", "#222222", 0.0)
    before = {"slide_count": 1, "slide_size": [1, 1], "slides": [{"slide": 1, "background": None, "shapes": [base]}]}
    after = {"slide_count": 1, "slide_size": [1, 1], "slides": [{"slide": 1, "background": None, "shapes": [changed]}]}
    assert any("fill color changed" in error for error in _design_audit(before, after)["errors"])



def test_deferred_qa_issue_signature_is_stable_and_order_independent():
    a = [
        {"slide": 2, "shape": "Body 1", "code": "TEXT_OVERFLOW", "message": "too tall"},
        {"slide": 1, "shape": "Title 1", "code": "EDIT_FAILED", "message": "busy"},
    ]
    b = list(reversed(a))
    assert _issue_signature(a) == _issue_signature(b)


def test_com_public_api_defaults_to_bounded_post_qa():
    import inspect
    signature = inspect.signature(edit_template_with_com)
    assert signature.parameters["max_post_qa_rounds"].default == 2


def test_com_editor_uses_complete_then_validate_not_mid_edit_abort():
    import inspect
    import canva_ppt_mcp.com_editor as module
    first = inspect.getsource(module._replace_text_first_pass)
    session = inspect.getsource(module._edit_with_com)
    public = inspect.getsource(module.edit_template_with_com)
    assert "defer all fit problems until post-QA" in first
    assert "PASS 1" in session
    assert "_validate_com_operations" in session
    assert "cycle_breaker_triggered" in session
    assert "automatic_restart_blocked" in public
    assert "COM edit changed template design; output rejected" not in public


def test_final_qa_failure_is_manifest_state_not_automatic_restart():
    import inspect
    import canva_ppt_mcp.com_editor as module
    public = inspect.getsource(module.edit_template_with_com)
    assert '"completed_with_unresolved_issues"' in public
    assert '"automatic_full_restart": False' in public
    assert '"retry_policy": "targeted_only_bounded"' in public


def test_bob_safe_server_boundary_returns_json_instead_of_tool_exception():
    server_path = Path(__file__).resolve().parents[1] / "src" / "canva_ppt_mcp" / "server.py"
    source = server_path.read_text(encoding="utf-8")
    assert "_normal_tool_interruption" in source
    assert 'except Exception as exc' in source
    assert '"mode_gate"' in source
    assert '"template_com"' in source
    assert '"do_not_retry"' in source
    assert '"mcp_transport_error"' in source


def test_semantic_qa_is_deferred_until_after_com_edit():
    import inspect
    import canva_ppt_mcp.com_editor as module
    source = inspect.getsource(module.edit_template_with_com)
    assert "semantic QA is deferred" in source
    assert "Semantic QA failed" not in source
    assert "SEMANTIC_" in source


def test_com_session_collects_slide_save_and_qa_runtime_errors():
    import inspect
    import canva_ppt_mcp.com_editor as module
    source = inspect.getsource(module._edit_with_com)
    assert '"slide_access_error"' in source
    assert '"save_issues"' in source
    assert '"qa_runtime_issues"' in source
    assert '"SAVE_ERROR"' in source
    assert '"QA_RUNTIME_ERROR"' in source
    assert '"COM_SESSION_INTERRUPTED"' in source
