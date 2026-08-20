from __future__ import annotations

import json
import os
import platform
import re
import shutil
import time
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Iterable

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

from .planner import create_plan
from .routing import resolve_template
from .semantic_qa import inspect_grounding

EMU = 914400
RPC_E_CALL_REJECTED = -2147418111
RPC_E_SERVERCALL_RETRYLATER = -2147417846
PROTECTED_PLACEHOLDERS = {
    PP_PLACEHOLDER.DATE,
    PP_PLACEHOLDER.FOOTER,
    PP_PLACEHOLDER.SLIDE_NUMBER,
}
PROTECTED_NAME_TOKENS = (
    "footer", "slide number", "slidenumber", "date", "copyright", "logo", "watermark",
)


@dataclass
class TextSlot:
    shape: Any
    name: str
    left: float
    top: float
    width: float
    height: float
    font_size: float
    placeholder_type: int | None
    original_text: str


class ComUnavailableError(RuntimeError):
    pass


def _ensure_com_environment() -> None:
    if platform.system() != "Windows":
        raise ComUnavailableError("COM template editing requires Windows and Microsoft PowerPoint")
    try:
        import pythoncom  # noqa: F401
        import pywintypes  # noqa: F401
        import win32com.client  # noqa: F401
    except ImportError as exc:
        raise ComUnavailableError("Install the Windows extra: pip install -e '.[windows]'") from exc


def _color_signature(color) -> tuple[str, str | int, float] | None:
    """Return a stable color token while tolerating PowerPoint theme normalization."""
    try:
        brightness = round(float(getattr(color, "brightness", 0.0) or 0.0), 4)
    except Exception:
        brightness = 0.0
    try:
        if color.type and color.rgb:
            return ("rgb", f"#{color.rgb}", brightness)
    except (AttributeError, TypeError, ValueError):
        pass
    try:
        theme = color.theme_color
        if theme is not None:
            return ("theme", int(theme), brightness)
    except (AttributeError, TypeError, ValueError):
        pass
    return None


def _hex_color(color) -> str | None:
    signature = _color_signature(color)
    return signature[1] if signature and signature[0] == "rgb" else None


def _walk_shapes_pptx(shapes, prefix: str = ""):
    for index, shape in enumerate(shapes, 1):
        path = f"{prefix}/{index}:{shape.name}" if prefix else f"{index}:{shape.name}"
        yield path, shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _walk_shapes_pptx(shape.shapes, path)


def _first_run_style(shape) -> dict[str, Any]:
    if not getattr(shape, "has_text_frame", False):
        return {"name": None, "size": None, "bold": None, "italic": None, "color": None}
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            return {
                "name": run.font.name,
                "size": round(run.font.size.pt, 2) if run.font.size else None,
                "bold": run.font.bold,
                "italic": run.font.italic,
                "color": _color_signature(run.font.color),
            }
    return {"name": None, "size": None, "bold": None, "italic": None, "color": None}


def _shape_signature(path: str, shape) -> dict[str, Any]:
    placeholder = None
    if getattr(shape, "is_placeholder", False):
        try:
            placeholder = int(shape.placeholder_format.type)
        except Exception:
            placeholder = None
    try:
        fill = _color_signature(shape.fill.fore_color)
    except Exception:
        fill = None
    try:
        line = _color_signature(shape.line.color)
    except Exception:
        line = None
    return {
        "path": path,
        "shape_type": int(shape.shape_type),
        "geometry": [int(shape.left), int(shape.top), int(shape.width), int(shape.height)],
        "rotation": round(float(shape.rotation or 0), 2),
        "placeholder": placeholder,
        "fill": fill,
        "line": line,
        "text_style": _first_run_style(shape),
    }


def design_fingerprint(path: str | Path) -> dict:
    """Fingerprint visual structure while deliberately ignoring text content.

    PowerPoint can normalize theme fonts/colors and OOXML runs when COM saves a
    file.  The fingerprint therefore records semantic style tokens instead of
    requiring byte-identical internal representation.
    """
    prs = Presentation(str(path))
    slides = []
    for slide_index, slide in enumerate(prs.slides, 1):
        signatures = [_shape_signature(shape_path, shape) for shape_path, shape in _walk_shapes_pptx(slide.shapes)]
        try:
            bg = _color_signature(slide.background.fill.fore_color)
        except Exception:
            bg = None
        slides.append({"slide": slide_index, "background": bg, "shapes": signatures})
    return {
        "slide_count": len(prs.slides),
        "slide_size": [int(prs.slide_width), int(prs.slide_height)],
        "slides": slides,
    }


def _color_change(before, after) -> tuple[str | None, str | None]:
    """Return (severity, reason). Theme<->RGB normalization is not a hard failure."""
    if before == after:
        return None, None
    if before is None or after is None:
        return None, "implicit/explicit color normalization"
    if before[0] == after[0]:
        return ("error", f"color changed: {before} -> {after}")
    # PowerPoint often materializes a theme color as an explicit RGB during save.
    return None, "theme/RGB color normalization"


def _font_size_change(before: float | None, after: float | None) -> tuple[str | None, str | None]:
    if before is None or after is None or before <= 0 or after <= 0:
        return None, None
    delta = after - before
    if abs(delta) <= 0.25:
        return None, None
    shrink_ratio = max(0.0, (before - after) / before)
    if delta > 0:
        # Growing text unexpectedly is a style change, not an overflow repair.
        return "error", f"font size increased: {before:g}pt -> {after:g}pt"
    if shrink_ratio <= 0.125 and before - after <= 4.0:
        return "warning", f"font size reduced to resolve overflow: {before:g}pt -> {after:g}pt"
    return "error", f"font size changed too much: {before:g}pt -> {after:g}pt"


def _design_audit(before: dict, after: dict) -> dict[str, list[str]]:
    """Separate real visual changes from benign PowerPoint text normalization."""
    errors: list[str] = []
    warnings: list[str] = []
    if before["slide_count"] != after["slide_count"]:
        errors.append(f"slide_count changed: {before['slide_count']} -> {after['slide_count']}")
    if before["slide_size"] != after["slide_size"]:
        errors.append("slide size changed")
    for left, right in zip(before["slides"], after["slides"]):
        slide_no = left["slide"]
        sev, reason = _color_change(left["background"], right["background"])
        if sev == "error":
            errors.append(f"slide {slide_no}: background {reason}")
        if len(left["shapes"]) != len(right["shapes"]):
            errors.append(f"slide {slide_no}: shape count changed")
            continue
        for a, b in zip(left["shapes"], right["shapes"]):
            path = a["path"]
            if path != b["path"] or a["shape_type"] != b["shape_type"]:
                errors.append(f"slide {slide_no}: shape identity/type changed at {path}")
                continue
            if any(abs(x - y) > 100 for x, y in zip(a["geometry"], b["geometry"])):
                errors.append(f"slide {slide_no}: shape geometry changed at {path}")
            if abs(a["rotation"] - b["rotation"]) > 0.01:
                errors.append(f"slide {slide_no}: rotation changed at {path}")
            if a["placeholder"] != b["placeholder"]:
                errors.append(f"slide {slide_no}: placeholder type changed at {path}")
            for field in ("fill", "line"):
                sev, reason = _color_change(a[field], b[field])
                if sev == "error":
                    errors.append(f"slide {slide_no}: {field} {reason} at {path}")
            sa, sb = a["text_style"], b["text_style"]
            # None -> explicit is usually theme/run materialization on save.
            if sa["name"] and sb["name"] and sa["name"].casefold() != sb["name"].casefold():
                errors.append(f"slide {slide_no}: font changed at {path}: {sa['name']} -> {sb['name']}")
            for field in ("bold", "italic"):
                if sa[field] is not None and sb[field] is not None and sa[field] != sb[field]:
                    errors.append(f"slide {slide_no}: {field} changed at {path}")
            sev, reason = _color_change(sa["color"], sb["color"])
            if sev == "error":
                errors.append(f"slide {slide_no}: text {reason} at {path}")
            sev, reason = _font_size_change(sa["size"], sb["size"])
            if sev == "error":
                errors.append(f"slide {slide_no}: {reason} at {path}")
            elif sev == "warning":
                warnings.append(f"slide {slide_no}: {reason} at {path}")
            if len(errors) >= 20:
                return {"errors": errors, "warnings": warnings}
    return {"errors": errors, "warnings": warnings}


def _design_diff(before: dict, after: dict) -> list[str]:
    """Backward-compatible hard-error list used by older callers/tests."""
    return _design_audit(before, after)["errors"]

def _validate_template_source(template: Path, output: Path) -> None:
    if template.resolve() == output.resolve():
        raise ValueError("template_path and output_path must be different; the original template is never overwritten")
    if output.suffix.lower() not in {".pptx", ".pptm"}:
        raise ValueError("COM edited output must end in .pptx or .pptm")
    output.parent.mkdir(parents=True, exist_ok=True)


def _com_retry(callable_, *args, attempts: int = 5, delay: float = 1.0, **kwargs):
    try:
        import pywintypes  # type: ignore
    except ImportError as exc:
        raise ComUnavailableError("pywin32 is required for COM template editing") from exc
    last = None
    for attempt in range(attempts):
        try:
            return callable_(*args, **kwargs)
        except pywintypes.com_error as exc:  # type: ignore[attr-defined]
            last = exc
            hresult = getattr(exc, "hresult", None)
            if hresult not in {RPC_E_CALL_REJECTED, RPC_E_SERVERCALL_RETRYLATER} or attempt == attempts - 1:
                raise
            time.sleep(delay)
    raise last  # pragma: no cover


def _connect_powerpoint():
    _ensure_com_environment()
    try:
        import pythoncom  # type: ignore
        import win32com.client  # type: ignore
    except ImportError as exc:
        raise ComUnavailableError("Install the Windows extra: pip install -e '.[windows]'") from exc
    pythoncom.CoInitialize()
    launched = False
    try:
        app = win32com.client.GetActiveObject("PowerPoint.Application")
    except Exception:
        app = win32com.client.DispatchEx("PowerPoint.Application")
        launched = True
    # Template COM edits are intentionally never headless. The user asked to
    # watch every PowerPoint modification happen in the real application.
    _com_retry(setattr, app, "Visible", True)
    try:
        _com_retry(app.Activate)
    except Exception:
        pass
    return pythoncom, app, launched


def _iter_com_shapes(shapes) -> Iterable[Any]:
    # msoGroup = 6
    for i in range(1, int(shapes.Count) + 1):
        shape = shapes.Item(i)
        if int(shape.Type) == 6:
            yield from _iter_com_shapes(shape.GroupItems)
        else:
            yield shape


def _placeholder_type(shape) -> int | None:
    # msoPlaceholder = 14
    if int(shape.Type) != 14:
        return None
    try:
        return int(shape.PlaceholderFormat.Type)
    except Exception:
        return None


def _safe_text_slot(shape) -> TextSlot | None:
    try:
        if int(shape.HasTextFrame) != -1:
            return None
    except Exception:
        return None
    name = str(shape.Name or "")
    lowered = name.lower()
    if any(token in lowered for token in PROTECTED_NAME_TOKENS):
        return None
    placeholder = _placeholder_type(shape)
    # COM numeric constants align with PowerPoint placeholder enum:
    # 13 slide number, 14 header, 15 footer, 16 date.
    if placeholder in {13, 14, 15, 16}:
        return None
    try:
        text = str(shape.TextFrame.TextRange.Text or "").strip()
    except Exception:
        text = ""
    # Preserve tiny numeric/ornamental tokens such as page indexes or "01".
    if text and len(text) <= 5 and re.fullmatch(r"[\d\s./|—–-]+", text):
        return None
    try:
        font_size = float(shape.TextFrame.TextRange.Font.Size)
        if font_size < 0:
            font_size = 0.0
    except Exception:
        font_size = 0.0
    return TextSlot(
        shape=shape,
        name=name,
        left=float(shape.Left), top=float(shape.Top),
        width=float(shape.Width), height=float(shape.Height),
        font_size=font_size,
        placeholder_type=placeholder,
        original_text=text,
    )


def _collect_slots(slide) -> tuple[TextSlot | None, list[TextSlot]]:
    slots = [slot for shape in _iter_com_shapes(slide.Shapes) if (slot := _safe_text_slot(shape))]
    if not slots:
        return None, []
    title = next((slot for slot in slots if slot.placeholder_type in {1, 3, 5}), None)
    if title is None:
        # Favor a large, high text element. Font size dominates; top position breaks ties.
        title = max(slots, key=lambda s: (s.font_size, -s.top, s.width))
    bodies = [slot for slot in slots if slot is not title]
    bodies.sort(key=lambda s: (s.top, s.left))
    return title, bodies


def _content_units(spec) -> list[str]:
    values: list[str] = []
    if spec.subtitle.strip():
        values.append(spec.subtitle.strip())
    for item in spec.items:
        heading = item.heading.strip()
        body = item.body.strip()
        if heading and body:
            values.append(f"{heading}\r{body}")
        elif heading or body:
            values.append(heading or body)
    return values


@dataclass
class _ComTextStyle:
    name: str | None = None
    size: float | None = None
    bold: int | None = None
    italic: int | None = None
    color_rgb: int | None = None
    autosize: int | None = None


def _capture_com_text_style(slot: TextSlot) -> _ComTextStyle:
    style = _ComTextStyle()
    try:
        font = slot.shape.TextFrame.TextRange.Font
        name = str(font.Name or "").strip()
        style.name = name or None
    except Exception:
        pass
    try:
        value = float(slot.shape.TextFrame.TextRange.Font.Size)
        style.size = value if value > 0 else None
    except Exception:
        pass
    for attr in ("Bold", "Italic"):
        try:
            value = int(getattr(slot.shape.TextFrame.TextRange.Font, attr))
            if value in {-1, 0}:
                setattr(style, attr.lower(), value)
        except Exception:
            pass
    try:
        style.color_rgb = int(slot.shape.TextFrame.TextRange.Font.Color.RGB)
    except Exception:
        pass
    try:
        style.autosize = int(slot.shape.TextFrame2.AutoSize)
    except Exception:
        pass
    return style


def _restore_com_text_style(slot: TextSlot, style: _ComTextStyle, *, size: float | None = None) -> None:
    """Reapply the template's visible text style after PowerPoint rewrites runs."""
    try:
        font = slot.shape.TextFrame.TextRange.Font
    except Exception:
        font = None
    if font is not None:
        # Mixed-style ranges can reject one property while accepting another,
        # so restore each attribute independently instead of failing the block.
        properties = []
        if style.name:
            properties.append((font, "Name", style.name))
        target_size = size if size is not None else style.size
        if target_size and target_size > 0:
            properties.append((font, "Size", float(target_size)))
        if style.bold in {-1, 0}:
            properties.append((font, "Bold", style.bold))
        if style.italic in {-1, 0}:
            properties.append((font, "Italic", style.italic))
        for target, attr, value in properties:
            try:
                _com_retry(setattr, target, attr, value)
            except Exception:
                pass
        if style.color_rgb is not None:
            try:
                _com_retry(setattr, font.Color, "RGB", style.color_rgb)
            except Exception:
                pass
    # Preserve the template's AutoSize mode. Never force shrink-to-fit.
    if style.autosize is not None:
        try:
            _com_retry(setattr, slot.shape.TextFrame2, "AutoSize", style.autosize)
        except Exception:
            pass


def _set_text_preserve_style(slot: TextSlot, text: str, style: _ComTextStyle, *, size: float | None = None) -> None:
    text_range = slot.shape.TextFrame.TextRange
    _com_retry(setattr, text_range, "Text", text)
    _restore_com_text_style(slot, style, size=size)


def _text_bounds(slot: TextSlot) -> dict[str, float | bool] | None:
    """Measure rendered text through PowerPoint without resizing the shape."""
    try:
        tf2 = slot.shape.TextFrame2
        tr2 = tf2.TextRange
        available_w = max(1.0, float(slot.shape.Width) - float(tf2.MarginLeft) - float(tf2.MarginRight))
        available_h = max(1.0, float(slot.shape.Height) - float(tf2.MarginTop) - float(tf2.MarginBottom))
        bound_w = float(tr2.BoundWidth)
        bound_h = float(tr2.BoundHeight)
        try:
            word_wrap = int(tf2.WordWrap) != 0
        except Exception:
            word_wrap = True
        return {
            "available_width": available_w,
            "available_height": available_h,
            "bound_width": bound_w,
            "bound_height": bound_h,
            "word_wrap": word_wrap,
        }
    except Exception:
        return None


def _is_text_overflow(slot: TextSlot) -> tuple[bool, dict[str, float | bool] | None]:
    metrics = _text_bounds(slot)
    if not metrics:
        return False, None
    height_over = float(metrics["bound_height"]) > float(metrics["available_height"]) * 1.03
    width_over = (not bool(metrics["word_wrap"]) and
                  float(metrics["bound_width"]) > float(metrics["available_width"]) * 1.03)
    return bool(height_over or width_over), metrics


def _compact_text(text: str, budget: int) -> str:
    """Grounded local compression: only removes words/clauses; never invents facts."""
    budget = max(8, int(budget))
    clean = re.sub(r"[ \t]+", " ", text).strip()
    if len(clean) <= budget:
        return clean
    replacements = (
        ("할 수 있습니다", "가능합니다"), ("할 수 있다", "가능하다"),
        ("및", "·"), ("그리고", "·"), ("대한", ""),
        ("통하여", "통해"), ("위하여", "위해"), ("기반으로", "기반"),
    )
    candidate = clean
    for source, target in replacements:
        candidate = candidate.replace(source, target)
    candidate = re.sub(r"\s+", " ", candidate).strip()
    if len(candidate) <= budget:
        return candidate

    # Keep a heading line when the planner supplied "heading\rbody".
    parts = re.split(r"\r\n|\r|\n", candidate, maxsplit=1)
    if len(parts) == 2:
        head = parts[0].strip()
        body_budget = max(6, budget - min(len(head), max(6, budget // 3)) - 1)
        body = _compact_text(parts[1], body_budget)
        if len(head) > budget // 3:
            head = _compact_text(head, max(6, budget // 3))
        joined = f"{head}\r{body}".strip()
        if len(joined) <= budget:
            return joined
        candidate = joined

    # Prefer complete clauses/sentences before word-boundary shortening.
    clauses = [x.strip() for x in re.split(r"(?<=[.!?다])\s+|[;；]", candidate) if x.strip()]
    if clauses:
        kept: list[str] = []
        used = 0
        for clause in clauses:
            cost = len(clause) + (1 if kept else 0)
            if kept and used + cost > budget:
                break
            if not kept and cost > budget:
                break
            kept.append(clause); used += cost
        if kept:
            return " ".join(kept)

    words = candidate.split()
    kept: list[str] = []
    for word in words:
        proposed = " ".join(kept + [word])
        if len(proposed) > budget:
            break
        kept.append(word)
    if kept:
        return " ".join(kept).rstrip(" ,.;:·")
    return candidate[:budget].rstrip(" ,.;:·")


def _ai_compact_text(text: str, budget: int, language: str) -> str | None:
    """Optional semantic rewrite when the OpenAI extra/key is already configured."""
    if not os.getenv("OPENAI_API_KEY"):
        return None
    try:
        from openai import OpenAI
        client = OpenAI()
        response = client.responses.create(
            model=os.getenv("PPT_MCP_MODEL", "gpt-5.6"),
            input=[
                {"role": "system", "content": (
                    "Shorten presentation text to fit an existing PowerPoint text box. "
                    "Use only facts, numbers, names and meaning already present in the input. "
                    "Do not add claims. Preserve the input language and important numeric values. "
                    f"Return only the rewritten text, at most {budget} characters."
                )},
                {"role": "user", "content": text},
            ],
        )
        value = str(response.output_text or "").strip()
        if value and len(value) <= budget:
            return value
    except Exception:
        return None
    return None


def _overflow_budget(text: str, metrics: dict[str, float | bool] | None) -> int:
    if not metrics:
        return max(8, int(len(text) * 0.82))
    ratios = [1.0]
    if float(metrics["bound_height"]) > 0:
        ratios.append(float(metrics["available_height"]) / float(metrics["bound_height"]))
    if not bool(metrics["word_wrap"]) and float(metrics["bound_width"]) > 0:
        ratios.append(float(metrics["available_width"]) / float(metrics["bound_width"]))
    ratio = max(0.35, min(ratios))
    return max(8, int(len(text) * ratio * 0.90))


def _replace_text_first_pass(slot: TextSlot, text: str, language: str = "ko") -> dict[str, Any]:
    """First pass: write the requested text and defer all fit problems until post-QA.

    This deliberately does *not* compact text, shrink fonts, or raise on overflow.
    The goal is to finish the whole deck first, save it, then validate globally.
    """
    style = _capture_com_text_style(slot)
    original_size = style.size or (slot.font_size if slot.font_size > 0 else None)
    _set_text_preserve_style(slot, text, style)
    overflow, metrics = _is_text_overflow(slot)
    return {
        "requested_text": text,
        "actual_text": text,
        "rewrite_method": "none",
        "font_size_before": original_size,
        "font_size_after": original_size,
        "font_shrunk": False,
        "overflow_resolved": not overflow,
        "deferred_overflow": bool(overflow),
        "metrics": metrics,
        "status": "deferred_overflow" if overflow else "applied",
    }


def _repair_text_after_validation(
    slot: TextSlot,
    text: str,
    language: str = "ko",
    *,
    repair_round: int = 1,
) -> dict[str, Any]:
    """Post-QA repair only. Never raises for a content-fit problem.

    Round 1 uses the normal measured budget. Round 2 makes the rewrite more
    compact, but typography is still protected by the same 12.5% / 4pt limit.
    """
    style = _capture_com_text_style(slot)
    original_size = style.size or (slot.font_size if slot.font_size > 0 else None)
    overflow, metrics = _is_text_overflow(slot)
    actual_text = str(slot.shape.TextFrame.TextRange.Text or "")
    rewrite_method = "none"

    if overflow and text.strip():
        budget = _overflow_budget(text, metrics)
        if repair_round >= 2:
            budget = max(8, int(budget * 0.78))
        local = _compact_text(text, budget)
        compacted = _ai_compact_text(text, budget, language) or local
        if compacted and compacted != actual_text:
            _set_text_preserve_style(slot, compacted, style)
            actual_text = compacted
            rewrite_method = "ai" if compacted != local else "local"
            overflow, metrics = _is_text_overflow(slot)

    final_size = original_size
    font_shrunk = False
    if overflow and original_size:
        floor = max(14.0, original_size * 0.875, original_size - 4.0)
        size = original_size
        while overflow and size - 0.5 >= floor - 0.01:
            size = round(size - 0.5, 2)
            _restore_com_text_style(slot, style, size=size)
            overflow, metrics = _is_text_overflow(slot)
        final_size = size
        font_shrunk = final_size < original_size - 0.01

    return {
        "requested_text": text,
        "actual_text": actual_text,
        "rewrite_method": rewrite_method,
        "font_size_before": original_size,
        "font_size_after": final_size,
        "font_shrunk": font_shrunk,
        "overflow_resolved": not overflow,
        "deferred_overflow": bool(overflow),
        "metrics": metrics,
        "status": "repaired" if not overflow else "unresolved_overflow",
    }


def _replace_text(slot: TextSlot, text: str, language: str = "ko") -> dict[str, Any]:
    """Backward-compatible repair helper; the first-pass editor does not call this."""
    return _repair_text_after_validation(slot, text, language, repair_round=1)


def _issue_signature(issues: list[dict[str, Any]]) -> tuple[tuple[Any, ...], ...]:
    """Stable signature used to stop repair cycles when nothing changes."""
    return tuple(sorted(
        (
            issue.get("slide"),
            issue.get("shape"),
            issue.get("code"),
            str(issue.get("message", ""))[:160],
        )
        for issue in issues
    ))


def _find_slot_by_name(slide, shape_name: str) -> TextSlot | None:
    for shape in _iter_com_shapes(slide.Shapes):
        slot = _safe_text_slot(shape)
        if slot is not None and slot.name == shape_name:
            return slot
    return None


def _validate_com_operations(pres, operations: list[dict[str, Any]]) -> list[dict[str, Any]]:
    """Validate only after the full first pass (and after each bounded repair pass)."""
    issues: list[dict[str, Any]] = []
    for slide_op in operations:
        slide_no = int(slide_op["slide"])
        slide = pres.Slides(slide_no)
        for op in slide_op.get("operations", []):
            shape_name = str(op.get("shape") or "")
            if op.get("status") == "edit_error":
                issues.append({
                    "slide": slide_no,
                    "shape": shape_name,
                    "role": op.get("role"),
                    "code": "EDIT_FAILED",
                    "message": str(op.get("error") or "COM text replacement failed"),
                    "requested_text": op.get("requested_text", ""),
                })
                continue
            slot = _find_slot_by_name(slide, shape_name)
            if slot is None:
                issues.append({
                    "slide": slide_no,
                    "shape": shape_name,
                    "role": op.get("role"),
                    "code": "TEXT_SLOT_MISSING",
                    "message": "edited text slot could not be found during post-QA",
                    "requested_text": op.get("requested_text", ""),
                })
                continue
            overflow, metrics = _is_text_overflow(slot)
            if overflow:
                issues.append({
                    "slide": slide_no,
                    "shape": shape_name,
                    "role": op.get("role"),
                    "code": "TEXT_OVERFLOW",
                    "message": "text does not fit the original template box after the full edit pass",
                    "requested_text": op.get("requested_text", ""),
                    "metrics": metrics,
                })
    return issues


def _update_operation(operations: list[dict[str, Any]], slide_no: int, shape_name: str, result: dict[str, Any]) -> None:
    for slide_op in operations:
        if int(slide_op.get("slide", -1)) != slide_no:
            continue
        for op in slide_op.get("operations", []):
            if op.get("shape") == shape_name:
                role = op.get("role")
                requested = op.get("requested_text", result.get("requested_text", ""))
                op.clear()
                op.update({"shape": shape_name, "role": role, "requested_text": requested, **result})
                return


def _repair_post_validation_issues(
    app,
    pres,
    operations: list[dict[str, Any]],
    issues: list[dict[str, Any]],
    step_delay: float,
    language: str,
    repair_round: int,
) -> list[dict[str, Any]]:
    """Repair only failed shapes after the entire presentation has been created."""
    repair_log: list[dict[str, Any]] = []
    for issue in issues:
        slide_no = int(issue["slide"])
        shape_name = str(issue.get("shape") or "")
        slide = pres.Slides(slide_no)
        slot = _find_slot_by_name(slide, shape_name)
        if slot is None:
            repair_log.append({**issue, "repair_status": "slot_missing"})
            continue
        requested = str(issue.get("requested_text") or "")
        _activate_edit_target(app, pres, slide_no, slot, step_delay)
        try:
            if issue.get("code") == "EDIT_FAILED":
                result = _replace_text_first_pass(slot, requested, language)
            elif issue.get("code") == "TEXT_OVERFLOW":
                result = _repair_text_after_validation(
                    slot, requested, language, repair_round=repair_round
                )
            else:
                repair_log.append({**issue, "repair_status": "not_auto_repairable"})
                continue
            try:
                _com_retry(slot.shape.Select)
            except Exception:
                pass
            _ui_pause(step_delay)
            _update_operation(operations, slide_no, shape_name, result)
            repair_log.append({
                "slide": slide_no,
                "shape": shape_name,
                "code": issue.get("code"),
                "repair_status": result.get("status"),
                "repair_round": repair_round,
                "actual_text": result.get("actual_text"),
                "font_size_after": result.get("font_size_after"),
            })
        except Exception as exc:
            # A repair failure is also deferred. Do not abort the deck or trigger
            # a full restart loop; record it for the final report.
            repair_log.append({
                "slide": slide_no,
                "shape": shape_name,
                "code": issue.get("code"),
                "repair_status": "repair_error",
                "repair_round": repair_round,
                "error": f"{type(exc).__name__}: {exc}",
            })
    return repair_log


def _validate_watch_delay(step_delay: float) -> float:
    value = float(step_delay)
    if not 0.20 <= value <= 5.0:
        raise ValueError("step_delay must be between 0.20 and 5.0 seconds so every COM edit remains visible")
    return value


def _ui_pause(seconds: float) -> None:
    """Give PowerPoint time to repaint selections before/after each visible edit."""
    try:
        import pythoncom  # type: ignore
        pythoncom.PumpWaitingMessages()
    except Exception:
        pass
    time.sleep(max(0.0, seconds))
    try:
        import pythoncom  # type: ignore
        pythoncom.PumpWaitingMessages()
    except Exception:
        pass


def _activate_edit_target(app, pres, slide_index: int, slot: TextSlot | None, step_delay: float) -> None:
    """Bring the exact slide/shape about to be edited into the foreground."""
    _com_retry(setattr, app, "Visible", True)
    try:
        _com_retry(app.Activate)
    except Exception:
        pass
    try:
        window = pres.Windows(1)
        _com_retry(window.Activate)
        # ppWindowMaximized = 3
        try:
            _com_retry(setattr, window, "WindowState", 3)
        except Exception:
            pass
        _com_retry(window.View.GotoSlide, slide_index)
    except Exception:
        # Presentation is still the explicitly opened target. If a particular
        # window API is unavailable, do not switch to ActivePresentation.
        pass
    _ui_pause(step_delay)
    if slot is not None:
        try:
            _com_retry(slot.shape.Select)
        except Exception:
            pass
        _ui_pause(max(0.20, step_delay * 0.65))


def _visible_replace_first_pass(
    app,
    pres,
    slide_index: int,
    slot: TextSlot,
    text: str,
    step_delay: float,
    language: str,
) -> dict[str, Any]:
    _activate_edit_target(app, pres, slide_index, slot, step_delay)
    result = _replace_text_first_pass(slot, text, language)
    try:
        _com_retry(slot.shape.Select)
    except Exception:
        pass
    _ui_pause(step_delay)
    return result


def _edit_with_com(
    output: Path,
    plan,
    step_delay: float,
    max_post_qa_rounds: int,
) -> tuple[list[dict[str, Any]], dict[str, Any], Any, Any]:
    """Finish as much of the visible COM edit as possible before any fatal reporting.

    Once the target presentation is open, slide/shape/save/QA problems become
    collected issues instead of exceptions. This is intentionally friendly to
    MCP clients such as Bob, which may otherwise retry a failed tool call from
    the beginning and repeat the same COM failure forever.
    """
    pythoncom = None
    app = None
    pres = None
    launched = False
    previous_alerts = None
    operations: list[dict[str, Any]] = []
    validation: dict[str, Any] = {
        "strategy": "complete_then_validate",
        "first_pass_issues": [],
        "repair_rounds": [],
        "cycle_breaker_triggered": False,
        "max_post_qa_rounds": max_post_qa_rounds,
        "final_issues": [],
        "session_issues": [],
        "save_issues": [],
        "qa_runtime_issues": [],
        "first_pass_completed": False,
    }

    def record_issue(bucket: str, code: str, message: str, *, slide=None, shape=None) -> dict[str, Any]:
        issue = {"slide": slide, "shape": shape, "code": code, "message": message}
        validation.setdefault(bucket, []).append(issue)
        return issue

    def safe_save(stage: str) -> bool:
        if pres is None:
            return False
        try:
            _com_retry(pres.Save)
            return True
        except Exception as exc:
            record_issue(
                "save_issues", "SAVE_ERROR",
                f"{stage}: {type(exc).__name__}: {exc}",
            )
            return False

    try:
        pythoncom, app, launched = _connect_powerpoint()
        try:
            previous_alerts = int(app.DisplayAlerts)
            _com_retry(setattr, app, "DisplayAlerts", 1)  # ppAlertsNone
        except Exception:
            previous_alerts = None

        # Opening the explicit target is the only point where there is no useful
        # partial deck to preserve yet, so an open failure is allowed to reach the
        # server safety boundary as a structured normal JSON response.
        pres = _com_retry(app.Presentations.Open, str(output), 0, 0, -1)
        opened = Path(str(pres.FullName)).resolve()
        if opened != output.resolve():
            raise RuntimeError(f"PowerPoint opened an unexpected target: {opened}")
        if int(pres.Slides.Count) != len(plan.slides):
            raise RuntimeError(
                f"template slide count changed before edit: COM={pres.Slides.Count}, plan={len(plan.slides)}"
            )

        # PASS 1: process every slide. A slide-level COM failure is data, not a stop.
        for index, spec in enumerate(plan.slides, 1):
            slide_ops: list[dict[str, Any]] = []
            try:
                try:
                    _activate_edit_target(app, pres, index, None, step_delay)
                except Exception as exc:
                    slide_ops.append({
                        "shape": None, "role": "slide", "requested_text": None,
                        "actual_text": None, "status": "ui_activation_error",
                        "error": f"{type(exc).__name__}: {exc}",
                    })
                slide = pres.Slides(index)
                title, bodies = _collect_slots(slide)
            except Exception as exc:
                slide_ops.append({
                    "shape": None, "role": "slide", "requested_text": None,
                    "actual_text": None, "status": "slide_access_error",
                    "error": f"{type(exc).__name__}: {exc}",
                })
                operations.append({
                    "slide": index, "title_shape": None, "editable_body_slots": 0,
                    "content_units": 0, "operations": slide_ops,
                })
                continue

            if title is not None:
                try:
                    edit_result = _visible_replace_first_pass(
                        app, pres, index, title, spec.title, step_delay, plan.language
                    )
                    slide_ops.append({"shape": title.name, "role": "title", **edit_result})
                except Exception as exc:
                    slide_ops.append({
                        "shape": title.name, "role": "title", "requested_text": spec.title,
                        "actual_text": None, "status": "edit_error",
                        "error": f"{type(exc).__name__}: {exc}",
                    })

            units = _content_units(spec)
            for body_index, slot in enumerate(bodies):
                replacement = units[body_index] if body_index < len(units) else ""
                try:
                    edit_result = _visible_replace_first_pass(
                        app, pres, index, slot, replacement, step_delay, plan.language
                    )
                    slide_ops.append({"shape": slot.name, "role": "body", **edit_result})
                except Exception as exc:
                    slide_ops.append({
                        "shape": slot.name, "role": "body", "requested_text": replacement,
                        "actual_text": None, "status": "edit_error",
                        "error": f"{type(exc).__name__}: {exc}",
                    })
            operations.append({
                "slide": index,
                "title_shape": title.name if title else None,
                "editable_body_slots": len(bodies),
                "content_units": len(units),
                "operations": slide_ops,
            })

        validation["first_pass_completed"] = True
        validation["first_pass_saved"] = safe_save("first_pass")

        # POST-QA only after the entire first pass.
        try:
            issues = _validate_com_operations(pres, operations)
        except Exception as exc:
            issues = [record_issue(
                "qa_runtime_issues", "QA_RUNTIME_ERROR",
                f"first post-QA: {type(exc).__name__}: {exc}",
            )]
        validation["first_pass_issues"] = list(issues)
        seen_signatures: set[tuple[tuple[Any, ...], ...]] = set()

        for repair_round in range(1, max_post_qa_rounds + 1):
            if not issues:
                break
            signature = _issue_signature(issues)
            if signature in seen_signatures:
                validation["cycle_breaker_triggered"] = True
                break
            seen_signatures.add(signature)
            try:
                repair_log = _repair_post_validation_issues(
                    app, pres, operations, issues, step_delay, plan.language, repair_round
                )
            except Exception as exc:
                repair_log = [{
                    "slide": None, "shape": None, "code": "QA_REPAIR_RUNTIME_ERROR",
                    "repair_status": "repair_runtime_error", "repair_round": repair_round,
                    "error": f"{type(exc).__name__}: {exc}",
                }]
                record_issue(
                    "qa_runtime_issues", "QA_REPAIR_RUNTIME_ERROR",
                    f"repair round {repair_round}: {type(exc).__name__}: {exc}",
                )
            safe_save(f"repair_round_{repair_round}")
            try:
                next_issues = _validate_com_operations(pres, operations)
            except Exception as exc:
                next_issues = [record_issue(
                    "qa_runtime_issues", "QA_RUNTIME_ERROR",
                    f"post repair round {repair_round}: {type(exc).__name__}: {exc}",
                )]
            validation["repair_rounds"].append({
                "round": repair_round,
                "input_issue_count": len(issues),
                "repairs": repair_log,
                "remaining_issue_count": len(next_issues),
                "remaining_issues": next_issues,
            })
            if _issue_signature(next_issues) == signature and next_issues:
                validation["cycle_breaker_triggered"] = True
                issues = next_issues
                break
            issues = next_issues

        # Save once more after all repair attempts, but never turn a save failure
        # into an MCP exception/restart loop.
        validation["final_saved"] = safe_save("final")
        combined = list(issues)
        for issue in validation.get("save_issues", []):
            if issue not in combined:
                combined.append(issue)
        for issue in validation.get("qa_runtime_issues", []):
            if issue not in combined:
                combined.append(issue)
        validation["final_issues"] = combined
        validation["content_passed"] = not combined
        try:
            _activate_edit_target(app, pres, len(plan.slides), None, step_delay)
        except Exception:
            pass
        return operations, validation, app, pres

    except Exception as exc:
        # If a presentation was already opened, preserve the partial visible result
        # and return an interrupted state instead of raising. If open/setup itself
        # failed, there is no useful partial document; the server safety boundary
        # will still convert the exception into a normal JSON response.
        if pres is None:
            raise
        record_issue(
            "session_issues", "COM_SESSION_INTERRUPTED",
            f"{type(exc).__name__}: {exc}",
        )
        safe_save("session_interrupted")
        validation["final_issues"] = list(validation.get("session_issues", [])) + list(validation.get("save_issues", []))
        validation["content_passed"] = False
        validation["session_interrupted"] = True
        return operations, validation, app, pres
    finally:
        if app is not None and previous_alerts is not None:
            try:
                _com_retry(setattr, app, "DisplayAlerts", previous_alerts)
            except Exception:
                pass
        if pythoncom is not None:
            try:
                pythoncom.CoUninitialize()
            except Exception:
                pass

def edit_template_with_com(*, topic: str, template_name: str, output_path: str,
                           audience: str = "", purpose: str = "", language: str = "ko",
                           research_text: str | None = None, source_urls: list[str] | None = None,
                           research_required: bool = True,
                           research_documents: list[dict[str, str]] | None = None,
                           template_dir: str | None = None,
                           step_delay: float = 0.55,
                           max_post_qa_rounds: int = 2) -> dict[str, Any]:
    """Visibly edit first, then validate the completed deck and repair only bad slides/shapes.

    Content/design QA failures are returned in the manifest instead of raising at
    the end. This prevents MCP clients from automatically restarting the exact
    same template job forever. Fatal setup/open/save failures can still raise.
    """
    if not topic or not topic.strip():
        raise ValueError("topic must not be empty")
    _ensure_com_environment()
    step_delay = _validate_watch_delay(step_delay)
    if not 0 <= int(max_post_qa_rounds) <= 3:
        raise ValueError("max_post_qa_rounds must be between 0 and 3")
    max_post_qa_rounds = int(max_post_qa_rounds)
    template = resolve_template(template_name, template_dir)
    output = Path(output_path).expanduser().resolve()
    _validate_template_source(template, output)

    manifest_dir = output.with_name(output.stem + "_qa")
    manifest_dir.mkdir(parents=True, exist_ok=True)

    before_source = template
    if template.suffix.lower() in {".pptx", ".pptm"}:
        shutil.copy2(template, output)
    else:
        if platform.system() != "Windows":
            raise ComUnavailableError(".potx/.potm COM template editing requires Windows PowerPoint")
        pythoncom = app = pres = None
        launched = False
        try:
            pythoncom, app, launched = _connect_powerpoint()
            pres = _com_retry(app.Presentations.Open, str(template), 0, 1, -1)
            _ui_pause(step_delay)
            _com_retry(pres.SaveAs, str(output), 24)  # ppSaveAsOpenXMLPresentation
        finally:
            if pres is not None:
                try: _com_retry(pres.Close)
                except Exception: pass
            if app is not None and launched:
                try: _com_retry(app.Quit)
                except Exception: pass
            if pythoncom is not None:
                try: pythoncom.CoUninitialize()
                except Exception: pass
        before_source = output

    # Keep a pristine working baseline for diagnostics/recovery, but do not
    # automatically restart the whole deck from it; bounded targeted repair is safer.
    baseline_copy = manifest_dir / ("template-baseline" + output.suffix.lower())
    shutil.copy2(output, baseline_copy)
    before = design_fingerprint(before_source)
    slide_count = before["slide_count"]
    if not 3 <= slide_count <= 30:
        raise ValueError("template must contain between 3 and 30 slides")

    plan = create_plan(
        topic=topic, audience=audience, purpose=purpose, slide_count=slide_count,
        language=language, content_json=None, research_text=research_text,
        source_urls=source_urls, research_required=research_required,
        style_preference=None, research_documents=research_documents,
    )
    semantic_issues = inspect_grounding(plan, topic, research_required)
    semantic_errors = [issue for issue in semantic_issues if issue.severity == "error"]
    # 3.4: semantic QA is deferred too. It must not abort the visible COM pass;
    # record it and let the completed deck return for review.

    operations, validation, _visible_app, _visible_presentation = _edit_with_com(
        output, plan, step_delay, max_post_qa_rounds
    )

    try:
        after = design_fingerprint(output)
        design_audit = _design_audit(before, after)
        design_issues = design_audit["errors"]
        design_warnings = design_audit["warnings"]
    except Exception as exc:
        # A fingerprint/OOXML read problem after COM editing is a final QA issue,
        # not a reason for Bob to restart the whole job.
        design_issues = [f"design audit could not complete: {type(exc).__name__}: {exc}"]
        design_warnings = []
    for slide_op in operations:
        for op in slide_op.get("operations", []):
            if op.get("font_shrunk"):
                warning = (
                    f"slide {slide_op['slide']}: overflow repair reduced font in {op['shape']}: "
                    f"{op.get('font_size_before')}pt -> {op.get('font_size_after')}pt"
                )
                if warning not in design_warnings:
                    design_warnings.append(warning)

    final_issues = list(validation.get("final_issues", []))
    final_issues.extend({
        "slide": None,
        "shape": None,
        "code": "DESIGN_CHANGE",
        "message": issue,
    } for issue in design_issues)
    final_issues.extend({
        "slide": issue.slide,
        "shape": issue.shape_name,
        "code": f"SEMANTIC_{issue.code}",
        "message": issue.message,
    } for issue in semantic_errors)

    passed = not final_issues
    if validation.get("session_interrupted"):
        completion_status = "interrupted_with_partial_result"
    elif passed and design_warnings:
        completion_status = "completed_with_warnings"
    elif passed:
        completion_status = "completed"
    else:
        completion_status = "completed_with_unresolved_issues"

    validation.update({
        "design_errors": design_issues,
        "design_warnings": design_warnings,
        "final_issues": final_issues,
        "passed": passed,
        "automatic_full_restart": False,
        "retry_policy": "targeted_only_bounded",
    })

    (manifest_dir / "template-com-plan.json").write_text(plan.model_dump_json(indent=2), encoding="utf-8")
    (manifest_dir / "template-com-validation.json").write_text(
        json.dumps(validation, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    manifest = {
        "mode": "template_com",
        "engine": "PowerPoint COM (pywin32)",
        "version": "3.4.0",
        "completion_status": completion_status,
        "passed": passed,
        "requires_manual_review": not passed,
        "automatic_restart_blocked": True,
        "do_not_retry": True,
        "tool_call_succeeded": True,
        "mcp_transport_error": False,
        "operation_completed": completion_status.startswith("completed"),
        "watch_mode": {
            "enabled": True,
            "powerpoint_visible": True,
            "step_delay_seconds": step_delay,
            "select_each_text_box": True,
            "keep_result_open": True,
            "suppress_modal_alerts_during_edit": True,
        },
        "deferred_qa": {
            "enabled": True,
            "edit_first_validate_after": True,
            "max_post_qa_rounds": max_post_qa_rounds,
            "cycle_breaker_triggered": validation.get("cycle_breaker_triggered", False),
            "full_restart_on_qa_failure": False,
            "bob_safe_no_tool_exception": True,
        },
        "template_path": str(template),
        "baseline_copy": str(baseline_copy),
        "output_path": str(output),
        "slide_count": slide_count,
        "design_preserved": not design_issues,
        "design_issues": design_issues,
        "design_warnings": design_warnings,
        "semantic_qa": {
            "passed": not semantic_errors,
            "issues": [issue.model_dump() for issue in semantic_issues],
        },
        "post_validation": validation,
        "operations": operations,
        "research_sources": [source.model_dump() for source in plan.research_sources],
        "plan_path": str(manifest_dir / "template-com-plan.json"),
        "validation_path": str(manifest_dir / "template-com-validation.json"),
    }
    (manifest_dir / "template-com-manifest.json").write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    return manifest
