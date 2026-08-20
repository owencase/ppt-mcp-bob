from __future__ import annotations

import json
import math
import os
import re
import shutil
import subprocess
import tempfile
from collections import Counter
from pathlib import Path

import fitz
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

from .models import QAIssue, QAReport


PLACEHOLDER_RE = re.compile(r"lorem ipsum|\bTODO\b|\[insert[^]]*\]|click to add|텍스트를 입력", re.I)
EMU = 914400


def _command(name: str) -> str:
    found = shutil.which(name)
    if not found:
        raise RuntimeError(f"Required renderer '{name}' was not found in PATH")
    return found


def render_presentation(pptx_path: str, output_dir: str, pages: set[int] | None = None) -> list[str]:
    out = Path(output_dir); out.mkdir(parents=True, exist_ok=True)
    if pages is None:
        for stale in out.glob("slide-*.png"):
            stale.unlink(missing_ok=True)
    pdf_dir = out / "pdf"; pdf_dir.mkdir(exist_ok=True)
    # An isolated writable LibreOffice profile avoids lock/permission failures
    # when the MCP client launches several jobs or runs in a sandbox.
    with tempfile.TemporaryDirectory(prefix="ppt-mcp-lo-") as temp:
        temp_path = Path(temp); profile = temp_path / "profile"; profile.mkdir()
        converted_dir = temp_path / "converted"; converted_dir.mkdir()
        bundled_fonts = Path(__file__).resolve().parent / "assets" / "fonts"
        fontconfig = temp_path / "fonts.conf"
        fontconfig.write_text(
            "<?xml version='1.0'?><!DOCTYPE fontconfig SYSTEM 'fonts.dtd'>"
            "<fontconfig><dir>/usr/share/fonts</dir>"
            f"<dir>{bundled_fonts}</dir><cachedir>{temp_path / 'font-cache'}</cachedir></fontconfig>",
            encoding="utf-8",
        )
        env = os.environ.copy(); env["FONTCONFIG_FILE"] = str(fontconfig)
        profile_uri = profile.resolve().as_uri()
        subprocess.run([_command("soffice"), f"-env:UserInstallation={profile_uri}", "--headless",
                        "--convert-to", "pdf", "--outdir", str(converted_dir), str(Path(pptx_path).resolve())],
                       check=True, capture_output=True, text=True, timeout=180, env=env)
        fresh_pdf = converted_dir / (Path(pptx_path).stem + ".pdf")
        if not fresh_pdf.exists(): raise RuntimeError("LibreOffice did not create a PDF")
        pdf = pdf_dir / fresh_pdf.name
        shutil.copy2(fresh_pdf, pdf)
    doc = fitz.open(pdf)
    rendered = []
    selected = pages or set(range(1, len(doc) + 1))
    for page_no in selected:
        page = doc[page_no - 1]
        pix = page.get_pixmap(matrix=fitz.Matrix(1.6, 1.6), alpha=False)
        target = out / f"slide-{page_no:02}.png"; pix.save(target)
        if not target.exists() or target.stat().st_size < 1024:
            raise RuntimeError(f"슬라이드 {page_no} 렌더 이미지가 비어 있거나 손상되었습니다: {target}")
        rendered.append(str(target.resolve()))
    return rendered


def _shape_box(shape):
    return tuple(v / EMU for v in (shape.left, shape.top, shape.width, shape.height))


def _walk_shapes(shapes, root=None):
    for shape in shapes:
        cluster = root or shape
        yield shape, cluster
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _walk_shapes(shape.shapes, cluster)


def _intersection(a, b):
    ax, ay, aw, ah = a; bx, by, bw, bh = b
    return max(0, min(ax + aw, bx + bw) - max(ax, bx)) * max(0, min(ay + ah, by + bh) - max(ay, by))


def _fill_hex(fill) -> str | None:
    try:
        if fill.type and fill.fore_color.rgb:
            return f"#{fill.fore_color.rgb}"
    except (AttributeError, TypeError, ValueError):
        pass
    return None


def _effective_background_info(slide, target) -> tuple[str | None, str]:
    tx, ty, tw, th = _shape_box(target); cx, cy = tx + tw / 2, ty + th / 2
    background = _fill_hex(slide.background.fill)
    source = "SLIDE_BACKGROUND"
    for shape in slide.shapes:
        if shape is target: break
        x, y, w, h = _shape_box(shape)
        if x <= cx <= x + w and y <= cy <= y + h:
            candidate = _fill_hex(getattr(shape, "fill", None))
            if candidate:
                background = candidate; source = shape.name
    return background, source


def _effective_background(slide, target) -> str | None:
    return _effective_background_info(slide, target)[0]


def _luminance(value: str) -> float:
    channels = [int(value[i:i + 2], 16) / 255 for i in (1, 3, 5)]
    linear = [c / 12.92 if c <= .04045 else ((c + .055) / 1.055) ** 2.4 for c in channels]
    return .2126 * linear[0] + .7152 * linear[1] + .0722 * linear[2]


def _contrast(a: str, b: str) -> float:
    l1, l2 = sorted((_luminance(a), _luminance(b)), reverse=True)
    return (l1 + .05) / (l2 + .05)


def _estimate_overflow(shape) -> bool:
    if not getattr(shape, "has_text_frame", False) or not shape.text.strip(): return False
    tf = shape.text_frame; width = max(.1, shape.width / EMU - .08); height = max(.1, shape.height / EMU)
    max_size = 10.0
    for p in tf.paragraphs:
        for r in p.runs:
            if r.font.size: max_size = max(max_size, r.font.size.pt)
    chars_per_line = max(4, width * 72 / (max_size * .53))
    lines = sum(max(1, math.ceil(sum(1.65 if ord(c) > 127 else 1 for c in line) / chars_per_line))
                for line in shape.text.splitlines() or [""])
    # PowerPoint's line box is close to 1.0 em for single-line labels. Add the
    # explicit text-frame top/bottom inset instead of a blanket 1.2 multiplier.
    required = lines * max_size / 72 * 1.12 + .04
    return required > height * 1.06


def _title_wrap_risk(shape) -> bool:
    if shape.name != "TITLE" or not getattr(shape, "has_text_frame", False):
        return False
    width = shape.width / EMU
    limit = 33 if width >= 11 else 22
    display_width = sum(1.65 if ord(char) > 127 else 1 for char in shape.text.strip())
    return display_width > limit or "\n" in shape.text


def inspect_presentation(pptx_path: str) -> list[QAIssue]:
    prs = Presentation(pptx_path); issues = []
    sw, sh = prs.slide_width / EMU, prs.slide_height / EMU
    for si, slide in enumerate(prs.slides, 1):
        text_shapes = []; has_visual = False
        flattened = list(_walk_shapes(slide.shapes))
        for shape, _cluster in flattened:
            x, y, w, h = _shape_box(shape)
            if x < -0.02 or y < -0.02 or x + w > sw + .02 or y + h > sh + .02:
                if not shape.name.startswith("DECOR_"):
                    issues.append(QAIssue(slide=si, code="OUT_OF_BOUNDS", message="요소가 슬라이드 경계를 벗어났습니다.", shape_name=shape.name))
            if shape.shape_type in {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.CHART, MSO_SHAPE_TYPE.GROUP, MSO_SHAPE_TYPE.AUTO_SHAPE} and not shape.name.startswith("TEXT"):
                has_visual = True
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                text_shapes.append(shape)
                if PLACEHOLDER_RE.search(shape.text):
                    issues.append(QAIssue(slide=si, code="PLACEHOLDER", message="placeholder 텍스트가 남아 있습니다.", shape_name=shape.name))
                if _estimate_overflow(shape):
                    issues.append(QAIssue(slide=si, code="TEXT_OVERFLOW", message="텍스트가 상자 높이를 넘을 가능성이 큽니다.", shape_name=shape.name))
                if _title_wrap_risk(shape):
                    issues.append(QAIssue(slide=si, code="TITLE_WRAP", message="슬라이드 제목이 한 줄 너비를 넘습니다.", shape_name=shape.name))
                bg, bg_source = _effective_background_info(slide, shape)
                if bg:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            try: fg = f"#{run.font.color.rgb}" if run.font.color.rgb else None
                            except (AttributeError, ValueError): fg = None
                            if fg:
                                threshold = 3.0 if run.font.size and run.font.size.pt >= 24 else 4.5
                                if _contrast(fg, bg) < threshold:
                                    code = "DECOR_TEXT_CONTRAST" if bg_source.startswith("DECOR_") else "LOW_CONTRAST"
                                    message = ("장식 요소가 텍스트 대비를 해칩니다." if code == "DECOR_TEXT_CONTRAST"
                                               else "텍스트와 배경의 대비가 부족합니다.")
                                    issues.append(QAIssue(slide=si, code=code, message=message, shape_name=shape.name))
                                    break
        seen_empty_clusters = set()
        for nested, cluster in _walk_shapes(slide.shapes):
            if getattr(nested, "is_placeholder", False) and getattr(nested, "has_text_frame", False) and not nested.text.strip():
                if id(cluster) not in seen_empty_clusters:
                    issues.append(QAIssue(slide=si, code="EMPTY_PLACEHOLDER",
                                          message="비어 있는 구조적 placeholder/슬롯이 남아 있습니다.", shape_name=cluster.name))
                    seen_empty_clusters.add(id(cluster))
        if not has_visual:
            issues.append(QAIssue(slide=si, code="NO_VISUAL", message="시각 요소가 없습니다."))
        for i, a in enumerate(text_shapes):
            if a.name.startswith("DECOR_"):
                continue
            for b in text_shapes[i + 1:]:
                if b.name.startswith("DECOR_"):
                    continue
                area = _intersection(_shape_box(a), _shape_box(b))
                smaller = min(a.width * a.height, b.width * b.height) / (EMU * EMU)
                if smaller and area / smaller > .08:
                    issues.append(QAIssue(slide=si, code="TEXT_OVERLAP", message="텍스트 요소가 겹칩니다.", shape_name=f"{a.name} / {b.name}"))
        # A visual placed above text is a likely accidental cover. Background
        # panels created before text remain valid and are intentionally ignored.
        top_level = list(slide.shapes)
        for text_index, text_shape in enumerate(top_level):
            if not getattr(text_shape, "has_text_frame", False) or not text_shape.text.strip():
                continue
            text_box = _shape_box(text_shape)
            text_area = max(.001, text_box[2] * text_box[3])
            for visual in top_level[text_index + 1:]:
                if visual.name.startswith("DECOR_"):
                    continue
                if visual.shape_type not in {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.CHART,
                                             MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.LINE}:
                    continue
                if _intersection(text_box, _shape_box(visual)) / text_area > .08:
                    issues.append(QAIssue(
                        slide=si, code="VISUAL_OVER_TEXT",
                        message="텍스트 위에 뒤늦게 추가된 시각 요소가 겹칩니다.",
                        shape_name=f"{text_shape.name} / {visual.name}",
                    ))
    if len(prs.slides) >= 5:
        dark = 0
        for slide in prs.slides:
            color = _fill_hex(slide.background.fill)
            if color and _luminance(color) < .22:
                dark += 1
        ratio = dark / len(prs.slides)
        if not .40 <= ratio <= .50:
            issues.append(QAIssue(
                slide=1, code="DARK_RHYTHM", severity="warning",
                message=f"다크 슬라이드 비중이 {ratio:.0%}입니다. 권장 범위는 40~50%입니다.",
            ))
    return issues


def inspect_rendered_text(pptx_path: str, qa_dir: str) -> list[QAIssue]:
    """Catch glyph/font failures that structural PPTX checks cannot see."""
    prs = Presentation(pptx_path)
    pdf_path = Path(qa_dir, "pdf", Path(pptx_path).stem + ".pdf")
    if not pdf_path.exists():
        return [QAIssue(slide=1, code="RENDER_MISSING", message="PDF 렌더 결과가 없습니다.")]
    doc = fitz.open(pdf_path); issues = []
    for index, (slide, page) in enumerate(zip(prs.slides, doc), 1):
        expected = "".join(
            shape.text for shape, _cluster in _walk_shapes(slide.shapes)
            if getattr(shape, "has_text_frame", False) and shape.text.strip()
        )
        actual = page.get_text()
        expected_chars = [c for c in expected if c.isalnum()]
        actual_counts = Counter(c for c in actual if c.isalnum())
        if len(expected_chars) >= 8:
            remaining = actual_counts.copy(); matched = 0
            for char in expected_chars:
                if remaining[char] > 0:
                    matched += 1; remaining[char] -= 1
            visible = matched / len(expected_chars)
            if visible < .45:
                issues.append(QAIssue(
                    slide=index, code="RENDER_TEXT_LOSS",
                    message="렌더링 후 글자 상당수가 사라졌습니다. 해당 언어 글꼴 설치/대체를 확인하세요.",
                ))
    return issues


def _shorten(text: str, ratio: float = .82) -> str:
    if len(text) < 30: return text
    target = max(24, int(len(text) * ratio))
    cut = text[:target].rsplit(" ", 1)[0]
    return cut.rstrip(" ,.;:") + "…"


def repair_slides(pptx_path: str, issues: list[QAIssue]) -> set[int]:
    prs = Presentation(pptx_path); changed = set()
    by_slide: dict[int, list[QAIssue]] = {}
    for issue in issues: by_slide.setdefault(issue.slide, []).append(issue)
    for slide_no, slide_issues in by_slide.items():
        slide = prs.slides[slide_no - 1]
        for issue in slide_issues:
            if issue.code in {"TEXT_OVERFLOW", "PLACEHOLDER"}:
                for shape, _cluster in _walk_shapes(slide.shapes):
                    if shape.name != issue.shape_name or not getattr(shape, "has_text_frame", False): continue
                    repaired = False
                    if issue.code == "PLACEHOLDER":
                        shape.text_frame.clear(); repaired = True
                    else:
                        for p in shape.text_frame.paragraphs:
                            for r in p.runs:
                                if r.font.size:
                                    floor = 44 if shape.name == "TITLE" else (20 if "HEAD" in shape.name or "LABEL" in shape.name else 14)
                                    new_size = max(floor, r.font.size.pt * .94)
                                    if new_size < r.font.size.pt - .05:
                                        r.font.size = Pt(new_size); repaired = True
                    if repaired:
                        changed.add(slide_no)
            elif issue.code == "EMPTY_PLACEHOLDER":
                for shape in list(slide.shapes):
                    if shape.name == issue.shape_name:
                        shape.element.getparent().remove(shape.element); changed.add(slide_no); break
            elif issue.code == "OUT_OF_BOUNDS":
                for shape in slide.shapes:
                    if shape.name == issue.shape_name:
                        shape.left = max(0, shape.left); shape.top = max(0, shape.top)
                        shape.width = min(shape.width, prs.slide_width - shape.left)
                        shape.height = min(shape.height, prs.slide_height - shape.top); changed.add(slide_no)
            elif issue.code in {"LOW_CONTRAST", "DECOR_TEXT_CONTRAST"}:
                for shape, _cluster in _walk_shapes(slide.shapes):
                    if shape.name != issue.shape_name or not getattr(shape, "has_text_frame", False): continue
                    bg = _effective_background(slide, shape) or "#FFFFFF"
                    replacement = "FFFFFF" if _contrast("#FFFFFF", bg) >= _contrast("#102A43", bg) else "102A43"
                    for p in shape.text_frame.paragraphs:
                        for r in p.runs: r.font.color.rgb = RGBColor.from_string(replacement)
                    changed.add(slide_no)
    if changed: prs.save(pptx_path)
    return changed


def qa_loop(pptx_path: str, qa_dir: str, max_rounds: int = 3, auto_fix: bool = True) -> QAReport:
    rendered = render_presentation(pptx_path, qa_dir); rounds = 0; issues = []
    for rounds in range(1, max_rounds + 1):
        issues = inspect_presentation(pptx_path) + inspect_rendered_text(pptx_path, qa_dir)
        errors = [x for x in issues if x.severity == "error"]
        if not errors: break
        if not auto_fix: break
        changed = repair_slides(pptx_path, errors)
        if not changed: break
        rendered = list(set(rendered + render_presentation(pptx_path, qa_dir, changed)))
    # A repair performed in the last allowed round must be re-inspected. Without
    # this pass the report can contain stale pre-repair errors.
    if auto_fix and issues and any(x.severity == "error" for x in issues):
        issues = inspect_presentation(pptx_path) + inspect_rendered_text(pptx_path, qa_dir)
    report = QAReport(passed=not any(x.severity == "error" for x in issues), rounds=rounds, issues=issues, rendered_slides=sorted(rendered))
    Path(qa_dir).mkdir(parents=True, exist_ok=True)
    Path(qa_dir, "qa-report.json").write_text(report.model_dump_json(indent=2), encoding="utf-8")
    return report
