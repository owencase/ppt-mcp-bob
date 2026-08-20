from __future__ import annotations

import math
import ipaddress
import socket
import tempfile
from pathlib import Path
from urllib.parse import urlparse
from urllib.request import Request, urlopen

from PIL import Image, ImageDraw, ImageOps
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

from .models import ContentItem, DeckPlan, DesignSystem, SlideSpec


EMU = 914400
SW, SH = 13.333, 7.5
MARGIN = 0.62


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value.lstrip("#"))


def tint(value: str, amount: float = 0.88) -> str:
    c = value.lstrip("#")
    vals = [int(c[i:i + 2], 16) for i in (0, 2, 4)]
    vals = [round(v + (255 - v) * amount) for v in vals]
    return "#" + "".join(f"{v:02X}" for v in vals)


def mix(a: str, b: str, amount: float) -> str:
    aa, bb = a.lstrip("#"), b.lstrip("#")
    av = [int(aa[i:i + 2], 16) for i in (0, 2, 4)]
    bv = [int(bb[i:i + 2], 16) for i in (0, 2, 4)]
    vals = [round(x + (y - x) * amount) for x, y in zip(av, bv)]
    return "#" + "".join(f"{v:02X}" for v in vals)


def _set_transparency(shape, amount: float) -> None:
    """Set solid-fill transparency through DrawingML (0=opaque, 1=invisible)."""
    try:
        solid = shape.fill._xPr.solidFill
        color = solid[0]
        for child in list(color):
            if child.tag.endswith("}alpha"):
                color.remove(child)
        alpha = OxmlElement("a:alpha")
        alpha.set("val", str(round((1 - max(0, min(1, amount))) * 100000)))
        color.append(alpha)
    except (AttributeError, IndexError):
        return


def _add_text(slide, text, x, y, w, h, font, size, color, bold=False,
              align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, name="TEXT", wrap=True):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = name
    tf = shape.text_frame
    tf.clear(); tf.word_wrap = wrap; tf.margin_left = tf.margin_right = Inches(0.03)
    tf.margin_top = tf.margin_bottom = Inches(0.02); tf.vertical_anchor = valign
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    contains_korean = any("가" <= c <= "힣" for c in text)
    run.font.name = "Malgun Gothic" if contains_korean else font
    run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = rgb(color)
    # Keep the requested Latin typeface while giving PowerPoint an explicit
    # Korean/Asian fallback. This prevents missing glyphs on Windows clients.
    rpr = run._r.get_or_add_rPr()
    rpr.set("lang", "ko-KR")
    ea = rpr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}ea")
    if ea is None:
        ea = OxmlElement("a:ea"); rpr.append(ea)
    ea.set("typeface", "Malgun Gothic" if contains_korean else font)
    return shape


def _add_round_rect(slide, x, y, w, h, fill, radius=True, name="VISUAL_CARD", transparency=0.0):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = name; shape.fill.solid(); shape.fill.fore_color.rgb = rgb(fill); shape.line.fill.background()
    if transparency: _set_transparency(shape, transparency)
    return shape


def _add_circle(slide, x, y, size, fill, name="VISUAL_MOTIF", transparency=0.0):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
    shape.name = name; shape.fill.solid(); shape.fill.fore_color.rgb = rgb(fill); shape.line.fill.background()
    if transparency: _set_transparency(shape, transparency)
    return shape


def _gradient_png(path: Path, start: str, end: str, accent: str, textured: bool) -> None:
    width, height = 1600, 900
    start_rgb = tuple(int(start.lstrip("#")[i:i + 2], 16) for i in (0, 2, 4))
    end_rgb = tuple(int(end.lstrip("#")[i:i + 2], 16) for i in (0, 2, 4))
    image = Image.new("RGB", (width, height))
    draw = ImageDraw.Draw(image)
    for y in range(height):
        t = y / max(1, height - 1)
        # A diagonal-like two-stop gradient; rasterized before PPT insertion.
        row = tuple(round(a + (b - a) * t) for a, b in zip(start_rgb, end_rgb))
        draw.line((0, y, width, y), fill=row)
    # Add a soft off-canvas glow without relying on PowerPoint gradient fills.
    glow = Image.new("RGBA", image.size, (0, 0, 0, 0))
    gd = ImageDraw.Draw(glow)
    ac = tuple(int(accent.lstrip("#")[i:i + 2], 16) for i in (0, 2, 4))
    for radius in range(410, 30, -18):
        alpha = max(0, round(2 + (410 - radius) * .025))
        gd.ellipse((width - 330 - radius, -170 - radius, width - 330 + radius, -170 + radius), fill=(*ac, alpha))
    image = Image.alpha_composite(image.convert("RGBA"), glow).convert("RGB")
    if textured:
        noise = Image.effect_noise(image.size, 18).convert("RGB")
        image = Image.blend(image, noise, .035)
    image.save(path, format="PNG", optimize=True)


def _place_image(slide, path: str, x: float, y: float, w: float, h: float):
    p = Path(path)
    if not p.is_file():
        return None
    with Image.open(p) as im:
        scale = max(w * 160 / im.width, h * 160 / im.height)
        target = (max(1, round(w * 160)), max(1, round(h * 160)))
        fitted = ImageOps.fit(im.convert("RGB"), target, method=Image.Resampling.LANCZOS)
        handle = tempfile.NamedTemporaryFile(prefix="ppt-mcp-image-", suffix=".jpg", delete=False)
        temp = Path(handle.name); handle.close(); fitted.save(temp, quality=92)
    try:
        pic = slide.shapes.add_picture(str(temp), Inches(x), Inches(y), Inches(w), Inches(h))
        pic.name = "VISUAL_IMAGE_CLAMPED"
    finally:
        temp.unlink(missing_ok=True)
    return pic


def _public_image_url(url: str) -> bool:
    parsed = urlparse(url)
    if parsed.scheme not in {"http", "https"} or not parsed.hostname:
        return False
    hostname = parsed.hostname.lower()
    if hostname.endswith((".wikimedia.org", ".wikipedia.org")):
        return True
    try:
        addresses = {info[4][0] for info in socket.getaddrinfo(hostname, parsed.port or 443)}
        return bool(addresses) and all(not ipaddress.ip_address(address).is_private
                                       and not ipaddress.ip_address(address).is_loopback
                                       and not ipaddress.ip_address(address).is_link_local
                                       for address in addresses)
    except (OSError, ValueError):
        return False


def _download_image(url: str, target: Path, max_bytes: int = 12 * 1024 * 1024) -> Path | None:
    if not _public_image_url(url):
        return None
    request = Request(url, headers={"User-Agent": "canva-ppt-mcp/2.2 image asset"})
    try:
        with urlopen(request, timeout=20) as response:
            data = response.read(max_bytes + 1)
        if len(data) > max_bytes:
            return None
        target.write_bytes(data)
        with Image.open(target) as image:
            image.verify()
        return target
    except Exception:
        target.unlink(missing_ok=True)
        return None


class AutoDeckRenderer:
    def __init__(self, plan: DeckPlan):
        self.plan = plan; self.ds = plan.design_system
        self.prs = Presentation(); self.prs.slide_width = Inches(SW); self.prs.slide_height = Inches(SH)
        self.current_index = 0
        self.current_layout = "title"
        self.current_dark = False
        self.asset_dir = Path(tempfile.mkdtemp(prefix="ppt-mcp-assets-"))
        self._gradient_cache: dict[str, Path] = {}

    @property
    def p(self): return self.ds.palette
    @property
    def t(self): return self.ds.typography

    @property
    def ko(self): return self.plan.language == "ko"

    def _dark_indices(self) -> set[int]:
        total = len(self.plan.slides)
        target = max(2, round(total * self.ds.dark_slide_ratio))
        selected = {0, total - 1}
        selected.update(i for i, slide in enumerate(self.plan.slides) if slide.layout == "image_focus")
        candidates = [i for i in range(1, total - 1) if i % 2 == 1 and i not in selected]
        candidates += [i for i in range(1, total - 1) if i not in candidates and i not in selected]
        selected.update(candidates[:max(0, target - len(selected))])
        return selected

    def _gradient_asset(self, dark: bool) -> Path:
        key = "dark" if dark else "light"
        if key not in self._gradient_cache:
            target = self.asset_dir / f"gradient-{key}.png"
            if dark:
                start = self.p.background_dark or self.p.primary
                end = mix(start, self.p.primary, .72)
            else:
                start = self.p.background_light
                end = tint(self.p.secondary[0], .55)
            _gradient_png(target, start, end, self.p.accent, self.ds.background_texture)
            self._gradient_cache[key] = target
        return self._gradient_cache[key]

    def _decorate_background(self, slide, dark: bool) -> None:
        if not self.ds.dynamic_composition or self.current_layout in {"title", "closing"}:
            return
        fill = self.p.accent if dark else self.p.primary
        mode = self.current_index % 3
        if mode == 0:
            _add_circle(slide, 10.12, .08, 3.05, fill, "DECOR_CROPPED_ORB", transparency=.86)
        elif mode == 1:
            field = _add_round_rect(slide, .35, 5.10, 1.85, 1.85, fill, radius=False,
                                    name="DECOR_ROTATED_FIELD", transparency=.89)
            field.rotation = 18
        else:
            block = _add_round_rect(slide, 10.82, 4.78, 1.95, 1.95,
                                    self.p.secondary[-1], radius=False,
                                    name="DECOR_ASYMMETRIC_FIELD", transparency=.88)
            block.rotation = 12

    def _base(self, dark=None, emphasis=False):
        if dark is None:
            dark = self.current_index in self._dark_indices()
        self.current_dark = bool(dark)
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = self.p.background_dark if dark and self.p.background_dark else self.p.primary
        base = bg if dark else self.p.background_light
        slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb(base)
        if self.ds.gradient_backgrounds and emphasis:
            slide.shapes.add_picture(str(self._gradient_asset(bool(dark))), 0, 0,
                                     width=self.prs.slide_width, height=self.prs.slide_height).name = "DECOR_GRADIENT_BACKGROUND"
        self._decorate_background(slide, bool(dark))
        return slide, bg if dark else self.p.background_light

    def _header(self, slide, spec, dark=False):
        fg = "#FFFFFF" if dark else self.p.primary
        display_width = sum(2 if ord(c) > 127 else 1 for c in spec.title)
        title_size = 40 if display_width > 74 else (44 if display_width > 34 else self.t.title_size)
        # LibreOffice occasionally wraps a Korean title into an invisible first
        # line when the text box is exactly one line high.  A little vertical
        # breathing room keeps the title stable in both PowerPoint and PDF.
        return _add_text(slide, spec.title, MARGIN, 0.34, 12.05, 1.55,
                         self.t.header_font, title_size, fg, True, name="TITLE", wrap=False)

    def _item_text(self, slide, item, x, y, w, h, dark=False, index=0):
        fg = "#FFFFFF" if dark else self.p.primary
        muted = "#DCE7EF" if dark else "#40566B"
        display_width = sum(2 if ord(c) > 127 else 1 for c in item.heading)
        one_line_capacity = max(8, int(w * 72 / (22 * .75)))
        heading_h = .96 if display_width > one_line_capacity else .48
        body_offset = heading_h + .1
        _add_text(slide, item.heading, x, y, w, heading_h, self.t.body_font, 22, fg, True, name=f"ITEM_{index}_HEAD")
        _add_text(slide, item.body, x, y + body_offset, w, h - body_offset, self.t.body_font, self.t.body_size, muted, name=f"ITEM_{index}_BODY")

    def title(self, spec):
        preset = self.ds.style_preset
        s, _ = self._base(True, emphasis=True); fg = "#FFFFFF"; muted = "#DCE7EF"

        if preset == "editorial":
            _add_text(s, "01", 9.0, .6, 3.3, 1.7, self.t.header_font, 92, tint(self.p.primary, .7), True, name="DECOR_NUMBER")
            _add_round_rect(s, 9.15, 3.05, 3.0, 3.35, tint(self.p.accent, .72), radius=False, name="DECOR_FIELD")
            _add_circle(s, 10.0, 3.85, 1.25, self.p.accent, "DECOR_ACCENT")
        elif preset == "organic":
            _add_circle(s, 9.15, .55, 2.9, tint(self.p.primary, .65), "DECOR_ORGANIC_1")
            _add_circle(s, 10.35, 2.45, 2.15, tint(self.p.accent, .42), "DECOR_ORGANIC_2")
            _add_circle(s, 8.65, 4.7, 1.55, self.p.accent, "DECOR_ORGANIC_3")
        elif preset == "luxury":
            field = _add_round_rect(s, 9.45, 1.05, 2.55, 2.55, self.p.accent, radius=False, name="DECOR_DIAMOND")
            field.rotation = 45
        elif preset == "geometric":
            _add_round_rect(s, 9.15, .65, 3.2, 5.95, self.p.secondary[-1], radius=False, name="DECOR_GEOMETRIC", transparency=.12)
            _add_round_rect(s, 10.05, 1.5, 1.45, 1.45, self.p.accent, radius=False, name="DECOR_ACCENT")
            _add_circle(s, 9.75, 4.25, 1.9, tint(self.p.secondary[0], .25), "DECOR_CIRCLE")
        elif preset == "swiss":
            _add_circle(s, 9.4, .65, 2.8, self.p.accent, "DECOR_SWISS")
            _add_text(s, "A", 9.72, 1.02, 2.15, 1.45, self.t.header_font, 72, "#FFFFFF", True, PP_ALIGN.CENTER, name="DECOR_LETTER")
        elif preset == "neon":
            _add_circle(s, 9.05, .35, 3.25, self.p.secondary[-1], "DECOR_NEON_1")
            _add_circle(s, 10.0, 1.25, 1.25, self.p.accent, "DECOR_NEON_2")
            _add_circle(s, 10.15, 5.0, 1.65, tint(self.p.accent, .18), "DECOR_NEON_3")
        else:
            _add_circle(s, 9.25, 0.22, 3.45, tint(self.p.accent, .15), "DECOR_MOTIF")
            _add_circle(s, 10.2, 4.9, 1.9, self.p.accent, "DECOR_ACCENT")

        _add_text(s, spec.title, 0.82, 1.48, 7.8, 2.05, self.t.header_font,
                  min(60, max(50, self.t.title_size + 4)), fg, True, name="TITLE", wrap=False)
        _add_text(s, spec.subtitle, 0.86, 3.65, 6.7, 1.0, self.t.body_font, 19, muted, name="SUBTITLE")

    def closing(self, spec):
        s, _ = self._base(True, emphasis=True)
        # Closing compositions intentionally differ from title motifs. This
        # prevents the recurring title/ending ornament that made earlier decks
        # look machine-generated.
        preset = self.ds.style_preset
        if preset == "editorial":
            mark = _add_round_rect(s, 9.8, .35, 2.25, 5.8, self.p.accent, False,
                                   "DECOR_CLOSING_EDITORIAL", transparency=.18)
            mark.rotation = -11
        elif preset == "neon":
            mark = _add_round_rect(s, 10.1, .25, 1.75, 6.8, self.p.accent, False,
                                   "DECOR_CLOSING_NEON", transparency=.35)
            mark.rotation = 8
        elif preset == "organic":
            _add_round_rect(s, 9.25, .72, 3.35, 5.85, tint(self.p.accent, .15), True,
                            "DECOR_CLOSING_ORGANIC", transparency=.18)
        elif preset == "luxury":
            _add_text(s, "결론" if self.ko else "FIN", 8.9, .25, 3.6, 2.0, self.t.header_font, 88,
                      self.p.accent, True, PP_ALIGN.CENTER, name="DECOR_CLOSING_LUXURY")
        elif preset == "geometric":
            mark = _add_round_rect(s, 9.75, .8, 2.25, 5.7, self.p.accent, False,
                                   "DECOR_CLOSING_GEOMETRIC", transparency=.15)
            mark.rotation = 22
        elif preset == "swiss":
            _add_circle(s, 9.45, .45, 2.65, self.p.accent, "DECOR_CLOSING_SWISS", transparency=.18)
        else:
            _add_round_rect(s, 9.35, .9, 2.7, 5.6, self.p.accent, True,
                            "DECOR_CLOSING_FIELD", transparency=.25)
        closing_width = sum(1.65 if ord(c) > 127 else 1 for c in spec.title)
        closing_size = 44 if closing_width > 26 else 52
        _add_text(s, spec.title, 0.9, 1.18, 8.7, 2.05, self.t.header_font, closing_size, "#FFFFFF", True, name="TITLE", wrap=False)
        _add_text(s, spec.subtitle, 0.92, 3.3, 7.8, .9, self.t.body_font, 18, "#DCE7EF", name="SUBTITLE")
        if spec.items:
            self._item_text(s, spec.items[0], .95, 4.50, 6.25, 1.68, True)

    def two_column(self, spec):
        s, _ = self._base(); dark = self.current_dark; self._header(s, spec, dark)
        items = (spec.items + [ContentItem(heading="핵심", body=spec.subtitle)])[:2]
        for i in range(2):
            x = 0.72 + i * 6.25
            panel = mix(self.p.background_dark or self.p.primary, self.p.secondary[0], .34) if dark else tint(self.p.secondary[0], .60)
            _add_round_rect(s, x, 2.02, 5.7, 4.38, panel, name=f"VISUAL_PANEL_{i}")
            size = .86 if i == 0 and self.ds.dynamic_composition else .58
            _add_circle(s, x + .42, 2.28, size, self.p.accent if i == 0 else ("#FFFFFF" if dark else self.p.primary), f"VISUAL_ICON_{i}")
            if i < len(items): self._item_text(s, items[i], x + .42, 3.20, 4.85, 2.30, dark, index=i)

    def icon_rows(self, spec):
        s, _ = self._base(); dark = self.current_dark; self._header(s, spec, dark)
        items = spec.items[:4] or [ContentItem(heading="핵심 관점", body=spec.subtitle)]
        fg = "#FFFFFF" if dark else self.p.primary
        muted = "#DCE7EF" if dark else "#40566B"
        row_step = 1.52 if len(items) <= 3 else 1.25
        body_height = 1.36 if len(items) <= 3 else 1.10
        for i, item in enumerate(items):
            y = 2.02 + i * row_step
            size = .98 if i == 0 and self.ds.dynamic_composition else .62
            fill = self.p.accent if i % 2 == 0 else ("#FFFFFF" if dark else self.p.primary)
            _add_circle(s, .78, y + .02, size, fill, f"VISUAL_ICON_{i}")
            label_color = self.p.primary if fill != self.p.primary else "#FFFFFF"
            _add_text(s, str(i + 1), .78, y + .18, size, .25, self.t.body_font, 12, label_color, True, PP_ALIGN.CENTER, name=f"ICON_LABEL_{i}")
            _add_text(s, item.heading, 1.95 if i == 0 else 1.72, y + .04, 2.7, 1.05,
                      self.t.body_font, 20, fg, True, name=f"ITEM_{i}_HEAD")
            _add_text(s, item.body, 4.78, y + .03, 7.48, body_height,
                      self.t.body_font, self.t.body_size, muted, name=f"ITEM_{i}_BODY")

    def big_stat(self, spec):
        s, _ = self._base(); dark = self.current_dark; self._header(s, spec, dark)
        item = spec.items[0] if spec.items else ContentItem(heading="핵심 지표", body=spec.subtitle, value="01")
        panel = mix(self.p.background_dark or self.p.primary, self.p.accent, .26) if dark else tint(self.p.accent, .82)
        _add_round_rect(s, .72, 2.02, 5.2, 4.45, panel, name="VISUAL_STAT")
        stat_color = self.p.accent if dark else self.p.primary
        _add_text(s, item.value or "01", .98, 2.12, 4.65, 1.75, self.t.header_font, 100, stat_color, True, name="STAT_VALUE")
        _add_text(s, item.heading, 1.08, 4.18, 4.35, 1.08, self.t.body_font, 24, "#FFFFFF" if dark else self.p.primary, True, name="STAT_LABEL")
        fg = "#FFFFFF" if dark else self.p.primary; muted = "#DCE7EF" if dark else "#40566B"
        _add_text(s, "왜 중요한가" if self.ko else "Why it matters", 6.42, 2.08, 5.35, .48,
                  self.t.body_font, 24, fg, True, name="ITEM_0_HEAD")
        _add_text(s, item.body, 6.42, 2.72, 5.35, 1.8, self.t.body_font, self.t.body_size, muted, name="ITEM_0_BODY")

    def grid_2x2(self, spec):
        s, _ = self._base(); dark = self.current_dark; self._header(s, spec, dark)
        items = spec.items[:4] or [ContentItem(heading="핵심", body=spec.subtitle)]
        coords = [(.72, 2.02, 7.15, 1.85), (8.18, 2.02, 4.43, 1.85),
                  (.72, 4.12, 5.7, 2.05), (6.92, 4.12, 5.69, 2.05)]
        for i, item in enumerate(items):
            x, y, w, h = coords[i]
            fill = mix(self.p.background_dark or self.p.primary, self.p.secondary[i % len(self.p.secondary)], .30) if dark else tint(self.p.secondary[i % len(self.p.secondary)], .68)
            _add_round_rect(s, x, y, w, h, fill, name=f"VISUAL_CELL_{i}")
            _add_circle(s, x + .3, y + .35, .44, self.p.accent, f"VISUAL_ICON_{i}")
            self._item_text(s, item, x + 1.0, y + .24, w - 1.35, h - .30, dark, index=i)

    def timeline(self, spec):
        s, _ = self._base(); dark = self.current_dark; self._header(s, spec, dark)
        items = spec.items[:4] or [ContentItem(heading="시작", body=spec.subtitle)]
        width = 11.6 / max(1, len(items))
        for i, item in enumerate(items):
            x = .72 + i * width
            size = 1.05 if i == 0 and self.ds.dynamic_composition else .72
            fill = self.p.accent if i == 0 else ("#FFFFFF" if dark else self.p.primary)
            _add_circle(s, x + .15, 2.15, size, fill, f"VISUAL_STEP_{i}")
            label_color = self.p.primary if fill != self.p.primary else "#FFFFFF"
            _add_text(s, f"{i+1:02}", x + .15, 2.41, size, .24, self.t.body_font, 12, label_color, True, PP_ALIGN.CENTER, name=f"STEP_NUM_{i}")
            # Timeline cards often carry complete source sentences.  Preserve
            # them instead of letting PowerPoint render a visual ellipsis.
            self._item_text(s, item, x + .05, 3.34, width - .35, 2.84, dark, index=i)

    def comparison(self, spec):
        s, _ = self._base(); dark = self.current_dark
        # A large half-slide color block is not a decorative stripe; it creates
        # one intentional 50/50 composition in the deck.
        block_fill = self.p.accent if dark else self.p.secondary[-1]
        _add_round_rect(s, 6.66, 0, 6.67, 7.5, block_fill, False,
                        "DECOR_COLOR_BLOCK", transparency=.82)
        self._header(s, spec, dark)
        fallback = ContentItem(heading="대안" if self.ko else "Alternative",
                               body="비교할 항목" if self.ko else "Comparison item")
        items = (spec.items + [fallback])[:2]
        for i, item in enumerate(items):
            x = .72 + i * 6.2
            fill = (mix(self.p.background_dark or self.p.primary, self.p.secondary[i % len(self.p.secondary)], .30)
                    if dark else tint(self.p.primary if i == 0 else self.p.accent, .82))
            _add_round_rect(s, x, 2.02, 5.72, 4.38, fill, name=f"VISUAL_COMPARE_{i}")
            mark_color = self.p.accent if dark else self.p.primary
            _add_text(s, "A" if i == 0 else "B", x + .42, 2.18, 1.1, .92, self.t.header_font, 48, mark_color, True, name=f"COMPARE_MARK_{i}")
            self._item_text(s, item, x + .42, 3.16, 4.85, 2.14, dark, index=i)

    def image_focus(self, spec):
        item = spec.items[0] if spec.items else ContentItem(heading="핵심 장면", body=spec.subtitle)
        image_path = item.image_path if item.image_path and Path(item.image_path).is_file() else None
        if not image_path and item.image_url:
            downloaded = _download_image(item.image_url, self.asset_dir / f"image-{self.current_index}.bin")
            image_path = str(downloaded) if downloaded else None
        if image_path:
            s, _ = self._base(True)
            _place_image(s, image_path, 0, 0, SW, SH)
            _add_round_rect(s, 0, 0, 7.15, SH, self.p.background_dark or self.p.primary,
                            False, "DECOR_IMAGE_OVERLAY", transparency=.15)
            self._header(s, spec, True)
            self._item_text(s, item, .78, 2.05, 5.45, 2.8, True, index=0)
        else:
            dark = self.current_index in self._dark_indices()
            s, _ = self._base(dark, emphasis=True)
            # Treat the fallback as a hero visual: one very large anchor plus
            # small supporting copy, rather than a fake image placeholder.
            _add_circle(s, 7.65, 1.05, 4.95, self.p.accent, "VISUAL_HERO_ICON", transparency=.10)
            _add_text(s, "✦", 8.42, 1.78, 3.35, 2.1, self.t.header_font, 112,
                      self.p.primary, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE,
                      name="VISUAL_HERO_GLYPH")
            self._header(s, spec, dark)
            self._item_text(s, item, .78, 2.05, 5.25, 2.8, dark, index=0)

    def chart(self, spec):
        s, _ = self._base(); dark = self.current_dark; self._header(s, spec, dark)
        chart_spec = spec.chart
        if not chart_spec:
            return self._item_text(s, spec.items[0] if spec.items else ContentItem(
                heading="핵심 수치" if self.ko else "Key metric", body=spec.subtitle),
                .78, 2.0, 11.7, 3.6, dark, index=0)
        data = CategoryChartData(); data.categories = chart_spec.categories
        for series in chart_spec.series:
            data.add_series(series.name, series.values)
        kind = {
            "bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE_MARKERS,
            "pie": XL_CHART_TYPE.PIE,
        }[chart_spec.chart_type]
        chart = s.shapes.add_chart(kind, Inches(.78), Inches(2.02), Inches(8.05), Inches(4.65), data).chart
        chart.has_legend = len(chart_spec.series) > 1 or chart_spec.chart_type == "pie"
        if chart.has_legend:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
        chart.has_title = False
        if chart_spec.chart_type != "pie":
            chart.value_axis.has_major_gridlines = True
            if chart_spec.value_suffix:
                chart.value_axis.tick_labels.number_format = f'0"{chart_spec.value_suffix}"'
        for series_index, series in enumerate(chart.series):
            series.format.fill.solid()
            color = self.p.accent if series_index == 0 else self.p.secondary[series_index % len(self.p.secondary)]
            series.format.fill.fore_color.rgb = rgb(color)
        if spec.items:
            self._item_text(s, spec.items[0], 9.25, 2.02, 3.25, 3.15, dark, index=0)

    def build(self, output_path: str):
        handlers = {name: getattr(self, name) for name in ["title", "closing", "two_column", "icon_rows", "big_stat", "grid_2x2", "timeline", "comparison", "image_focus", "chart"]}
        try:
            for i, spec in enumerate(self.plan.slides):
                layout = spec.layout
                if i > 0 and layout == self.plan.slides[i - 1].layout and layout not in {"title", "closing"}:
                    layout = "two_column" if layout != "two_column" else "icon_rows"
                self.current_index = i
                self.current_layout = layout
                handlers.get(layout, self.two_column)(spec)
            self.prs.save(output_path)
        finally:
            for asset in self.asset_dir.glob("*"):
                asset.unlink(missing_ok=True)
            self.asset_dir.rmdir()


def render_auto_deck(plan: DeckPlan, output_path: str) -> None:
    AutoDeckRenderer(plan).build(output_path)


def write_speaker_notes(pptx_path: str, plan: DeckPlan) -> None:
    prs = Presentation(pptx_path)
    for slide, spec in zip(prs.slides, plan.slides):
        if not spec.speaker_notes:
            continue
        notes = slide.notes_slide.notes_text_frame
        existing = notes.text.strip()
        if spec.speaker_notes not in existing:
            notes.text = (existing + "\n\n" + spec.speaker_notes).strip()
    prs.save(pptx_path)
