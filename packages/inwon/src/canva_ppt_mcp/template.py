from __future__ import annotations

from collections import Counter
from contextlib import contextmanager
from copy import deepcopy
from pathlib import Path
import shutil
import subprocess
import tempfile
import zipfile
import xml.etree.ElementTree as ET
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.oxml.ns import qn

from .render import _download_image


ROLE_TYPES = {
    "title": {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE},
    "body": {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.SUBTITLE},
}
PROTECTED_PLACEHOLDERS = {PP_PLACEHOLDER.DATE, PP_PLACEHOLDER.FOOTER, PP_PLACEHOLDER.SLIDE_NUMBER}


@contextmanager
def _as_pptx(path: str | Path):
    source = Path(path).resolve()
    if source.suffix.lower() != ".potx":
        yield str(source); return
    soffice = shutil.which("soffice")
    if not soffice:
        raise RuntimeError(".potx conversion requires LibreOffice (soffice)")
    with tempfile.TemporaryDirectory(prefix="ppt-mcp-potx-") as temp:
        temp_path = Path(temp)
        profile = temp_path / "profile"; profile.mkdir()
        subprocess.run([soffice, f"-env:UserInstallation={profile.as_uri()}", "--headless",
                        "--convert-to", "pptx", "--outdir", str(temp_path), str(source)],
                       check=True, capture_output=True, text=True, timeout=180)
        converted = temp_path / (source.stem + ".pptx")
        if not converted.exists():
            raise RuntimeError("LibreOffice did not produce a PPTX from the POTX template")
        yield str(converted)


def _hex_from_color(color: Any) -> str | None:
    try:
        if color.type and color.rgb:
            return f"#{color.rgb}"
    except (AttributeError, ValueError):
        return None
    return None


def classify_slide(slide) -> str:
    flattened = _all_shapes(slide)
    names = {s.name for s in flattened}
    if any(n.startswith("VISUAL_COMPARE_") for n in names): return "comparison"
    if any(n.startswith("VISUAL_STEP_") for n in names): return "timeline"
    if any(n.startswith("VISUAL_CELL_") for n in names): return "grid_2x2"
    if "STAT_VALUE" in names: return "big_stat"
    if any(n.startswith("VISUAL_PANEL_") for n in names): return "two_column"
    if sum(n.startswith("VISUAL_ICON_") for n in names) >= 2: return "icon_rows"
    if "SUBTITLE" in names and "TITLE" in names: return "title"
    images = sum(1 for s in flattened if s.shape_type == MSO_SHAPE_TYPE.PICTURE)
    charts = sum(1 for s in flattened if getattr(s, "has_chart", False))
    texts = [s for s in flattened if getattr(s, "has_text_frame", False) and s.text.strip()]
    groups = sum(1 for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.GROUP)
    if len(texts) <= 2 and not charts and images <= 1:
        return "title"
    if charts:
        return "big_stat"
    if groups >= 3 or len(texts) >= 6:
        return "grid_2x2"
    if images and len(texts) >= 2:
        return "two_column"
    return "icon_rows"


def classify_layout(layout) -> str:
    placeholders = list(layout.placeholders)
    title_count = sum(p.placeholder_format.type in ROLE_TYPES["title"] for p in placeholders)
    body_count = sum(p.placeholder_format.type in ROLE_TYPES["body"] for p in placeholders)
    picture_count = sum(p.placeholder_format.type == PP_PLACEHOLDER.PICTURE for p in placeholders)
    name = layout.name.lower()
    if picture_count:
        return "image_focus"
    if "section" in name or (title_count and not body_count):
        return "title"
    if body_count >= 2 or "two" in name or "comparison" in name:
        return "two_column"
    if body_count == 1:
        return "icon_rows"
    return "title"


def _theme_tokens(path: str) -> tuple[list[str], list[str]]:
    colors: list[str] = []; fonts: list[str] = []
    try:
        with zipfile.ZipFile(path) as archive:
            theme_names = [name for name in archive.namelist() if name.startswith("ppt/theme/theme") and name.endswith(".xml")]
            for name in theme_names:
                root = ET.fromstring(archive.read(name))
                for element in root.iter():
                    tag = element.tag.rsplit("}", 1)[-1]
                    if tag == "srgbClr" and element.attrib.get("val"):
                        value = "#" + element.attrib["val"].upper()
                        if value not in colors: colors.append(value)
                    if tag in {"latin", "ea", "cs"} and element.attrib.get("typeface"):
                        value = element.attrib["typeface"]
                        if value and value not in fonts: fonts.append(value)
    except (OSError, zipfile.BadZipFile, ET.ParseError):
        pass
    return colors, fonts


def inspect_template(path: str | Path) -> dict:
    with _as_pptx(path) as load_path:
        prs = Presentation(load_path)
        theme_colors, theme_fonts = _theme_tokens(load_path)
    colors: Counter[str] = Counter()
    fonts: Counter[str] = Counter()
    slides = []
    for index, slide in enumerate(prs.slides, 1):
        slot_count = 0
        image_count = 0
        for shape in _all_shapes(slide):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_count += 1
            if getattr(shape, "has_text_frame", False):
                if shape.is_placeholder or shape.text.strip():
                    slot_count += 1
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.name:
                            fonts[run.font.name] += 1
                        color = _hex_from_color(run.font.color)
                        if color:
                            colors[color] += 1
            try:
                color = _hex_from_color(shape.fill.fore_color)
                if color:
                    colors[color] += 1
            except (AttributeError, TypeError, ValueError):
                pass
        slides.append({
            "slide": index,
            "layout_name": slide.slide_layout.name,
            "classified_as": classify_slide(slide),
            "slot_count": slot_count,
            "image_count": image_count,
        })
    return {
        "path": str(Path(path).resolve()),
        "slide_count": len(prs.slides),
        "slide_size_inches": [round(prs.slide_width / 914400, 2), round(prs.slide_height / 914400, 2)],
        "slides": slides,
        "layouts": slides,
        "available_layouts": [{
            "layout_index": index, "layout_name": layout.name,
            "classified_as": classify_layout(layout),
            "placeholder_count": len(layout.placeholders),
        } for index, layout in enumerate(prs.slide_layouts)],
        "palette": list(dict.fromkeys([c for c, _ in colors.most_common(12)] + theme_colors))[:16],
        "fonts": list(dict.fromkeys([f for f, _ in fonts.most_common(8)] + theme_fonts))[:12],
    }


def _clone_slide(prs: Presentation, source_slide):
    destination = prs.slides.add_slide(source_slide.slide_layout)
    for shape in list(destination.shapes):
        destination.shapes._spTree.remove(shape.element)
    rel_map = {}
    for rel in source_slide.part.rels.values():
        if rel.reltype.endswith("/slideLayout") or rel.reltype.endswith("/notesSlide"):
            continue
        target = rel.target_ref if rel.is_external else rel.target_part
        rel_map[rel.rId] = destination.part.rels._add_relationship(rel.reltype, target, rel.is_external)
    for shape in source_slide.shapes:
        element = deepcopy(shape.element)
        for node in element.iter():
            for attr_name, attr_value in list(node.attrib.items()):
                if attr_value in rel_map and attr_name in {qn("r:embed"), qn("r:link"), qn("r:id")}:
                    node.set(attr_name, rel_map[attr_value])
        destination.shapes._spTree.insert_element_before(element, "p:extLst")
    source_bg = source_slide._element.cSld.bg
    if source_bg is not None:
        current_bg = destination._element.cSld.bg
        if current_bg is not None: destination._element.cSld.remove(current_bg)
        destination._element.cSld.insert(0, deepcopy(source_bg))
    return destination


def _remove_slide(prs: Presentation, slide) -> None:
    slide_id = next(x for x in prs.slides._sldIdLst if x.rId == slide.part.partname.rsplit("/", 1)[-1]) if False else None
    for sld_id in list(prs.slides._sldIdLst):
        if prs.part.related_part(sld_id.rId) is slide.part:
            prs.part.drop_rel(sld_id.rId)
            prs.slides._sldIdLst.remove(sld_id)
            break


def _set_text_preserve_style(shape, text: str) -> None:
    tf = shape.text_frame
    style = None
    if tf.paragraphs and tf.paragraphs[0].runs:
        r = tf.paragraphs[0].runs[0]
        style = (r.font.name, r.font.size, r.font.bold, r.font.italic, _hex_from_color(r.font.color))
    original_alignment = tf.paragraphs[0].alignment if tf.paragraphs else None
    original_level = tf.paragraphs[0].level if tf.paragraphs else 0
    tf.clear(); p = tf.paragraphs[0]
    if original_alignment is not None:
        p.alignment = original_alignment
    p.level = original_level
    run = p.add_run(); run.text = text
    if style:
        run.font.name, run.font.size, run.font.bold, run.font.italic = style[:4]
        if style[4]:
            from pptx.dml.color import RGBColor
            run.font.color.rgb = RGBColor.from_string(style[4][1:])


def _editable_text_slots(slide):
    slots = []
    def walk(shapes, cluster=None):
        for shape in shapes:
            root = cluster or shape
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                walk(shape.shapes, root)
            elif getattr(shape, "has_text_frame", False) and (shape.is_placeholder or shape.text.strip()):
                if shape.is_placeholder and shape.placeholder_format.type in PROTECTED_PLACEHOLDERS:
                    continue
                lowered = shape.name.lower()
                if any(token in lowered for token in ("footer", "slide number", "date", "copyright")):
                    continue
                if shape.top > 6.85 * 914400 and shape.height < .45 * 914400:
                    continue
                slots.append((shape, root))
    walk(slide.shapes)
    return slots


def _all_shapes(slide):
    result = []
    def walk(shapes):
        for shape in shapes:
            result.append(shape)
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP: walk(shape.shapes)
    walk(slide.shapes)
    return result


def _fill_named_slots(slide, spec) -> bool:
    """Fill decks created by this engine without flattening paired slots."""
    by_name = {s.name: s for s in _all_shapes(slide)}
    if "TITLE" not in by_name:
        return False
    _set_text_preserve_style(by_name["TITLE"], spec.title)
    if "SUBTITLE" in by_name:
        _set_text_preserve_style(by_name["SUBTITLE"], spec.subtitle)
    if "STAT_VALUE" in by_name:
        value = spec.items[0].value if spec.items else ""
        _set_text_preserve_style(by_name["STAT_VALUE"], value or "01")
    for index in range(6):
        head = by_name.get(f"ITEM_{index}_HEAD"); body = by_name.get(f"ITEM_{index}_BODY")
        if index < len(spec.items):
            item = spec.items[index]
            if head and not ("STAT_VALUE" in by_name and index == 0):
                _set_text_preserve_style(head, item.heading)
            if body: _set_text_preserve_style(body, item.body)
        else:
            # Remove the full item cluster, including its visual marker/frame.
            prefixes = (f"ITEM_{index}_", f"VISUAL_ICON_{index}", f"VISUAL_CELL_{index}",
                        f"VISUAL_PANEL_{index}", f"VISUAL_STEP_{index}", f"ICON_LABEL_{index}",
                        f"STEP_NUM_{index}")
            for shape in list(_all_shapes(slide)):
                if any(shape.name.startswith(prefix) for prefix in prefixes):
                    shape.element.getparent().remove(shape.element)
    return True


def _replace_image_slot(slide, spec) -> bool:
    item = next((item for item in spec.items if item.image_path or item.image_url), None)
    if not item:
        return False
    path = Path(item.image_path).resolve() if item.image_path else None
    temporary = None
    if (not path or not path.is_file()) and item.image_url:
        handle = tempfile.NamedTemporaryFile(prefix="ppt-template-image-", suffix=".bin", delete=False)
        temporary = Path(handle.name); handle.close()
        path = _download_image(item.image_url, temporary)
    if not path or not path.is_file():
        if temporary: temporary.unlink(missing_ok=True)
        return False
    try:
        picture_placeholders = [shape for shape in slide.placeholders
                                if shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE]
        if picture_placeholders:
            picture_placeholders[0].insert_picture(str(path)); return True
        pictures = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
        named = [shape for shape in pictures if any(token in shape.name.lower()
                 for token in ("image", "photo", "picture", "hero"))]
        candidates = named or (pictures if spec.layout == "image_focus" else [])
        if not candidates:
            return False
        target = max(candidates, key=lambda shape: shape.width * shape.height)
        replacement = slide.shapes.add_picture(str(path), target.left, target.top, target.width, target.height)
        replacement.name = target.name
        target.element.addprevious(replacement.element)
        target.element.getparent().remove(target.element)
        return True
    finally:
        if temporary: temporary.unlink(missing_ok=True)


def apply_content_to_template(template_path: str, output_path: str, plan) -> dict:
    with _as_pptx(template_path) as load_path:
        prs = Presentation(load_path)
        originals = list(prs.slides)
        inventory = [(slide, classify_slide(slide)) for slide in originals]
        layout_inventory = [(layout, classify_layout(layout)) for layout in prs.slide_layouts]
        if not inventory and not layout_inventory:
            raise ValueError("템플릿에 사용할 슬라이드나 레이아웃이 없습니다.")
        mapping = []
        for out_index, spec in enumerate(plan.slides, 1):
            desired_role = "title" if spec.layout == "closing" else spec.layout
            candidates = [x for x in inventory if x[1] == desired_role]
            if not candidates:
                candidates = inventory
            if candidates:
                source, source_role = candidates[(out_index - 1) % len(candidates)]
                slide = _clone_slide(prs, source)
                source_slide_number = originals.index(source) + 1
            else:
                layout_candidates = [x for x in layout_inventory if x[1] == desired_role] or layout_inventory
                source_layout, source_role = layout_candidates[(out_index - 1) % len(layout_candidates)]
                slide = prs.slides.add_slide(source_layout)
                source_slide_number = None
            _replace_image_slot(slide, spec)
            if _fill_named_slots(slide, spec):
                mapping.append({"output_slide": out_index, "source_slide": source_slide_number,
                                "requested_layout": spec.layout, "source_layout": source_role})
                continue
            slots = sorted(_editable_text_slots(slide), key=lambda pair: (pair[0].top, pair[0].left))
            shapes = [pair[0] for pair in slots]
            clusters = {id(shape): cluster for shape, cluster in slots}
            title_shape = next((s for s in shapes if s.is_placeholder and s.placeholder_format.type in ROLE_TYPES["title"]), None)
            if title_shape is None and shapes:
                title_shape = shapes.pop(0)
            elif title_shape in shapes:
                shapes.remove(title_shape)
            if title_shape:
                _set_text_preserve_style(title_shape, spec.title)
            body_texts = ([spec.subtitle] if spec.subtitle else []) + [
                (f"{i.heading}\n{i.body}" if i.heading else i.body) for i in spec.items if i.heading or i.body
            ]
            for shape, text in zip(shapes, body_texts):
                _set_text_preserve_style(shape, text)
            removed = set()
            for shape in shapes[len(body_texts):]:
                # If the slot lives in a group, remove the top-level group so
                # its associated image, icon, and text disappear together.
                cluster = clusters[id(shape)]
                if id(cluster) not in removed:
                    cluster.element.getparent().remove(cluster.element); removed.add(id(cluster))
            mapping.append({"output_slide": out_index, "source_slide": source_slide_number,
                            "requested_layout": spec.layout, "source_layout": source_role})
        for slide in originals:
            _remove_slide(prs, slide)
        prs.save(output_path)
    return {"mode": "template", "mapping": mapping, "template": inspect_template(template_path)}
