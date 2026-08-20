from __future__ import annotations

import sys
import tempfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "src"))

from pptx import Presentation
from pptx.util import Inches, Pt

from canva_ppt_mcp.com_editor import design_fingerprint
from canva_ppt_mcp.routing import (
    confirm_presentation_mode,
    consume_mode_confirmation,
    infer_presentation_mode,
    prepare_presentation_task,
)


def main() -> None:
    failures: list[str] = []

    def check(name, fn):
        try:
            fn(); print(f"PASS {name}")
        except Exception as exc:
            failures.append(f"{name}: {type(exc).__name__}: {exc}")
            print(f"FAIL {failures[-1]}")

    check("generate intent", lambda: (_ for _ in ()).throw(AssertionError())
          if infer_presentation_mode("IBM 소개 PPT 만들어줘") != "generate" else None)
    check("template intent", lambda: (_ for _ in ()).throw(AssertionError())
          if infer_presentation_mode("이 디자인 그대로 내용만 수정해줘") != "template_com" else None)

    def token_case():
        with tempfile.TemporaryDirectory() as temp:
            prepared = prepare_presentation_task("PPT 만들어줘", temp)
            confirmed = confirm_presentation_mode(prepared["confirmation_id"], "generate")
            consume_mode_confirmation(requested_mode="generate", execution_token=confirmed["execution_token"])
            try:
                consume_mode_confirmation(requested_mode="generate", execution_token=confirmed["execution_token"])
            except RuntimeError:
                return
            raise AssertionError("execution token was reusable")
    check("one-time confirmation token", token_case)

    def fingerprint_case():
        with tempfile.TemporaryDirectory() as temp:
            path = Path(temp) / "text-only-change.pptx"
            prs = Presentation(); slide = prs.slides.add_slide(prs.slide_layouts[6])
            box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
            run = box.text_frame.paragraphs[0].add_run(); run.text = "OLD"; run.font.size = Pt(24)
            prs.save(path)
            before = design_fingerprint(path)
            prs = Presentation(path); prs.slides[0].shapes[0].text_frame.paragraphs[0].runs[0].text = "NEW"; prs.save(path)
            assert design_fingerprint(path) == before
    check("design fingerprint ignores content", fingerprint_case)

    if failures:
        raise SystemExit("\n".join(failures))
    print("All core hybrid-mode checks passed.")


if __name__ == "__main__":
    main()
