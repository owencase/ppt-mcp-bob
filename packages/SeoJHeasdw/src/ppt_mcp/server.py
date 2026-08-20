"""MCP 서버.

MCP 에는 세 가지 기본 요소가 있습니다. 대부분 tool 만 쓰는데, 나머지 둘이
있어야 서버가 제 몫을 합니다.

  tool      모델이 *실행*하는 것        — 부작용이 있는 동작
  resource  모델이 *읽는* 것            — 부작용 없는 데이터
  prompt    사용자가 *고르는* 것        — 잘 만든 시작점

테마 팔레트를 예로 들면 이렇습니다. tool 로 만들면 모델이 "색을 알아내려고"
동작을 실행해야 합니다. resource 로 두면 그냥 읽습니다. 읽기와 쓰기를 구분하는
것이 이 설계의 요점입니다.

tool 은 셋뿐입니다. 적을수록 모델이 덜 헷갈립니다.
"""
from __future__ import annotations

import json
from typing import Any

try:  # MCP SDK 2.x
    from mcp.server.mcpserver import MCPServer
except ImportError:  # MCP SDK 1.x — 이름만 다르고 표면은 같습니다
    from mcp.server.fastmcp import FastMCP as MCPServer

from pptx import Presentation

from .models import DeckSpec
from .render import CM_TO_EMU, render_deck, resolve_output_path
from .theme import THEMES

mcp = MCPServer("ppt-mcp-seojh")


# ── tools ──────────────────────────────────────────────────────────────────

@mcp.tool()
def create_deck(spec: DeckSpec, output_path: str) -> dict[str, Any]:
    """덱 하나를 통째로 만들어 .pptx 로 저장한다.

    좌표나 폰트 크기를 지정하지 않는다. 슬라이드 종류(kind)와 내용만 주면
    레이아웃은 서버가 정한다. 넘치는 텍스트는 폰트를 줄여 맞추고, 줄여도
    안 되면 어느 슬라이드가 문제인지 알려주며 실패한다.

    사용 가능한 kind: title, section, bullets, chart.
    테마 팔레트는 resource `theme://list` 에서 읽을 수 있다.
    """
    path = render_deck(spec, output_path)
    return {
        "output_path": str(path),
        "slide_count": len(spec.slides),
        "theme": spec.theme,
        "kinds": [slide.kind for slide in spec.slides],
    }


@mcp.tool()
def validate_deck(spec: DeckSpec) -> dict[str, Any]:
    """파일을 만들지 않고 스펙만 검사한다.

    긴 덱을 만들기 전에 먼저 불러 형식을 확인하는 용도다. 인자가
    DeckSpec 으로 파싱되었다는 것 자체가 통과를 뜻한다.
    """
    return {
        "valid": True,
        "slide_count": len(spec.slides),
        "kinds": [slide.kind for slide in spec.slides],
    }


@mcp.tool()
def describe_deck(pptx_path: str) -> dict[str, Any]:
    """기존 .pptx 를 열어 슬라이드 수와 도형 목록을 돌려준다.

    자기가 만든 결과를 확인하거나, 남이 만든 덱을 살펴볼 때 쓴다.
    """
    path = resolve_output_path(pptx_path)
    if not path.exists():
        raise ValueError(f"파일이 없습니다: {path}")
    presentation = Presentation(str(path))
    return {
        "path": str(path),
        "slide_count": len(presentation.slides),
        "width_cm": round(presentation.slide_width / CM_TO_EMU, 2),
        "height_cm": round(presentation.slide_height / CM_TO_EMU, 2),
        "slides": [
            {
                "index": index,
                "shapes": [
                    {"name": shape.name, "text": getattr(shape, "text", "")[:60]}
                    for shape in slide.shapes
                ],
            }
            for index, slide in enumerate(presentation.slides)
        ],
    }


# ── resources ──────────────────────────────────────────────────────────────
# 읽기 전용 데이터. tool 로 만들면 모델이 색을 알아내려고 '실행'해야 합니다.

@mcp.resource("theme://list")
def theme_list() -> str:
    """쓸 수 있는 테마 전부와 각 팔레트."""
    return json.dumps([theme.as_dict() for theme in THEMES.values()],
                      ensure_ascii=False, indent=2)


@mcp.resource("theme://{name}")
def theme_detail(name: str) -> str:
    """테마 하나의 팔레트."""
    if name not in THEMES:
        raise ValueError(f"'{name}' 테마는 없습니다. 사용 가능: {', '.join(sorted(THEMES))}")
    return json.dumps(THEMES[name].as_dict(), ensure_ascii=False, indent=2)


# ── prompt ─────────────────────────────────────────────────────────────────
# 사용자가 고르는 잘 만든 시작점. 매번 프롬프트를 다시 쓰지 않게 합니다.

@mcp.prompt()
def deck_from_topic(topic: str, slide_count: int = 6) -> str:
    """주제 하나로 덱 구조를 잡는 프롬프트."""
    return (
        f"'{topic}' 주제로 {slide_count}장짜리 발표 자료를 만들어 주세요.\n\n"
        "구성 지침:\n"
        "- 첫 장은 kind='title'\n"
        "- 내용이 세 덩어리 이상이면 kind='section' 으로 나누기\n"
        "- 숫자가 있으면 kind='bullets' 대신 kind='chart'\n"
        "- 불릿은 한 장에 6개 이하, 문장이 아니라 요점으로\n\n"
        "테마는 resource `theme://list` 를 읽고 주제에 맞게 고른 뒤,\n"
        "create_deck 을 한 번만 호출하세요."
    )


def main() -> None:
    mcp.run()


if __name__ == "__main__":
    main()
