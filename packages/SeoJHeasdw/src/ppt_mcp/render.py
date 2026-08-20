"""DeckSpec → .pptx 렌더링.

python-pptx 를 아는 유일한 모듈입니다. 레이아웃 계산은 layout.py 가 이미 끝냈고,
여기서는 그 결과를 도형으로 옮기기만 합니다.

두 가지를 지킵니다.

- **원자적 저장.** 임시 파일에 다 쓴 뒤 os.replace 로 바꿔칩니다. 렌더 도중
  죽어도 반쯤 쓰인 .pptx 가 남지 않습니다. 레퍼런스 구현은 tool 호출마다
  대상 파일을 바로 덮어써서, 중간에 실패하면 깨진 파일이 남습니다.
- **경로 게이트.** 파일을 만드는 지점이 여기 하나뿐이라 검사도 한 곳에서 합니다.
"""
from __future__ import annotations

import os
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

from . import layout
from .models import BulletsSlide, ChartSlide, DeckSpec, SectionSlide, TitleSlide
from .theme import Theme, get_theme

CM_TO_EMU = 360_000
BLANK_LAYOUT = 6


def _emu(value_cm: float) -> int:
    return int(value_cm * CM_TO_EMU)


def _rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def resolve_output_path(raw: str) -> Path:
    """출력 경로를 검증합니다.

    PPT_MCP_OUTPUT_DIR 이 설정되어 있으면 그 안으로만 쓸 수 있습니다. MCP 서버는
    LLM 이 준 경로를 그대로 받으므로, 어디든 쓸 수 있게 두면 안 됩니다.
    """
    path = Path(raw).expanduser().resolve()
    if path.suffix.lower() != ".pptx":
        raise ValueError(f"경로는 .pptx 로 끝나야 합니다: {path.name}")
    allowed = os.getenv("PPT_MCP_OUTPUT_DIR")
    if allowed:
        root = Path(allowed).expanduser().resolve()
        if root != path and root not in path.parents:
            raise ValueError(f"출력은 PPT_MCP_OUTPUT_DIR 안에만 쓸 수 있습니다: {root}")
    return path


def _add_text(slide, box: layout.Box, text: str, *, size_pt: float, color: str,
              bold: bool = False, align=PP_ALIGN.LEFT) -> None:
    shape = slide.shapes.add_textbox(_emu(box.left), _emu(box.top),
                                     _emu(box.width), _emu(box.height))
    frame = shape.text_frame
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)


def _paint_background(slide, theme: Theme) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(theme.bg)


def _render_title(slide, spec: TitleSlide, theme: Theme) -> None:
    boxes = layout.title_slide()
    size = layout.fit_font_size(spec.title, boxes["title"], 44, layout.MIN_TITLE_PT)
    _add_text(slide, boxes["title"], spec.title, size_pt=size, color=theme.accent,
              bold=True, align=PP_ALIGN.CENTER)
    if spec.subtitle:
        size = layout.fit_font_size(spec.subtitle, boxes["subtitle"], 20, layout.MIN_BODY_PT)
        _add_text(slide, boxes["subtitle"], spec.subtitle, size_pt=size,
                  color=theme.muted, align=PP_ALIGN.CENTER)


def _render_section(slide, spec: SectionSlide, theme: Theme) -> None:
    box = layout.section_slide()["title"]
    size = layout.fit_font_size(spec.title, box, 40, layout.MIN_TITLE_PT)
    _add_text(slide, box, spec.title, size_pt=size, color=theme.accent,
              bold=True, align=PP_ALIGN.CENTER)


def _render_bullets(slide, spec: BulletsSlide, theme: Theme) -> None:
    boxes = layout.bullets_slide(len(spec.points))
    size = layout.fit_font_size(spec.title, boxes["title"], 30, layout.MIN_TITLE_PT)
    _add_text(slide, boxes["title"], spec.title, size_pt=size, color=theme.accent, bold=True)

    slot = boxes["bullet_slot"]
    # 모든 불릿이 같은 크기여야 보기 좋습니다. 가장 긴 것에 맞춰 한 번만 정합니다.
    longest = max(spec.points, key=len)
    size = layout.fit_font_size(longest, slot, 20, layout.MIN_BODY_PT)
    for index, point in enumerate(spec.points):
        box = layout.Box(slot.left, slot.top + index * slot.height, slot.width, slot.height)
        _add_text(slide, box, f"·  {point}", size_pt=size, color=theme.text)


def _render_chart(slide, spec: ChartSlide, theme: Theme) -> None:
    boxes = layout.chart_slide()
    size = layout.fit_font_size(spec.title, boxes["title"], 30, layout.MIN_TITLE_PT)
    _add_text(slide, boxes["title"], spec.title, size_pt=size, color=theme.accent, bold=True)

    data = CategoryChartData()
    data.categories = list(spec.series)
    data.add_series("", tuple(spec.series.values()))
    box = boxes["chart"]
    frame = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, _emu(box.left),
                                   _emu(box.top), _emu(box.width), _emu(box.height), data)
    chart = frame.chart
    chart.has_legend = False
    chart.plots[0].vary_by_categories = False
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = _rgb(theme.accent)


_RENDERERS = {
    TitleSlide: _render_title,
    SectionSlide: _render_section,
    BulletsSlide: _render_bullets,
    ChartSlide: _render_chart,
}


def render_deck(spec: DeckSpec, output_path: str) -> Path:
    """덱을 만들고 원자적으로 저장합니다. 저장된 경로를 돌려줍니다."""
    target = resolve_output_path(output_path)
    target.parent.mkdir(parents=True, exist_ok=True)
    theme = get_theme(spec.theme)

    presentation = Presentation()
    presentation.slide_width = Emu(_emu(layout.SLIDE_W))
    presentation.slide_height = Emu(_emu(layout.SLIDE_H))
    blank = presentation.slide_layouts[BLANK_LAYOUT]

    for index, slide_spec in enumerate(spec.slides, 1):
        slide = presentation.slides.add_slide(blank)
        _paint_background(slide, theme)
        try:
            _RENDERERS[type(slide_spec)](slide, slide_spec, theme)
        except ValueError as exc:
            # 몇 번째 슬라이드인지 붙여줘야 호출자가 그 장만 고칠 수 있습니다.
            raise ValueError(f"{index}번째 슬라이드({slide_spec.kind}): {exc}") from None

    # 같은 디렉터리의 임시 파일에 쓴 뒤 교체합니다. 같은 파일시스템이라야
    # os.replace 가 원자적입니다.
    handle, temp_name = tempfile.mkstemp(dir=target.parent, suffix=".pptx.tmp")
    os.close(handle)
    temp_path = Path(temp_name)
    try:
        presentation.save(str(temp_path))
        os.replace(temp_path, target)
    except BaseException:
        temp_path.unlink(missing_ok=True)
        raise
    return target
