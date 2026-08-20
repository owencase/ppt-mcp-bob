"""
ppt-bridge/bridge.py
────────────────────
Python bridge between the MCP Server and Microsoft PowerPoint.

Protocol
--------
stdin  → single JSON line: {"action": "<name>", "params": {...}}
stdout → single JSON line: {"success": true/false, "message": "...", "data": {...}}
stderr → 사람이 읽는 로그 전용

stdout 은 프로토콜 전용입니다. 핸들러가 도는 동안에는 stdout 이 버퍼로
바꿔치기되므로, 디버깅용 print() 를 남겨도 응답 JSON 이 깨지지 않습니다.
그 출력은 stderr 로 옮겨져서 그대로 볼 수 있습니다 (main() 참고).

All measurements that arrive from the MCP server are in centimetres.
python-pptx works in EMUs (1 cm = 360_000 EMU), so we convert internally.

Dependencies
------------
  pip install -r requirements.txt
"""

from __future__ import annotations

import io
import json
import sys
from typing import Any

# ── python-pptx imports ────────────────────────────────────────────────────
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
CM_TO_EMU = 360_000

# Predefined themes  {name: {bg, text, accent1, accent2}}
THEMES: dict[str, dict[str, str]] = {
    "minimal_dark":   {"bg": "1E1E2E", "text": "FFFFFF", "accent1": "CBA6F7", "accent2": "89DCEB"},
    "minimal_light":  {"bg": "FFFFFF", "text": "1E1E2E", "accent1": "3B82F6", "accent2": "8B5CF6"},
    "tech_blue":      {"bg": "0F172A", "text": "E2E8F0", "accent1": "38BDF8", "accent2": "22D3EE"},
    "marketing_warm": {"bg": "FFF7ED", "text": "1C1917", "accent1": "F97316", "accent2": "EAB308"},
}

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def cm(value: float) -> int:
    """Convert centimetres to EMU."""
    return int(value * CM_TO_EMU)


def hex_to_rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def ok(message: str = "OK", data: Any = None) -> dict:
    result: dict[str, Any] = {"success": True, "message": message}
    if data is not None:
        result["data"] = data
    return result


def err(message: str, data: Any = None) -> dict:
    result: dict[str, Any] = {"success": False, "error": message}
    if data is not None:
        result["data"] = data
    return result


def load_prs(file_path: str) -> Presentation:
    return Presentation(file_path)


def save_prs(prs: Presentation, file_path: str) -> None:
    prs.save(file_path)

# ---------------------------------------------------------------------------
# Action handlers
# ---------------------------------------------------------------------------

def handle_create_presentation(params: dict) -> dict:
    file_path: str = params["file_path"]
    width_cm: float = params.get("width_cm", 33.867)   # 16:9 widescreen
    height_cm: float = params.get("height_cm", 19.05)

    prs = Presentation()
    prs.slide_width = Emu(cm(width_cm))
    prs.slide_height = Emu(cm(height_cm))
    save_prs(prs, file_path)
    return ok(f"Presentation created: {file_path}")


def handle_add_slide(params: dict) -> dict:
    file_path: str = params["file_path"]
    layout_index: int = int(params.get("layout_index", 6))

    prs = load_prs(file_path)
    layout = prs.slide_layouts[layout_index]
    prs.slides.add_slide(layout)
    save_prs(prs, file_path)
    slide_count = len(prs.slides)
    return ok(f"Slide added. Total slides: {slide_count}", {"slide_index": slide_count - 1})


def handle_add_text_box(params: dict) -> dict:
    file_path: str = params["file_path"]
    slide_index: int = int(params["slide_index"])
    text: str = params["text"]
    font_size_pt: float = float(params.get("font_size_pt", 24))
    bold: bool = bool(params.get("bold", False))
    color_hex: str = str(params.get("color_hex", "000000"))
    align_str: str = str(params.get("align", "left"))

    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    alignment = align_map.get(align_str, PP_ALIGN.LEFT)

    prs = load_prs(file_path)
    slide = prs.slides[slide_index]

    txBox = slide.shapes.add_textbox(
        cm(float(params["left_cm"])),
        cm(float(params["top_cm"])),
        cm(float(params["width_cm"])),
        cm(float(params["height_cm"])),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold
    run.font.color.rgb = hex_to_rgb(color_hex)

    save_prs(prs, file_path)
    return ok("Text box added")


def handle_add_image(params: dict) -> dict:
    file_path: str = params["file_path"]
    slide_index: int = int(params["slide_index"])
    image_path: str = params["image_path"]

    prs = load_prs(file_path)
    slide = prs.slides[slide_index]
    slide.shapes.add_picture(
        image_path,
        cm(float(params["left_cm"])),
        cm(float(params["top_cm"])),
        cm(float(params["width_cm"])),
        cm(float(params["height_cm"])),
    )
    save_prs(prs, file_path)
    return ok("Image added")


def handle_set_background_color(params: dict) -> dict:
    file_path: str = params["file_path"]
    slide_index: int = int(params["slide_index"])
    color_hex: str = str(params["color_hex"])

    prs = load_prs(file_path)
    slide = prs.slides[slide_index]

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(color_hex)

    save_prs(prs, file_path)
    return ok(f"Background set to #{color_hex}")


def handle_add_shape(params: dict) -> dict:
    file_path: str = params["file_path"]
    slide_index: int = int(params["slide_index"])
    shape_type: str = str(params.get("shape_type", "rectangle"))
    fill_hex: str = str(params.get("fill_color_hex", "4472C4"))
    line_hex: str | None = params.get("line_color_hex")

    prs = load_prs(file_path)
    slide = prs.slides[slide_index]

    shape_id = (
        MSO_SHAPE.ROUNDED_RECTANGLE
        if shape_type == "rounded_rectangle"
        else MSO_SHAPE.RECTANGLE
    )

    shape = slide.shapes.add_shape(
        shape_id,
        cm(float(params["left_cm"])),
        cm(float(params["top_cm"])),
        cm(float(params["width_cm"])),
        cm(float(params["height_cm"])),
    )

    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fill_hex)

    if line_hex:
        shape.line.color.rgb = hex_to_rgb(line_hex)
    else:
        shape.line.fill.background()  # no border

    save_prs(prs, file_path)
    return ok(f"{shape_type} shape added")


def handle_apply_theme(params: dict) -> dict:
    file_path: str = params["file_path"]
    theme_name: str = params["theme"]

    if theme_name not in THEMES:
        return err(f"Unknown theme '{theme_name}'. Available: {list(THEMES.keys())}")

    theme = THEMES[theme_name]
    prs = load_prs(file_path)

    for slide in prs.slides:
        # background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(theme["bg"])

    save_prs(prs, file_path)
    return ok(f"Theme '{theme_name}' applied to all {len(prs.slides)} slides", {"theme": theme})


def handle_save_presentation(params: dict) -> dict:
    file_path: str = params["file_path"]
    prs = load_prs(file_path)
    save_prs(prs, file_path)
    return ok(f"Saved: {file_path}")


def handle_get_presentation_info(params: dict) -> dict:
    file_path: str = params["file_path"]
    prs = load_prs(file_path)

    slides_info = []
    for i, slide in enumerate(prs.slides):
        shapes_info = [
            {"name": s.name, "shape_type": str(s.shape_type), "left_cm": round(s.left / CM_TO_EMU, 2),
             "top_cm": round(s.top / CM_TO_EMU, 2)}
            for s in slide.shapes
        ]
        slides_info.append({"slide_index": i, "shapes": shapes_info})

    data = {
        "slide_count": len(prs.slides),
        "width_cm": round(prs.slide_width / CM_TO_EMU, 2),
        "height_cm": round(prs.slide_height / CM_TO_EMU, 2),
        "slides": slides_info,
    }
    return ok("Presentation info retrieved", data)

# ---------------------------------------------------------------------------
# Dispatch table
# ---------------------------------------------------------------------------
HANDLERS = {
    "create_presentation":    handle_create_presentation,
    "add_slide":              handle_add_slide,
    "add_text_box":           handle_add_text_box,
    "add_image":              handle_add_image,
    "set_background_color":   handle_set_background_color,
    "add_shape":              handle_add_shape,
    "apply_theme":            handle_apply_theme,
    "save_presentation":      handle_save_presentation,
    "get_presentation_info":  handle_get_presentation_info,
}

# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

def dispatch(raw: str) -> dict:
    """요청 JSON 문자열 하나를 받아 응답 dict 를 돌려줍니다.

    프로세스나 stdout 을 건드리지 않는 순수 함수라서 테스트에서 바로 부를 수
    있습니다. 실제 I/O 는 main() 이 담당합니다.
    """
    try:
        request = json.loads(raw)
    except json.JSONDecodeError as exc:
        return err(f"Invalid JSON input: {exc}")

    if not isinstance(request, dict):
        return err("Request must be a JSON object")

    action: str = request.get("action", "")
    params: dict = request.get("params", {})

    handler = HANDLERS.get(action)
    if handler is None:
        available = sorted(HANDLERS)
        # available 을 data 에도 실어 둡니다. 사람이 읽는 메시지와 별개로
        # 기계가 읽을 수 있어야, mcp-server 의 tool 목록과 어긋났는지
        # 테스트가 자동으로 잡아낼 수 있습니다 (mcp-server/test 참고).
        return err(f"Unknown action '{action}'. Available: {available}", {"available": available})

    try:
        return handler(params)
    except KeyError as exc:
        # 필수 파라미터 누락. 무엇이 빠졌는지 이름을 그대로 알려줍니다.
        return err(f"Missing required parameter in '{action}': {exc}")
    except Exception as exc:  # noqa: BLE001
        return err(f"Unhandled exception in '{action}': {type(exc).__name__}: {exc}")


def as_utf8(stream):
    """콘솔 인코딩이 UTF-8 이 아닌 환경(주로 Windows)을 위해 스트림을 감쌉니다.

    버퍼가 없는 스트림(테스트에서 넘기는 StringIO 등)은 그대로 돌려줍니다.
    """
    buffer = getattr(stream, "buffer", None)
    if buffer is None:
        return stream
    return io.TextIOWrapper(buffer, encoding="utf-8", errors="replace")


def main() -> None:
    # 인코딩 강제는 여기서 합니다. import 시점에 전역 sys.stdin/stdout 을 바꾸면
    # 모듈을 import 하는 것만으로 부작용이 생겨 테스트에서 쓸 수 없게 됩니다.
    sys.stdin = as_utf8(sys.stdin)
    sys.stdout = as_utf8(sys.stdout)

    raw = sys.stdin.read().strip()

    # stdout 은 프로토콜 전용입니다. 핸들러가 도는 동안에는 진짜 stdout 을
    # 치워두고 버퍼로 바꿔칩니다. 그래야 누군가 디버깅용 print() 를 남겨도
    # 응답 JSON 한 줄이 오염되지 않습니다.
    real_stdout = sys.stdout
    buffer = io.StringIO()
    sys.stdout = buffer
    try:
        response = dispatch(raw)
    finally:
        sys.stdout = real_stdout
        stray = buffer.getvalue()
        if stray:
            # 삼키지 않고 stderr 로 흘려보냅니다. print() 로 디버깅하던 사람이
            # 출력을 그대로 볼 수 있어야 하니까요.
            print("[bridge] stdout 으로 나간 출력을 stderr 로 옮겼습니다:", file=sys.stderr)
            print(stray, end="" if stray.endswith("\n") else "\n", file=sys.stderr)

    # ensure_ascii 기본값(True)을 유지합니다. 출력이 순수 ASCII 가 되어
    # 콘솔 인코딩이 무엇이든 깨지지 않습니다.
    real_stdout.write(json.dumps(response) + "\n")
    real_stdout.flush()


if __name__ == "__main__":
    main()
