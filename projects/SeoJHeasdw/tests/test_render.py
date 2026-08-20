"""렌더링 테스트. 실제로 .pptx 를 만들어 열어 봅니다."""
import os

import pytest
from pptx import Presentation

from ppt_mcp.models import DeckSpec, SectionSlide
from ppt_mcp.render import render_deck, resolve_output_path

DECK = {
    "theme": "tech_blue",
    "slides": [
        {"kind": "title", "title": "AI 혁신의 시대", "subtitle": "2026 기술 전망"},
        {"kind": "section", "title": "핵심 변화"},
        {"kind": "bullets", "title": "세 가지 축", "points": ["자동화", "개인화", "예측"]},
        {"kind": "chart", "title": "성장률", "series": {"2024": 12, "2025": 31, "2026": 58}},
    ],
}


class TestRender:
    def test_every_kind_renders(self, tmp_path):
        target = tmp_path / "deck.pptx"
        render_deck(DeckSpec.model_validate(DECK), str(target))
        presentation = Presentation(str(target))
        assert len(presentation.slides) == 4

    def test_slide_size_is_16_by_9(self, tmp_path):
        target = tmp_path / "deck.pptx"
        render_deck(DeckSpec.model_validate(DECK), str(target))
        presentation = Presentation(str(target))
        ratio = presentation.slide_width / presentation.slide_height
        assert abs(ratio - 16 / 9) < 0.01

    def test_failure_names_the_slide_number(self, tmp_path):
        # 2번째 슬라이드만 문제일 때, 몇 번째인지 알려줘야 그 장만 고칩니다.
        #
        # model_construct 로 검증을 우회합니다. 정상 경로에서는 스펙 제약이
        # 먼저 막아서 이 에러까지 갈 일이 없지만(아래 불변식 테스트 참고),
        # 렌더 계층의 에러 문구 자체는 따로 확인해 둘 가치가 있습니다.
        spec = DeckSpec.model_construct(theme="tech_blue", slides=[
            SectionSlide.model_construct(kind="section", title="정상"),
            SectionSlide.model_construct(kind="section", title="가" * 400),
        ])
        with pytest.raises(ValueError, match="2번째 슬라이드"):
            render_deck(spec, str(tmp_path / "x.pptx"))


class TestAtomicSave:
    def test_failed_render_leaves_no_partial_file(self, tmp_path):
        # 중간에 실패해도 반쯤 쓰인 .pptx 가 남으면 안 됩니다.
        target = tmp_path / "deck.pptx"
        bad = {"theme": "tech_blue", "slides": [{"kind": "section", "title": "가" * 400}]}
        with pytest.raises(ValueError):
            render_deck(DeckSpec.model_validate(bad), str(target))
        assert not target.exists()
        assert list(tmp_path.iterdir()) == []      # 임시 파일도 안 남음

    def test_rerender_replaces_cleanly(self, tmp_path):
        target = tmp_path / "deck.pptx"
        render_deck(DeckSpec.model_validate(DECK), str(target))
        render_deck(DeckSpec.model_validate(DECK), str(target))
        assert [p.name for p in tmp_path.iterdir()] == ["deck.pptx"]


class TestPathGate:
    def test_non_pptx_is_rejected(self, tmp_path):
        with pytest.raises(ValueError, match=r"\.pptx"):
            resolve_output_path(str(tmp_path / "deck.pdf"))

    def test_writes_outside_output_dir_are_blocked(self, tmp_path, monkeypatch):
        allowed = tmp_path / "allowed"; allowed.mkdir()
        monkeypatch.setenv("PPT_MCP_OUTPUT_DIR", str(allowed))
        resolve_output_path(str(allowed / "ok.pptx"))          # 안쪽은 통과
        with pytest.raises(ValueError, match="PPT_MCP_OUTPUT_DIR"):
            resolve_output_path(str(tmp_path / "escape.pptx"))  # 바깥은 차단

    def test_traversal_cannot_escape(self, tmp_path, monkeypatch):
        allowed = tmp_path / "allowed"; allowed.mkdir()
        monkeypatch.setenv("PPT_MCP_OUTPUT_DIR", str(allowed))
        with pytest.raises(ValueError, match="PPT_MCP_OUTPUT_DIR"):
            resolve_output_path(str(allowed / ".." / "escape.pptx"))


class TestSpecAndLayoutAgree:
    """스펙 제약과 레이아웃 용량이 서로 맞는지.

    models.py 의 상한(제목 80자, 불릿 6개 × 120자)과 layout.py 의 상자 크기는
    따로 정해집니다. 둘이 어긋나면 **검증을 통과한 스펙이 렌더에서 터집니다.**
    호출자 입장에서는 시키는 대로 했는데 실패하는 셈이라 최악입니다.

    이 테스트가 두 파일을 묶어 둡니다. 상자를 줄이거나 글자 수 상한을 올리면
    여기서 걸립니다.
    """

    def test_max_length_title_still_fits(self, tmp_path):
        spec = DeckSpec.model_validate({"theme": "minimal_dark", "slides": [
            {"kind": "title", "title": "가" * 80, "subtitle": "나" * 120},
            {"kind": "section", "title": "다" * 60},
            {"kind": "bullets", "title": "라" * 80,
             "points": ["마" * 120 for _ in range(6)]},
        ]})
        render_deck(spec, str(tmp_path / "max.pptx"))   # 예외가 나면 실패

    def test_max_chart_series_fits(self, tmp_path):
        spec = DeckSpec.model_validate({"theme": "minimal_light", "slides": [
            {"kind": "chart", "title": "바" * 80,
             "series": {f"{2015 + i}": i * 3.5 for i in range(12)}},
        ]})
        render_deck(spec, str(tmp_path / "chart.pptx"))
