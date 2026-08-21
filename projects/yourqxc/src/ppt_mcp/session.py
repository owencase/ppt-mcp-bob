"""열려 있는 덱을 들고 있는 세션 레지스트리.

슬라이드 단위 편집(add/update/delete/move)을 하려면 저장 전 상태를 서버가 기억해야
한다. 세션은 프로세스 안에서만 유지되며, save_deck을 불러야 디스크에 남는다.
"""

from __future__ import annotations

import time
from dataclasses import dataclass, field
from pathlib import Path
from uuid import uuid4

from pptx import Presentation

from .builder import DeckBuilder
from .config import DECK_SUFFIXES, Settings, resolve_read_path
from .models import DeckSpec


@dataclass
class DeckSession:
    deck_id: str
    builder: DeckBuilder
    origin: str
    saved_path: Path | None = None
    created_at: float = field(default_factory=time.time)
    updated_at: float = field(default_factory=time.time)
    _saved_at: float = field(default=0.0, repr=False)

    def touch(self) -> None:
        self.updated_at = time.time()

    @property
    def slide_count(self) -> int:
        return len(self.builder.prs.slides)

    def summary(self) -> dict:
        return {
            "deck_id": self.deck_id,
            "title": self.builder.spec.title,
            "slide_count": self.slide_count,
            "theme": self.builder.theme.name,
            "origin": self.origin,
            "saved_path": str(self.saved_path) if self.saved_path else None,
            "unsaved_changes": self.saved_path is None or self.updated_at > self._saved_at,
        }

    def mark_saved(self, path: Path) -> None:
        self.saved_path = path
        self._saved_at = time.time()


class SessionRegistry:
    """deck_id → DeckSession."""

    def __init__(self, settings: Settings, limit: int = 16) -> None:
        self.settings = settings
        self.limit = limit
        self._sessions: dict[str, DeckSession] = {}

    def _new_id(self) -> str:
        return f"deck_{uuid4().hex[:8]}"

    def _guard_capacity(self) -> None:
        if len(self._sessions) >= self.limit:
            raise RuntimeError(
                f"열어 둔 덱이 너무 많습니다({len(self._sessions)}/{self.limit}). "
                "close_deck으로 정리한 뒤 다시 시도하세요."
            )

    def create(self, spec: DeckSpec) -> DeckSession:
        """빈 덱(또는 스펙에 담긴 슬라이드까지)을 새로 연다."""
        self._guard_capacity()
        builder = DeckBuilder(spec, self.settings)
        for slide_spec in spec.slides:
            builder.add_slide(slide_spec)
        session = DeckSession(self._new_id(), builder, origin="new")
        self._sessions[session.deck_id] = session
        return session

    def open_file(self, path_ref: str, spec: DeckSpec) -> DeckSession:
        """기존 파일을 열어 이어서 편집한다. 파일의 레이아웃·마스터를 그대로 쓴다."""
        self._guard_capacity()
        path = resolve_read_path(path_ref, self.settings, suffixes=DECK_SUFFIXES)
        prs = Presentation(str(path))
        builder = DeckBuilder(spec, self.settings, presentation=prs, use_template=True)
        session = DeckSession(self._new_id(), builder, origin=str(path))
        self._sessions[session.deck_id] = session
        return session

    def get(self, deck_id: str) -> DeckSession:
        session = self._sessions.get(deck_id)
        if session is None:
            known = ", ".join(sorted(self._sessions)) or "없음"
            raise KeyError(f"열려 있는 덱이 아닙니다: {deck_id} (현재 열린 덱: {known})")
        return session

    def close(self, deck_id: str) -> DeckSession:
        session = self.get(deck_id)
        del self._sessions[deck_id]
        return session

    def list(self) -> list[dict]:
        return [s.summary() for s in sorted(self._sessions.values(), key=lambda x: x.created_at)]


__all__ = ["SessionRegistry", "DeckSession"]
