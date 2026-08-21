"""MCP 서버 정의 — 도구 11개.

크게 두 갈래로 쓴다.
  · 한 번에 만들기: create_deck에 덱 전체 스펙을 넘긴다. 대부분 이걸로 끝난다.
  · 고쳐 가며 만들기: open_deck → add/update/delete/move_slide → save_deck.
"""

from __future__ import annotations

from pydantic import BaseModel, Field

from mcp.server import MCPServer
from mcp.types import ToolAnnotations

from .builder import DeckBuilder
from .config import DECK_SUFFIXES, Settings, resolve_read_path, resolve_write_path
from .inspection import DeckReport, inspect_deck
from .models import SLIDE_TYPES, AnySlide, DeckSpec
from .session import SessionRegistry
from .template import TemplateMap
from .theme import describe_themes

INSTRUCTIONS = """\
PowerPoint(.pptx) 덱을 만들고 고치는 서버입니다.

기본 흐름
  1. 처음이면 describe_options로 테마·슬라이드 종류를 확인하세요.
  2. 덱 한 벌을 한 번에 만들 때는 create_deck에 slides 배열을 통째로 넘깁니다.
  3. 만든 뒤 일부만 손볼 때는 open_deck으로 열고 add/update/delete/move_slide를
     쓴 다음 반드시 save_deck을 호출하세요. save_deck 전에는 디스크에 아무것도
     쓰이지 않습니다.

좋은 덱을 만드는 요령
  · 한 장에 메시지 하나. 불릿은 6개 이하, 한 줄은 40자 안쪽이 읽기 좋습니다.
  · 표지(title) → 목차(agenda) → 간지(section)로 뼈대를 잡고 본문을 채우세요.
  · 숫자를 강조할 땐 bullets보다 kpi, 추세는 chart, 항목 대조는 comparison이 낫습니다.
  · 발표용이면 각 슬라이드 notes에 말할 내용을 적어 두세요.
  · 사내 서식이 있으면 template에 .potx/.pptx 경로를 주세요. 테마보다 우선합니다.

제약
  · 파일은 서버에 설정된 작업 디렉터리 안에만 저장됩니다.
  · 이미지 URL은 기본으로 차단돼 있습니다. 로컬 경로나 base64를 쓰세요.
"""

_SLIDE_TYPE_GUIDE: dict[str, str] = {
    "title": "표지. 덱 맨 앞에 한 장.",
    "agenda": "번호가 붙은 목차. 항목 4~7개.",
    "section": "장 구분 간지. 덱이 길 때 흐름을 끊어 준다.",
    "bullets": "제목 + 불릿. 가장 기본. 문자열 앞의 '-'로 하위 단계를 만든다.",
    "two_column": "좌우 2단. 항목을 나란히 놓을 때.",
    "comparison": "색이 구분된 카드 두 장. As-Is/To-Be, 장단점 대조에.",
    "table": "표. 행이 8개를 넘으면 두 장으로 나누는 편이 낫다.",
    "chart": "차트. 추세는 line, 항목 비교는 column, 구성비는 pie.",
    "image": "이미지. placement로 전체/좌/우/꽉채움을 고른다.",
    "quote": "인용문 한 장. 고객 목소리나 핵심 문장 강조에.",
    "kpi": "숫자 카드 2~4개. 지표를 크게 보여줄 때.",
    "timeline": "가로 타임라인. 로드맵이나 단계별 계획에.",
    "blank": "위치를 직접 지정하는 자유 배치. 다른 종류가 안 맞을 때만.",
}


# --- 도구 결과 모델 -------------------------------------------------------

class DeckResult(BaseModel):
    """덱 저장 결과."""

    path: str = Field(description="저장된 파일의 절대 경로.")
    slide_count: int = Field(description="슬라이드 수.")
    theme: str = Field(description="적용된 테마 이름.")
    used_template: bool = Field(description="템플릿 파일을 기반으로 만들었는지 여부.")
    slide_types: list[str] = Field(description="슬라이드 종류를 순서대로 나열.")
    warnings: list[str] = Field(
        default_factory=list,
        description="내용이 넘쳐 글자를 줄인 슬라이드 안내. 비어 있으면 문제 없다. "
                    "항목이 있으면 해당 슬라이드를 두 장으로 나누는 것을 검토할 것.")


class SessionResult(BaseModel):
    """열려 있는 덱의 현재 상태."""

    deck_id: str = Field(description="이후 편집 도구에 넘길 식별자.")
    title: str = Field(description="덱 제목.")
    slide_count: int = Field(description="현재 슬라이드 수.")
    theme: str = Field(description="적용된 테마 이름.")
    origin: str = Field(description="'new' 또는 열어 온 파일 경로.")
    saved_path: str | None = Field(None, description="마지막으로 저장된 경로. 저장 전이면 null.")
    message: str = Field(description="방금 수행한 작업 요약.")
    warnings: list[str] = Field(
        default_factory=list,
        description="방금 작업에서 나온 안내(내용 과다로 글자를 줄인 경우 등).")


class OptionsResult(BaseModel):
    """이 서버가 지원하는 선택지."""

    themes: list[dict] = Field(description="사용 가능한 테마와 대표 색.")
    slide_types: list[dict] = Field(description="슬라이드 종류와 쓰임새.")
    template_layouts: list[dict] | None = Field(
        None, description="template을 준 경우 그 파일의 레이아웃 목록.")
    output_dir: str = Field(description="파일이 저장되는 작업 디렉터리.")
    remote_images_allowed: bool = Field(description="이미지 URL 다운로드 허용 여부.")


class OpenDeckList(BaseModel):
    decks: list[dict] = Field(description="열려 있는 덱 목록.")


# --- 서버 조립 -------------------------------------------------------------

def build_server(settings: Settings | None = None) -> MCPServer:
    settings = settings or Settings.from_env()
    registry = SessionRegistry(settings)
    mcp = MCPServer(name="ppt-mcp", version="0.1.0", instructions=INSTRUCTIONS)

    def _result(builder: DeckBuilder, path) -> DeckResult:
        return DeckResult(
            path=str(path),
            slide_count=len(builder.prs.slides),
            theme=builder.theme.name,
            used_template=builder.uses_template,
            slide_types=[s.type for s in builder.spec.slides],
            warnings=list(builder.warnings),
        )

    def _session_result(session, message: str, since: int = 0) -> SessionResult:
        info = session.summary()
        return SessionResult(
            deck_id=info["deck_id"], title=info["title"], slide_count=info["slide_count"],
            theme=info["theme"], origin=info["origin"], saved_path=info["saved_path"],
            message=message, warnings=session.builder.warnings[since:],
        )

    # --- 탐색 -------------------------------------------------------------
    @mcp.tool(
        description="지원하는 테마, 슬라이드 종류와 쓰임새, (template을 주면) 그 파일의 "
                    "레이아웃 목록을 돌려준다. 덱을 만들기 전에 한 번 확인하면 좋다.",
        annotations=ToolAnnotations(title="선택지 확인", readOnlyHint=True, idempotentHint=True),
    )
    def describe_options(
        template: str | None = Field(
            None, description="레이아웃을 확인할 .potx/.pptx 경로. 생략 가능."),
    ) -> OptionsResult:
        layouts = None
        if template:
            from pptx import Presentation
            path = resolve_read_path(template, settings, suffixes=DECK_SUFFIXES)
            layouts = TemplateMap(Presentation(str(path))).describe()
        return OptionsResult(
            themes=describe_themes(),
            slide_types=[{"type": t, "usage": _SLIDE_TYPE_GUIDE[t]} for t in SLIDE_TYPES],
            template_layouts=layouts,
            output_dir=str(settings.output_dir),
            remote_images_allowed=settings.allow_remote_images,
        )

    @mcp.tool(
        description="기존 .pptx를 읽어 슬라이드별 제목·텍스트·발표자 노트와 표/차트/그림 "
                    "개수를 돌려준다. 덱을 요약하거나 고치기 전에 현재 내용을 파악할 때 쓴다.",
        annotations=ToolAnnotations(title="덱 읽기", readOnlyHint=True, idempotentHint=True),
    )
    def read_deck(
        path: str = Field(description="읽을 .pptx/.potx 경로."),
        include_text: bool = Field(True, description="본문 텍스트까지 포함할지. 개요만 필요하면 false."),
    ) -> DeckReport:
        return inspect_deck(path, settings, include_text=include_text)

    @mcp.tool(
        description="현재 열려 있는(아직 저장하지 않았을 수 있는) 덱 목록을 돌려준다.",
        annotations=ToolAnnotations(title="열린 덱 목록", readOnlyHint=True, idempotentHint=True),
    )
    def list_open_decks() -> OpenDeckList:
        return OpenDeckList(decks=registry.list())

    # --- 한 번에 만들기 ----------------------------------------------------
    @mcp.tool(
        description="덱 한 벌을 통째로 만들어 .pptx로 저장한다. 슬라이드 구성이 이미 정해졌다면 "
                    "이 도구 한 번이면 끝난다. 저장까지 한 번에 이뤄진다.",
        annotations=ToolAnnotations(title="덱 만들기", readOnlyHint=False, idempotentHint=False),
    )
    def create_deck(
        spec: DeckSpec = Field(description="덱 전체 스펙. slides 배열에 슬라이드를 순서대로 담는다."),
        output_path: str | None = Field(
            None, description="저장할 파일명 또는 경로. 생략하면 덱 제목으로 작업 디렉터리에 저장한다."),
    ) -> DeckResult:
        builder = DeckBuilder(spec, settings)
        builder.build()
        path = resolve_write_path(output_path, settings, default_stem=_safe_stem(spec.title))
        builder.prs.save(str(path))
        return _result(builder, path)

    # --- 고쳐 가며 만들기 --------------------------------------------------
    @mcp.tool(
        description="편집할 덱을 연다. path를 주면 기존 파일을 열어 그 서식을 이어 쓰고, "
                    "생략하면 새 덱을 시작한다. 반환된 deck_id로 슬라이드를 추가·수정한 뒤 "
                    "save_deck을 호출해야 파일로 남는다.",
        annotations=ToolAnnotations(title="덱 열기", readOnlyHint=False, idempotentHint=False),
    )
    def open_deck(
        path: str | None = Field(None, description="열 기존 .pptx 경로. 새로 시작하면 생략."),
        spec: DeckSpec | None = Field(
            None, description="새 덱의 초기 설정(제목·테마·템플릿 등). slides를 채우면 그것까지 만들어 둔다."),
    ) -> SessionResult:
        if spec is None:
            spec = DeckSpec(title="제목 없는 덱", slides=[])
        if path and "page_numbers" not in spec.model_fields_set:
            # 남의 덱에 이어 붙일 때 새 슬라이드에만 페이지 번호가 찍히면 어색하다.
            spec = spec.model_copy(update={"page_numbers": False})
        if path:
            session = registry.open_file(path, spec)
            message = f"'{session.origin}'을 열었습니다. 현재 {session.slide_count}장."
        else:
            session = registry.create(spec)
            message = f"새 덱을 열었습니다. 현재 {session.slide_count}장."
        return _session_result(session, message)

    @mcp.tool(
        description="열려 있는 덱에 슬라이드를 한 장 추가한다. index를 주면 그 위치에 끼워 넣는다.",
        annotations=ToolAnnotations(title="슬라이드 추가", readOnlyHint=False, idempotentHint=False),
    )
    def add_slide(
        deck_id: str = Field(description="open_deck이 돌려준 식별자."),
        slide: AnySlide = Field(description="추가할 슬라이드 하나."),
        index: int | None = Field(None, description="끼워 넣을 위치(0부터). 생략하면 맨 뒤."),
    ) -> SessionResult:
        session = registry.get(deck_id)
        seen = len(session.builder.warnings)
        session.builder.spec.slides.append(slide)
        session.builder.add_slide(slide, index=index)
        session.touch()
        where = "맨 뒤" if index is None else f"{index}번 자리"
        return _session_result(session, f"{where}에 {slide.type} 슬라이드를 추가했습니다.", seen)

    @mcp.tool(
        description="지정한 위치의 슬라이드를 새 내용으로 통째로 바꾼다. 일부만 고칠 때도 "
                    "바뀐 내용을 포함한 슬라이드 전체를 넘겨야 한다.",
        annotations=ToolAnnotations(title="슬라이드 교체", readOnlyHint=False, idempotentHint=True),
    )
    def update_slide(
        deck_id: str = Field(description="open_deck이 돌려준 식별자."),
        index: int = Field(description="바꿀 슬라이드 위치(0부터)."),
        slide: AnySlide = Field(description="그 자리에 들어갈 새 슬라이드."),
    ) -> SessionResult:
        session = registry.get(deck_id)
        seen = len(session.builder.warnings)
        session.builder.delete_slide(index)
        session.builder.add_slide(slide, index=index)
        if index < len(session.builder.spec.slides):
            session.builder.spec.slides[index] = slide
        session.touch()
        return _session_result(session, f"{index}번 슬라이드를 {slide.type}(으)로 바꿨습니다.", seen)

    @mcp.tool(
        description="지정한 위치의 슬라이드를 삭제한다. 되돌릴 수 없으니 위치를 먼저 확인할 것.",
        annotations=ToolAnnotations(title="슬라이드 삭제", readOnlyHint=False,
                                    destructiveHint=True, idempotentHint=False),
    )
    def delete_slide(
        deck_id: str = Field(description="open_deck이 돌려준 식별자."),
        index: int = Field(description="지울 슬라이드 위치(0부터)."),
    ) -> SessionResult:
        session = registry.get(deck_id)
        session.builder.delete_slide(index)
        if index < len(session.builder.spec.slides):
            del session.builder.spec.slides[index]
        session.touch()
        return _session_result(session, f"{index}번 슬라이드를 삭제했습니다.")

    @mcp.tool(
        description="슬라이드 순서를 바꾼다.",
        annotations=ToolAnnotations(title="슬라이드 이동", readOnlyHint=False, idempotentHint=False),
    )
    def move_slide(
        deck_id: str = Field(description="open_deck이 돌려준 식별자."),
        from_index: int = Field(description="옮길 슬라이드의 현재 위치(0부터)."),
        to_index: int = Field(description="옮겨 놓을 위치(0부터)."),
    ) -> SessionResult:
        session = registry.get(deck_id)
        session.builder.move_slide(from_index, to_index)
        slides = session.builder.spec.slides
        if 0 <= from_index < len(slides):
            slides.insert(max(0, min(to_index, len(slides) - 1)), slides.pop(from_index))
        session.touch()
        return _session_result(session, f"{from_index}번을 {to_index}번 자리로 옮겼습니다.")

    @mcp.tool(
        description="열려 있는 덱을 .pptx 파일로 저장한다. 편집 후 반드시 호출해야 한다. "
                    "덱은 계속 열려 있으므로 이어서 편집할 수 있다.",
        annotations=ToolAnnotations(title="덱 저장", readOnlyHint=False, idempotentHint=True),
    )
    def save_deck(
        deck_id: str = Field(description="open_deck이 돌려준 식별자."),
        output_path: str | None = Field(
            None, description="저장 경로. 생략하면 이전 저장 위치, 그것도 없으면 덱 제목으로 저장."),
    ) -> DeckResult:
        session = registry.get(deck_id)
        target = output_path or (str(session.saved_path) if session.saved_path else None)
        path = resolve_write_path(target, settings,
                                  default_stem=_safe_stem(session.builder.spec.title))
        session.builder._set_core_properties()
        session.builder.prs.save(str(path))
        session.mark_saved(path)
        return _result(session.builder, path)

    @mcp.tool(
        description="덱을 닫아 메모리에서 내린다. 저장하지 않은 변경은 사라진다.",
        annotations=ToolAnnotations(title="덱 닫기", readOnlyHint=False,
                                    destructiveHint=True, idempotentHint=False),
    )
    def close_deck(
        deck_id: str = Field(description="닫을 덱의 식별자."),
    ) -> SessionResult:
        session = registry.close(deck_id)
        saved = "저장됨" if session.saved_path else "저장하지 않고 닫음"
        return _session_result(session, f"덱을 닫았습니다({saved}).")

    return mcp


def _safe_stem(title: str) -> str:
    """제목을 파일명으로 쓸 수 있게 다듬는다."""
    cleaned = "".join(ch if ch.isalnum() or ch in " -_()." else "_" for ch in title).strip()
    cleaned = "_".join(cleaned.split())
    return (cleaned[:60] or "deck")


__all__ = ["build_server", "INSTRUCTIONS"]
