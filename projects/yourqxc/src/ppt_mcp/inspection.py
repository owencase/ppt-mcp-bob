"""기존 .pptx 읽기 — 요약, 텍스트 추출, 발표자 노트.

'이 덱 요약해 줘', '3번 슬라이드 문구 바꿔 줘' 같은 요청에서 에이전트가 먼저
현재 상태를 파악하는 데 쓴다.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pydantic import BaseModel, Field

from .config import DECK_SUFFIXES, Settings, resolve_read_path

EMU_PER_INCH = 914400.0


class SlideReport(BaseModel):
    """슬라이드 한 장의 요약."""

    index: int = Field(description="0부터 시작하는 슬라이드 번호.")
    layout: str = Field(description="사용된 레이아웃 이름.")
    title: str | None = Field(None, description="제목 플레이스홀더 텍스트(없으면 첫 텍스트 상자).")
    texts: list[str] = Field(default_factory=list, description="슬라이드 안의 모든 텍스트 덩어리.")
    notes: str | None = Field(None, description="발표자 노트.")
    tables: int = Field(0, description="표 개수.")
    charts: int = Field(0, description="차트 개수.")
    pictures: int = Field(0, description="그림 개수.")


class DeckReport(BaseModel):
    """덱 전체 요약."""

    path: str = Field(description="읽은 파일의 절대 경로.")
    slide_count: int = Field(description="슬라이드 수.")
    slide_size: str = Field(description="슬라이드 크기(인치). 예: '13.33 x 7.50'.")
    aspect: str = Field(description="가로세로 비율 추정값. 예: '16:9'.")
    title: str | None = Field(None, description="파일 속성의 제목.")
    author: str | None = Field(None, description="파일 속성의 작성자.")
    layouts: list[str] = Field(default_factory=list, description="이 파일에서 쓸 수 있는 레이아웃 이름들.")
    slides: list[SlideReport] = Field(default_factory=list, description="슬라이드별 요약.")


def _shape_text(shape) -> str:
    if not shape.has_text_frame:
        return ""
    return "\n".join(
        "".join(run.text for run in para.runs) for para in shape.text_frame.paragraphs
    ).strip()


def _has_letters(text: str) -> bool:
    """따옴표·불릿 기호만 있는 텍스트를 제목 후보에서 걸러낸다."""
    return any(ch.isalnum() for ch in text)


def _largest_text(slide, slide_height: float) -> str | None:
    """제목 플레이스홀더가 없을 때 제목으로 볼 만한 텍스트.

    화면 위쪽 고정 구간으로 자르면 간지처럼 제목이 가운데 있는 슬라이드에서 틀린다.
    그래서 '가장 위에 있는 글자'를 기준으로 1.2인치 안쪽만 후보로 삼고, 그중 가장 큰
    글자를 고른다. KPI 숫자나 본문처럼 아래쪽에 있는 큰 글자는 자연스럽게 빠진다.
    """
    candidates: list[tuple[float, float, str]] = []  # (top, size, text)
    for shape in slide.shapes:
        if not shape.has_text_frame or shape.top is None:
            continue
        top = shape.top / EMU_PER_INCH
        for para in shape.text_frame.paragraphs:
            text = "".join(run.text for run in para.runs).strip()
            if not text or not _has_letters(text):
                continue
            size = max((run.font.size.pt for run in para.runs if run.font.size), default=0.0)
            candidates.append((top, size, text))
    if not candidates:
        return None
    topmost = min(c[0] for c in candidates)
    band = [c for c in candidates if c[0] <= topmost + 1.2]
    return max(band, key=lambda c: (c[1], -c[0]))[2]


def _guess_aspect(width_in: float, height_in: float) -> str:
    ratio = width_in / height_in if height_in else 0
    for label, value in (("16:9", 16 / 9), ("16:10", 1.6), ("4:3", 4 / 3), ("A4", 1.414)):
        if abs(ratio - value) < 0.03:
            return label
    return f"{ratio:.2f}:1"


def inspect_deck(path_ref: str, settings: Settings, *, include_text: bool = True) -> DeckReport:
    path: Path = resolve_read_path(path_ref, settings, suffixes=DECK_SUFFIXES)
    prs = Presentation(str(path))
    width = prs.slide_width / EMU_PER_INCH
    height = prs.slide_height / EMU_PER_INCH

    slides: list[SlideReport] = []
    for index, slide in enumerate(prs.slides):
        texts, tables, charts, pictures = [], 0, 0, 0
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                tables += 1
            if getattr(shape, "has_chart", False):
                charts += 1
            if shape.shape_type is not None and "PICTURE" in str(shape.shape_type):
                pictures += 1
            text = _shape_text(shape)
            if text:
                texts.append(text)

        title = None
        if slide.shapes.title is not None:
            title = _shape_text(slide.shapes.title) or None
        if title is None:
            # 제목 플레이스홀더가 없으면 가장 큰 글자를 제목으로 본다
            # (표지의 '주간 보고' 같은 작은 라벨이 먼저 잡히는 걸 막는다).
            title = _largest_text(slide, height)

        notes = None
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip() or None

        slides.append(SlideReport(
            index=index, layout=slide.slide_layout.name, title=title,
            texts=texts if include_text else [], notes=notes,
            tables=tables, charts=charts, pictures=pictures,
        ))

    props = prs.core_properties
    return DeckReport(
        path=str(path),
        slide_count=len(prs.slides),
        slide_size=f"{width:.2f} x {height:.2f}",
        aspect=_guess_aspect(width, height),
        title=props.title or None,
        author=props.author or None,
        layouts=[layout.name for layout in prs.slide_layouts],
        slides=slides,
    )


__all__ = ["inspect_deck", "DeckReport", "SlideReport"]
