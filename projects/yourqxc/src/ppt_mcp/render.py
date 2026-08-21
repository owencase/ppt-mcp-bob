"""python-pptx 위에 올리는 저수준 그리기 헬퍼.

여기서 신경 쓰는 것 두 가지:
  1. 한글 폰트 — run.font.name은 <a:latin>만 건드리므로, 한글이 제대로 나오려면
     <a:ea>를 직접 넣어야 한다.
  2. 도형 기본 그림자 제거 — 테마 기본값이 남으면 납작한 디자인이 지저분해진다.
"""

from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from .models import Bullet, BulletLike
from .theme import Theme, hex_to_rgb

_ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
_ANCHOR = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}


def resolve_color(theme: Theme, token: str):
    """색 토큰 또는 hex 문자열 → RGBColor."""
    mapping = {
        "default": theme.text,
        "muted": theme.text_muted,
        "accent": theme.accent,
        "accent_alt": theme.accent_alt,
        "positive": theme.positive,
        "negative": theme.negative,
        "surface": theme.surface,
        "background": theme.background,
        "divider": theme.divider,
    }
    return hex_to_rgb(mapping.get(token, token))


def apply_typeface(run, theme: Theme) -> None:
    """<a:latin>·<a:ea>·<a:cs>를 모두 지정해 한글이 대체 폰트로 새지 않게 한다."""
    run.font.name = theme.font_latin  # 스키마상 올바른 위치에 <a:latin> 생성
    rPr = run._r.get_or_add_rPr()
    prev = rPr.find(qn("a:latin"))
    if prev is None:
        return
    for tag, face in (("a:ea", theme.font_ea), ("a:cs", theme.font_latin)):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            prev.addnext(el)  # latin → ea → cs 순서 유지
        el.set("typeface", face)
        prev = el


def style_run(run, theme: Theme, *, size: int, bold: bool = False,
              color: str = "default", italic: bool = False) -> None:
    font = run.font
    font.size = Pt(size)
    font.bold = bold
    font.italic = italic
    font.color.rgb = resolve_color(theme, color)
    apply_typeface(run, theme)


def add_textbox(slide, x: float, y: float, w: float, h: float, *,
                anchor: str = "top", wrap: bool = True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = _ANCHOR[anchor]
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    return box


def write_paragraph(tf, text: str, theme: Theme, *, size: int, bold: bool = False,
                    color: str = "default", align: str = "left", level: int = 0,
                    italic: bool = False, space_after: float = 6.0,
                    line_spacing: float = 1.18, first: bool = False):
    """텍스트 프레임에 문단 하나를 쓴다. first=True면 기본 빈 문단을 재사용."""
    para = tf.paragraphs[0] if first else tf.add_paragraph()
    para.level = min(level, 8)
    para.alignment = _ALIGN[align]
    para.space_after = Pt(space_after)
    para.line_spacing = line_spacing
    run = para.add_run()
    run.text = text
    style_run(run, theme, size=size, bold=bold, color=color, italic=italic)
    return para


def enable_shrink_on_overflow(tf, font_scale: int = 90, line_reduction: int = 10) -> None:
    """PowerPoint 자체 '넘치면 줄이기'를 안전망으로 켜 둔다."""
    bodyPr = tf._txBody.bodyPr
    for tag in ("a:normAutofit", "a:spAutoFit", "a:noAutofit"):
        existing = bodyPr.find(qn(tag))
        if existing is not None:
            bodyPr.remove(existing)
    el = bodyPr.makeelement(qn("a:normAutofit"), {})
    if font_scale < 100:  # 100%면 속성 없이 두고 PowerPoint가 알아서 판단하게 한다
        el.set("fontScale", str(font_scale * 1000))
    if line_reduction:
        el.set("lnSpcReduction", str(line_reduction * 1000))
    bodyPr.append(el)


def add_rect(slide, x: float, y: float, w: float, h: float, *,
             fill: str | None = None, line: str | None = None,
             line_width: float = 1.0, theme: Theme, rounded: bool = False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = resolve_color(theme, fill)
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = resolve_color(theme, line)
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    if shape.has_text_frame:
        shape.text_frame.word_wrap = True
        shape.text_frame.margin_left = Inches(0.22)
        shape.text_frame.margin_right = Inches(0.22)
        shape.text_frame.margin_top = Inches(0.14)
        shape.text_frame.margin_bottom = Inches(0.14)
    return shape


def paint_background(slide, theme: Theme) -> None:
    """슬라이드 배경을 테마 색으로 채운다(템플릿 사용 시에는 호출하지 않는다)."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(theme.background)


def set_notes(slide, text: str | None) -> None:
    if not text:
        return
    slide.notes_slide.notes_text_frame.text = text


def normalize_bullets(items: list[BulletLike] | None) -> list[Bullet]:
    """문자열/객체가 섞인 불릿 목록을 Bullet 객체로 정규화.

    문자열은 앞의 '-'·'*'·'•' 개수나 두 칸 들여쓰기로 단계를 읽는다.
    예: '- 하위 항목' → level 1, '-- 더 하위' → level 2.
    """
    result: list[Bullet] = []
    for item in items or []:
        if isinstance(item, Bullet):
            result.append(item)
            continue
        text = str(item)
        level = 0
        stripped = text.lstrip()
        leading_spaces = len(text) - len(stripped)
        level += leading_spaces // 2
        while stripped[:1] in {"-", "*", "•"}:
            marker_run = len(stripped) - len(stripped.lstrip("-*•"))
            if marker_run == 0:
                break
            level += marker_run
            stripped = stripped[marker_run:].lstrip()
        if not stripped:
            continue
        result.append(Bullet(text=stripped, level=min(level, 3)))
    return result


__all__ = ["resolve_color", "apply_typeface", "style_run", "add_textbox",
           "write_paragraph", "add_rect", "paint_background",
           "set_notes", "normalize_bullets", "enable_shrink_on_overflow", "Inches", "Pt"]
