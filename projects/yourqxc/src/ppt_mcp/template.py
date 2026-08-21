"""사내 템플릿(.potx/.pptx)의 레이아웃을 슬라이드 종류에 매핑.

템플릿을 쓰면 마스터에 있는 로고·배경·서식이 그대로 살아난다. 슬라이드 종류에
맞는 레이아웃을 이름으로 추정하고, 쓰지 않은 빈 플레이스홀더는 지워서 편집 화면에
'제목을 입력하십시오' 같은 안내문이 남지 않게 한다.
"""

from __future__ import annotations

from dataclasses import dataclass

from pptx.enum.shapes import PP_PLACEHOLDER

#: 슬라이드 종류별로 레이아웃 이름에서 찾을 키워드(앞에 있을수록 우선).
_LAYOUT_KEYWORDS: dict[str, tuple[str, ...]] = {
    "title": ("title slide", "제목 슬라이드", "표지", "cover", "title"),
    "section": ("section header", "구역 머리글", "간지", "section", "divider"),
    "agenda": ("agenda", "목차", "contents", "title and content", "제목 및 내용"),
    "bullets": ("title and content", "제목 및 내용", "content", "본문", "body"),
    "two_column": ("two content", "콘텐츠 2개", "comparison", "비교", "2단"),
    "comparison": ("comparison", "비교", "two content", "콘텐츠 2개"),
    "quote": ("quote", "인용", "title only", "제목만"),
    "image": ("picture", "그림", "content with caption", "title only", "제목만"),
    "table": ("title and content", "제목 및 내용", "title only", "제목만"),
    "chart": ("title and content", "제목 및 내용", "title only", "제목만"),
    "kpi": ("title only", "제목만", "title and content", "제목 및 내용"),
    "timeline": ("title only", "제목만", "title and content", "제목 및 내용"),
    "blank": ("blank", "빈 화면", "빈"),
}

_TITLE_TYPES = {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}
_BODY_TYPES = {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT}


@dataclass(frozen=True)
class LayoutInfo:
    index: int
    name: str
    placeholders: list[str]


class TemplateMap:
    """프레젠테이션의 레이아웃 목록을 슬라이드 종류로 이어 주는 어댑터."""

    def __init__(self, presentation) -> None:
        self.prs = presentation
        self._layouts = list(presentation.slide_layouts)

    # --- 조회 -----------------------------------------------------------
    def describe(self) -> list[dict]:
        out = []
        for i, layout in enumerate(self._layouts):
            names = []
            for ph in layout.placeholders:
                try:
                    kind = str(ph.placeholder_format.type).split(" ")[0]
                except Exception:  # 일부 템플릿은 type이 비어 있다
                    kind = "UNKNOWN"
                names.append(f"{kind}(idx={ph.placeholder_format.idx})")
            out.append({"index": i, "name": layout.name, "placeholders": names})
        return out

    def by_name(self, name: str | None):
        if not name:
            return None
        target = name.strip().lower()
        for layout in self._layouts:
            if layout.name.strip().lower() == target:
                return layout
        for layout in self._layouts:  # 부분 일치까지 허용
            if target in layout.name.strip().lower():
                return layout
        return None

    def blank(self):
        """플레이스홀더가 가장 적은 레이아웃 = 사실상 빈 레이아웃."""
        named = self.by_name("blank") or self.by_name("빈")
        if named is not None:
            return named
        return min(self._layouts, key=lambda l: len(l.placeholders))

    def for_type(self, slide_type: str, override: str | None = None):
        forced = self.by_name(override)
        if forced is not None:
            return forced
        for keyword in _LAYOUT_KEYWORDS.get(slide_type, ()):
            match = self.by_name(keyword)
            if match is not None:
                return match
        return self.blank()


# --- 슬라이드 위 플레이스홀더 조작 --------------------------------------

def find_placeholder(slide, kinds: set) -> object | None:
    for ph in slide.placeholders:
        try:
            if ph.placeholder_format.type in kinds:
                return ph
        except Exception:
            continue
    return None


def title_placeholder(slide):
    return find_placeholder(slide, _TITLE_TYPES)


def body_placeholder(slide):
    return find_placeholder(slide, _BODY_TYPES)


def subtitle_placeholder(slide):
    return find_placeholder(slide, {PP_PLACEHOLDER.SUBTITLE})


def drop_empty_placeholders(slide) -> int:
    """내용이 채워지지 않은 플레이스홀더를 슬라이드에서 제거."""
    removed = 0
    for ph in list(slide.placeholders):
        has_text = ph.has_text_frame and ph.text_frame.text.strip()
        if has_text:
            continue
        # 그림·표·차트가 들어간 플레이스홀더는 건드리지 않는다
        if getattr(ph, "shape_type", None) is not None and not ph.has_text_frame:
            continue
        ph._element.getparent().remove(ph._element)
        removed += 1
    return removed


__all__ = ["TemplateMap", "LayoutInfo", "title_placeholder", "body_placeholder",
           "subtitle_placeholder", "drop_empty_placeholders", "find_placeholder"]
