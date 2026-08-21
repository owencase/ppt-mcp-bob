"""DeckSpec → python-pptx Presentation.

테마 모드(템플릿 없음)에서는 빈 레이아웃 위에 직접 도형을 그려 일관된 결과를 낸다.
템플릿 모드에서는 종류에 맞는 레이아웃을 골라 제목/본문 플레이스홀더를 채우고,
표·차트·KPI처럼 플레이스홀더로 표현할 수 없는 것만 도형으로 얹는다.
"""

from __future__ import annotations

from dataclasses import replace

from pptx import Presentation
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from . import models as M
from .assets import load_image
from .config import DECK_SUFFIXES, Settings, resolve_read_path
from .render import (add_rect, add_textbox, enable_shrink_on_overflow, normalize_bullets,
                     paint_background, resolve_color, set_notes, style_run, write_paragraph)
from .template import TemplateMap, body_placeholder, drop_empty_placeholders, subtitle_placeholder, title_placeholder
from .textfit import estimate_block_height_pt, fit_font_size, fit_single_line
from .theme import Geometry, Theme, get_theme, hex_to_rgb

_CHART_TYPES = {
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
    "area": XL_CHART_TYPE.AREA,
}

_BULLET_GLYPH = {0: "■", 1: "–", 2: "·", 3: "·"}


class DeckBuilder:
    """스펙 한 벌을 받아 Presentation을 만든다. 슬라이드 단위 추가에도 재사용된다."""

    def __init__(self, spec: M.DeckSpec, settings: Settings, presentation=None,
                 use_template: bool = True) -> None:
        """presentation을 주면 이미 열려 있는 덱에 이어 그린다(슬라이드 단위 편집).

        use_template=True면 그 파일의 레이아웃·마스터를 따르고, False면 테마로 직접 그린다.
        기존 파일을 여는 경우는 원래 디자인에 맞추는 편이 자연스러워 True가 기본이다.
        """
        self.spec = spec
        self.settings = settings
        self.theme = self._resolve_theme(spec)
        if presentation is not None:
            self.prs = presentation
            self.template = TemplateMap(presentation) if use_template else None
        else:
            self.prs, self.template = self._open_presentation(spec)
        self.geo = self._derive_geometry()
        self.warnings: list[str] = []

    # --- 준비 -----------------------------------------------------------
    def _resolve_theme(self, spec: M.DeckSpec) -> Theme:
        theme = get_theme(spec.theme or self.settings.default_theme)
        theme = theme.with_overrides(spec.theme_overrides)
        return theme.with_fonts(spec.font_latin, spec.font_ea)

    def _open_presentation(self, spec: M.DeckSpec):
        template_ref = spec.template or (
            str(self.settings.default_template) if self.settings.default_template else None)
        if template_ref:
            path = resolve_read_path(template_ref, self.settings, suffixes=DECK_SUFFIXES)
            prs = Presentation(str(path))
            for slide in list(prs.slides):  # 템플릿에 딸린 예시 슬라이드는 비운다
                rid = prs.slides._sldIdLst[-1].rId
                prs.part.drop_rel(rid)
                del prs.slides._sldIdLst[-1]
            return prs, TemplateMap(prs)

        prs = Presentation()
        if spec.aspect == "16:9":
            prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
        else:
            prs.slide_width, prs.slide_height = Inches(10), Inches(7.5)
        return prs, None

    def _derive_geometry(self) -> Geometry:
        width = self.prs.slide_width / 914400
        height = self.prs.slide_height / 914400
        base = self.theme.geometry
        margin = round(base.margin_x * width / 13.333, 3)
        return replace(base, slide_w=round(width, 3), slide_h=round(height, 3), margin_x=margin)

    @property
    def uses_template(self) -> bool:
        return self.template is not None

    # --- 진입점 ---------------------------------------------------------
    def build(self) -> Presentation:
        if len(self.spec.slides) > self.settings.max_slides:
            raise ValueError(
                f"슬라이드가 너무 많습니다({len(self.spec.slides)}). 최대 {self.settings.max_slides}장.")
        for slide_spec in self.spec.slides:
            self.add_slide(slide_spec)
        self._set_core_properties()
        return self.prs

    def add_slide(self, spec, index: int | None = None):
        """슬라이드 한 장을 그린다. index를 주면 그 위치로 옮긴다."""
        renderer = getattr(self, f"_render_{spec.type}", None)
        if renderer is None:
            raise ValueError(f"지원하지 않는 슬라이드 종류입니다: {spec.type}")

        layout = (self.template.for_type(spec.type, spec.layout) if self.uses_template
                  else self.prs.slide_layouts[6])
        slide = self.prs.slides.add_slide(layout)
        if not self.uses_template:
            paint_background(slide, self.theme)

        renderer(slide, spec)
        self._add_chrome(slide, spec)
        set_notes(slide, spec.notes)
        if self.uses_template:
            drop_empty_placeholders(slide)
        if index is not None:
            self.move_slide(len(self.prs.slides) - 1, index)
            return self.prs.slides[index]
        self._renumber_slide_parts()
        return slide

    def _renumber_slide_parts(self) -> None:
        """슬라이드 파트 이름을 현재 순서대로 slide1..slideN으로 다시 매긴다.

        python-pptx는 새 슬라이드 파트 이름을 sldIdLst 길이 + 1로 정하기 때문에,
        중간을 지운 뒤 추가하면 기존 파트와 이름이 겹쳐 저장 시 zip 항목이 중복된다.
        (prs.slides가 이 정리를 해 주지만 lazyproperty라 최초 1회만 실행된다.)
        """
        id_list = self.prs.slides._sldIdLst
        self.prs.part.rename_slide_parts([sldId.rId for sldId in id_list])

    def delete_slide(self, index: int) -> None:
        """슬라이드를 지우고 딸린 관계(rel)도 정리한다."""
        id_list = self.prs.slides._sldIdLst
        items = list(id_list)
        if not 0 <= index < len(items):
            raise IndexError(f"슬라이드 {index}가 없습니다(총 {len(items)}장).")
        element = items[index]
        self.prs.part.drop_rel(element.rId)
        id_list.remove(element)
        self._renumber_slide_parts()

    def move_slide(self, source: int, target: int) -> None:
        id_list = self.prs.slides._sldIdLst
        items = list(id_list)
        if not 0 <= source < len(items):
            raise IndexError(f"슬라이드 {source}가 없습니다(총 {len(items)}장).")
        target = max(0, min(target, len(items) - 1))
        element = items[source]
        id_list.remove(element)
        id_list.insert(target, element)
        self._renumber_slide_parts()

    def _set_core_properties(self) -> None:
        props = self.prs.core_properties
        props.title = self.spec.title
        if self.spec.subtitle:
            props.subject = self.spec.subtitle
        if self.spec.author:
            props.author = self.spec.author

    # --- 공통 조각 ------------------------------------------------------
    def _title_block(self, slide, title: str | None, subtitle: str | None = None) -> float:
        """제목(+부제)을 얹고, 본문이 시작될 y 좌표를 돌려준다."""
        geo, theme = self.geo, self.theme
        content_top = geo.content_top

        if title:
            placeholder = title_placeholder(slide) if self.uses_template else None
            if placeholder is not None:
                placeholder.text_frame.text = title
                for para in placeholder.text_frame.paragraphs:
                    for run in para.runs:
                        style_run(run, theme, size=theme.size_slide_title, bold=True)
                content_top = max(
                    content_top,
                    (placeholder.top + placeholder.height) / 914400 + 0.28,
                )
            else:
                size = fit_single_line(title, geo.content_w,
                                       base_pt=theme.size_slide_title, min_pt=20)
                box = add_textbox(slide, geo.margin_x, geo.title_top, geo.content_w, 0.82,
                                  anchor="top")
                write_paragraph(box.text_frame, title, theme, size=size, bold=True,
                                space_after=0, first=True)
                add_rect(slide, geo.margin_x, geo.title_top + 0.9, 0.86, 0.055,
                         fill="accent", theme=theme)
                content_top = geo.content_top

        if subtitle:
            box = add_textbox(slide, geo.margin_x, content_top - 0.06, geo.content_w, 0.4)
            write_paragraph(box.text_frame, subtitle, theme, size=theme.size_small,
                            color="muted", space_after=0, first=True)
            content_top += 0.46
        return content_top

    def _content_height(self, top: float) -> float:
        return max(0.6, self.geo.slide_h - self.geo.margin_bottom - top)

    def _add_chrome(self, slide, spec) -> None:
        """푸터와 페이지 번호. 표지·간지에는 넣지 않는다."""
        if spec.type in {"title", "section"}:
            return
        geo, theme = self.geo, self.theme
        y = geo.slide_h - 0.52
        if self.spec.footer:
            box = add_textbox(slide, geo.margin_x, y, geo.content_w * 0.7, 0.3)
            write_paragraph(box.text_frame, self.spec.footer, theme,
                            size=theme.size_caption, color="muted", space_after=0, first=True)
        if self.spec.page_numbers:
            number = len(self.prs.slides)
            box = add_textbox(slide, geo.slide_w - geo.margin_x - 1.0, y, 1.0, 0.3)
            write_paragraph(box.text_frame, str(number), theme, size=theme.size_caption,
                            color="muted", align="right", space_after=0, first=True)

    def _write_bullets(self, slide, bullets, top: float, left: float, width: float,
                       height: float, *, base_size: int | None = None, min_size: int = 11) -> None:
        items = normalize_bullets(bullets)
        if not items:
            return
        theme = self.theme
        base = base_size or theme.size_body
        size = fit_font_size([(b.text, b.level) for b in items], width, height,
                             base_pt=base, min_pt=min_size)
        box = add_textbox(slide, left, top, width, height)
        tf = box.text_frame
        for i, bullet in enumerate(items):
            text = f"{_BULLET_GLYPH.get(bullet.level, '·')}  {bullet.text}"
            para = write_paragraph(
                tf, text, theme,
                size=max(min_size, size - bullet.level),
                bold=bullet.bold,
                color=bullet.color if bullet.color != "default" else ("default" if bullet.level == 0 else "muted"),
                level=0, space_after=6.0 if bullet.level == 0 else 3.0, first=i == 0,
            )
            self._indent(para, bullet.level)
        self._autofit(tf, [(b.text, b.level) for b in items], size, width, height)

    def _autofit(self, tf, entries, size: int, box_w: float, box_h: float) -> None:
        """최소 글자 크기로도 넘치면 PowerPoint 자동 축소 비율을 직접 계산해 넣는다.

        normAutofit의 기본값을 그냥 켜 두면 PowerPoint가 파일을 열 때 저장된 비율을
        그대로 쓰기 때문에, 실제로 필요한 비율을 우리가 계산해 줘야 넘치지 않는다.
        내용이 그만큼 많다는 뜻이므로 경고도 함께 남긴다.
        """
        available_pt = box_h * 72.0
        if not entries or estimate_block_height_pt(entries, box_w, size) <= available_pt:
            enable_shrink_on_overflow(tf, font_scale=100, line_reduction=0)
            return

        # 비율을 선형으로 계산하면 안 된다. 문단 사이 간격은 글자와 함께 줄지 않고,
        # 글자가 작아지면 줄바꿈도 다시 일어난다. 실제로 들어가는 비율을 찾아 내려간다.
        scale = 40
        for candidate in range(95, 35, -5):
            if estimate_block_height_pt(entries, box_w, size * candidate / 100) <= available_pt:
                scale = candidate
                break
        enable_shrink_on_overflow(tf, font_scale=scale, line_reduction=10)
        self.warnings.append(
            f"{len(self.prs.slides)}번 슬라이드: 내용이 많아 글자를 {scale}%까지 줄였습니다. "
            "두 장으로 나누면 읽기 좋습니다."
        )

    @staticmethod
    def _indent(para, level: int) -> None:
        """글머리 기호와 본문이 어긋나지 않도록 들여쓰기와 매달린 들여쓰기를 준다."""
        pPr = para._p.get_or_add_pPr()
        pPr.set("marL", str(int((0.34 * level + 0.26) * 914400)))
        pPr.set("indent", str(int(-0.26 * 914400)))

    # --- 슬라이드 렌더러 -------------------------------------------------
    def _render_title(self, slide, spec: M.TitleSlide) -> None:
        geo, theme = self.geo, self.theme
        if self.uses_template:
            ph = title_placeholder(slide)
            if ph is not None:
                ph.text_frame.text = spec.title
                for run in ph.text_frame.paragraphs[0].runs:
                    style_run(run, theme, size=theme.size_cover_title, bold=True)
                sub = subtitle_placeholder(slide) or body_placeholder(slide)
                tail = " · ".join(x for x in (spec.subtitle, spec.presenter, spec.date) if x)
                if sub is not None and tail:
                    sub.text_frame.text = tail
                    for run in sub.text_frame.paragraphs[0].runs:
                        style_run(run, theme, size=theme.size_cover_subtitle, color="muted")
                return

        top = geo.slide_h * 0.28
        if spec.eyebrow:
            box = add_textbox(slide, geo.margin_x, top - 0.75, geo.content_w, 0.35)
            write_paragraph(box.text_frame, spec.eyebrow, theme, size=theme.size_small,
                            bold=True, color="accent", space_after=0, first=True)
        add_rect(slide, geo.margin_x, top - 0.24, 1.7, 0.09, fill="accent", theme=theme)

        size = fit_font_size([(spec.title, 0)], geo.content_w, 1.9,
                             base_pt=theme.size_cover_title, min_pt=26)
        box = add_textbox(slide, geo.margin_x, top + 0.1, geo.content_w, 1.9)
        write_paragraph(box.text_frame, spec.title, theme, size=size, bold=True,
                        space_after=0, line_spacing=1.1, first=True)

        if spec.subtitle:
            box = add_textbox(slide, geo.margin_x, top + 2.05, geo.content_w * 0.85, 0.95)
            write_paragraph(box.text_frame, spec.subtitle, theme,
                            size=theme.size_cover_subtitle, color="muted", first=True)

        tail = " · ".join(x for x in (spec.presenter, spec.date) if x)
        if tail:
            box = add_textbox(slide, geo.margin_x, geo.slide_h - 1.0, geo.content_w, 0.4)
            write_paragraph(box.text_frame, tail, theme, size=theme.size_small,
                            color="muted", space_after=0, first=True)

    def _render_section(self, slide, spec: M.SectionSlide) -> None:
        geo, theme = self.geo, self.theme
        if self.uses_template:
            ph = title_placeholder(slide)
            if ph is not None:  # 템플릿의 간지 서식을 그대로 살린다
                ph.text_frame.text = spec.title
                for run in ph.text_frame.paragraphs[0].runs:
                    style_run(run, theme, size=theme.size_section_title, bold=True)
                body = body_placeholder(slide) or subtitle_placeholder(slide)
                tail = " · ".join(x for x in (spec.number, spec.subtitle) if x)
                if body is not None and tail:
                    body.text_frame.text = tail
                    for run in body.text_frame.paragraphs[0].runs:
                        style_run(run, theme, size=theme.size_body, color="muted")
                return

        if not self.uses_template:
            add_rect(slide, 0, 0, geo.slide_w, geo.slide_h, fill="surface", theme=theme)
            add_rect(slide, 0, 0, 0.28, geo.slide_h, fill="accent", theme=theme)

        top = geo.slide_h * 0.36
        left = geo.margin_x + 0.2
        if spec.number:
            box = add_textbox(slide, left, top - 0.7, 3.0, 0.5)
            write_paragraph(box.text_frame, spec.number, theme, size=theme.size_heading,
                            bold=True, color="accent", space_after=0, first=True)
        size = fit_single_line(spec.title, geo.content_w - 0.4,
                               base_pt=theme.size_section_title, min_pt=22)
        box = add_textbox(slide, left, top, geo.content_w - 0.4, 1.2)
        write_paragraph(box.text_frame, spec.title, theme, size=size, bold=True,
                        space_after=0, first=True)
        if spec.subtitle:
            box = add_textbox(slide, left, top + 1.15, geo.content_w * 0.8, 0.8)
            write_paragraph(box.text_frame, spec.subtitle, theme, size=theme.size_body,
                            color="muted", first=True)

    def _render_agenda(self, slide, spec: M.AgendaSlide) -> None:
        geo, theme = self.geo, self.theme
        top = self._title_block(slide, spec.title)
        height = self._content_height(top)
        items = spec.items
        columns = 2 if len(items) > 5 else 1
        per_column = (len(items) + columns - 1) // columns
        col_w = (geo.content_w - geo.gutter) / columns if columns > 1 else geo.content_w
        row_h = min(0.82, height / max(1, per_column))
        size = min(theme.size_body, int(row_h * 72 * 0.42))

        for i, text in enumerate(items):
            col, row = divmod(i, per_column)
            x = geo.margin_x + col * (col_w + geo.gutter)
            y = top + row * row_h
            num = add_textbox(slide, x, y, 0.6, row_h * 0.9)
            write_paragraph(num.text_frame, f"{i + 1:02d}", theme, size=size, bold=True,
                            color="accent", space_after=0, first=True)
            box = add_textbox(slide, x + 0.62, y, col_w - 0.62, row_h * 0.9)
            write_paragraph(box.text_frame, text, theme, size=size, space_after=0, first=True)
            add_rect(slide, x, y + row_h * 0.92, col_w, 0.012, fill="divider", theme=theme)

    def _render_bullets(self, slide, spec: M.BulletsSlide) -> None:
        top = self._title_block(slide, spec.title, spec.subtitle)
        geo = self.geo
        body = body_placeholder(slide) if self.uses_template else None
        if body is not None:
            self._fill_body_placeholder(body, spec.bullets)
            return
        self._write_bullets(slide, spec.bullets, top, geo.margin_x, geo.content_w,
                            self._content_height(top))

    def _fill_body_placeholder(self, placeholder, bullets) -> None:
        """템플릿 본문 플레이스홀더에 불릿을 채운다(템플릿의 글머리 서식을 그대로 씀)."""
        items = normalize_bullets(bullets)
        tf = placeholder.text_frame
        tf.clear()
        size = self.theme.size_body
        for i, bullet in enumerate(items):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.level = bullet.level
            run = para.add_run()
            run.text = bullet.text
            style_run(run, self.theme, size=max(12, size - bullet.level),
                      bold=bullet.bold, color=bullet.color)
        # 플레이스홀더 크기는 템플릿이 정하므로, 그 실제 치수를 기준으로 축소율을 잡는다.
        box_w = (placeholder.width or 0) / 914400 or self.geo.content_w
        box_h = (placeholder.height or 0) / 914400 or self._content_height(self.geo.content_top)
        self._autofit(tf, [(b.text, b.level) for b in items], size, box_w, box_h)

    def _render_two_column(self, slide, spec: M.TwoColumnSlide) -> None:
        top = self._title_block(slide, spec.title, spec.subtitle)
        geo = self.geo
        height = self._content_height(top)
        for index, column in enumerate((spec.left, spec.right)):
            x = geo.margin_x + index * (geo.column_w + geo.gutter)
            self._render_column(slide, column, x, top, geo.column_w, height)

    def _render_column(self, slide, column: M.Column, x: float, y: float,
                       width: float, height: float) -> None:
        theme = self.theme
        cursor = y
        if column.heading:
            box = add_textbox(slide, x, cursor, width, 0.42)
            write_paragraph(box.text_frame, column.heading, theme, size=theme.size_heading,
                            bold=True, color="accent", space_after=0, first=True)
            add_rect(slide, x, cursor + 0.46, width, 0.014, fill="divider", theme=theme)
            cursor += 0.66
        remaining = max(0.5, height - (cursor - y))
        if column.text:
            size = fit_font_size([(column.text, 0)], width, remaining,
                                 base_pt=theme.size_body, min_pt=11)
            box = add_textbox(slide, x, cursor, width, remaining)
            write_paragraph(box.text_frame, column.text, theme, size=size, first=True)
            self._autofit(box.text_frame, [(column.text, 0)], size, width, remaining)
        if column.bullets:
            self._write_bullets(slide, column.bullets, cursor, x, width, remaining)

    def _render_comparison(self, slide, spec: M.ComparisonSlide) -> None:
        top = self._title_block(slide, spec.title)
        geo, theme = self.geo, self.theme
        height = self._content_height(top)
        tones = {"neutral": "accent", "positive": "positive", "negative": "negative"}
        for index, (column, tone) in enumerate(
                ((spec.left, spec.left_tone), (spec.right, spec.right_tone))):
            x = geo.margin_x + index * (geo.column_w + geo.gutter)
            add_rect(slide, x, top, geo.column_w, height, fill="surface", theme=theme)
            add_rect(slide, x, top, geo.column_w, 0.06, fill=tones[tone], theme=theme)
            inner_x, inner_w = x + 0.28, geo.column_w - 0.56
            cursor = top + 0.3
            if column.heading:
                box = add_textbox(slide, inner_x, cursor, inner_w, 0.44)
                write_paragraph(box.text_frame, column.heading, theme,
                                size=theme.size_heading, bold=True, color=tones[tone],
                                space_after=0, first=True)
                cursor += 0.58
            remaining = height - (cursor - top) - 0.3
            if column.text:
                box = add_textbox(slide, inner_x, cursor, inner_w, remaining)
                write_paragraph(box.text_frame, column.text, theme, size=theme.size_body, first=True)
            if column.bullets:
                self._write_bullets(slide, column.bullets, cursor, inner_x, inner_w, remaining)

    def _render_kpi(self, slide, spec: M.KpiSlide) -> None:
        top = self._title_block(slide, spec.title, spec.subtitle)
        geo, theme = self.geo, self.theme
        items = spec.items
        count = len(items)
        gap = 0.3
        card_w = (geo.content_w - gap * (count - 1)) / count
        card_h = min(2.6, self._content_height(top))
        tones = {"neutral": "muted", "positive": "positive", "negative": "negative"}

        inner_w = card_w - 0.5
        for i, item in enumerate(items):
            x = geo.margin_x + i * (card_w + gap)
            add_rect(slide, x, top, card_w, card_h, fill="surface", theme=theme)
            add_rect(slide, x, top, card_w, 0.06, fill="accent", theme=theme)

            # 값 → 라벨 → 증감 순으로 겹치지 않게 세로로 쌓는다.
            delta_h = 0.34 if item.delta else 0.0
            value_y = top + 0.34
            value_h = min(card_h * 0.40, 1.3)
            label_y = value_y + value_h
            label_h = max(0.32, top + card_h - 0.26 - delta_h - label_y)

            value_size = fit_single_line(item.value, inner_w,
                                         base_pt=min(theme.size_kpi_value, int(value_h * 72 * 0.72)),
                                         min_pt=20)
            box = add_textbox(slide, x + 0.25, value_y, inner_w, value_h)
            write_paragraph(box.text_frame, item.value, theme, size=value_size, bold=True,
                            space_after=0, first=True)

            label_size = fit_font_size([(item.label, 0)], inner_w, label_h,
                                       base_pt=theme.size_kpi_label, min_pt=10)
            box = add_textbox(slide, x + 0.25, label_y, inner_w, label_h)
            write_paragraph(box.text_frame, item.label, theme, size=label_size,
                            color="muted", space_after=0, first=True)

            if item.delta:
                box = add_textbox(slide, x + 0.25, top + card_h - delta_h - 0.2, inner_w, delta_h)
                write_paragraph(box.text_frame, item.delta, theme, size=theme.size_small,
                                bold=True, color=tones[item.tone], space_after=0, first=True)

    def _render_timeline(self, slide, spec: M.TimelineSlide) -> None:
        top = self._title_block(slide, spec.title)
        geo, theme = self.geo, self.theme
        steps = spec.steps
        count = len(steps)
        height = self._content_height(top)
        axis_y = top + height * 0.34
        slot_w = geo.content_w / count

        add_rect(slide, geo.margin_x + slot_w * 0.5, axis_y, geo.content_w - slot_w,
                 0.018, fill="divider", theme=theme)
        for i, step in enumerate(steps):
            center = geo.margin_x + slot_w * (i + 0.5)
            box = add_textbox(slide, center - slot_w / 2 + 0.1, axis_y - 0.62, slot_w - 0.2, 0.4)
            write_paragraph(box.text_frame, step.label, theme, size=theme.size_small,
                            bold=True, color="accent", align="center", space_after=0, first=True)
            marker = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(center - 0.11), Inches(axis_y - 0.1),
                Inches(0.22), Inches(0.22))
            marker.fill.solid()
            marker.fill.fore_color.rgb = resolve_color(theme, "accent")
            marker.line.fill.background()
            marker.shadow.inherit = False

            box = add_textbox(slide, center - slot_w / 2 + 0.1, axis_y + 0.34, slot_w - 0.2, 0.62)
            write_paragraph(box.text_frame, step.title, theme, size=theme.size_small + 2,
                            bold=True, align="center", space_after=2, first=True)
            if step.description:
                box = add_textbox(slide, center - slot_w / 2 + 0.1, axis_y + 1.0,
                                  slot_w - 0.2, height - (axis_y + 1.0 - top))
                write_paragraph(box.text_frame, step.description, theme,
                                size=theme.size_caption, color="muted", align="center", first=True)

    def _render_quote(self, slide, spec: M.QuoteSlide) -> None:
        geo, theme = self.geo, self.theme
        if not self.uses_template:
            add_rect(slide, 0, 0, geo.slide_w, geo.slide_h, fill="surface", theme=theme)
        width = geo.content_w * 0.82
        left = (geo.slide_w - width) / 2
        mark_h, text_h, tail_h = 1.0, 2.5, 0.5
        top = (geo.slide_h - (mark_h + text_h + tail_h)) / 2

        box = add_textbox(slide, left, top, width, mark_h)
        write_paragraph(box.text_frame, "\u201c", theme, size=54, bold=True, color="accent",
                        align="center", space_after=0, first=True)

        size = fit_font_size([(spec.text, 0)], width, text_h, base_pt=30, min_pt=16,
                             line_spacing=1.32)
        box = add_textbox(slide, left, top + mark_h, width, text_h, anchor="top")
        write_paragraph(box.text_frame, spec.text, theme, size=size, align="center",
                        line_spacing=1.32, space_after=0, first=True)

        tail = " \u00b7 ".join(x for x in (spec.attribution, spec.role) if x)
        if tail:
            box = add_textbox(slide, left, top + mark_h + text_h, width, tail_h)
            write_paragraph(box.text_frame, f"\u2014 {tail}", theme, size=theme.size_small,
                            color="muted", align="center", space_after=0, first=True)

    def _render_image(self, slide, spec: M.ImageSlide) -> None:
        geo, theme = self.geo, self.theme
        stream = load_image(spec.image, self.settings)

        if spec.placement == "full_bleed":
            self._place_image(slide, stream, 0, 0, geo.slide_w, geo.slide_h, cover=True)
            if spec.title:
                add_rect(slide, 0, geo.slide_h - 1.25, geo.slide_w, 1.25,
                         fill="background", theme=theme)
                box = add_textbox(slide, geo.margin_x, geo.slide_h - 1.05, geo.content_w, 0.6)
                write_paragraph(box.text_frame, spec.title, theme,
                                size=theme.size_slide_title - 4, bold=True, space_after=0, first=True)
            return

        top = self._title_block(slide, spec.title)
        height = self._content_height(top)
        if spec.caption:
            height -= 0.42

        if spec.placement == "full" or not spec.bullets:
            self._place_image(slide, stream, geo.margin_x, top, geo.content_w, height)
            caption_x, caption_w = geo.margin_x, geo.content_w
        else:
            image_first = spec.placement == "left"
            image_x = geo.margin_x if image_first else geo.margin_x + geo.column_w + geo.gutter
            text_x = geo.margin_x + geo.column_w + geo.gutter if image_first else geo.margin_x
            self._place_image(slide, stream, image_x, top, geo.column_w, height)
            self._write_bullets(slide, spec.bullets, top, text_x, geo.column_w, height)
            caption_x, caption_w = image_x, geo.column_w

        if spec.caption:
            box = add_textbox(slide, caption_x, top + height + 0.1, caption_w, 0.36)
            write_paragraph(box.text_frame, spec.caption, theme, size=theme.size_caption,
                            color="muted", space_after=0, first=True)

    def _place_image(self, slide, stream, x: float, y: float, w: float, h: float,
                     *, cover: bool = False) -> None:
        """상자 안에 비율을 유지하며 배치. cover=True면 상자를 꽉 채우고 넘치는 쪽을 자른다."""
        picture = slide.shapes.add_picture(stream, Inches(x), Inches(y))
        native_w, native_h = picture.width / 914400, picture.height / 914400
        scale = (max if cover else min)(w / native_w, h / native_h)
        draw_w, draw_h = native_w * scale, native_h * scale
        picture.width, picture.height = Inches(draw_w), Inches(draw_h)
        picture.left = Inches(x + (w - draw_w) / 2)
        picture.top = Inches(y + (h - draw_h) / 2)
        if cover:  # 상자 밖으로 나간 부분을 잘라 낸다
            picture.crop_left = picture.crop_right = max(0.0, (draw_w - w) / draw_w / 2)
            picture.crop_top = picture.crop_bottom = max(0.0, (draw_h - h) / draw_h / 2)
            picture.width, picture.height = Inches(w), Inches(h)
            picture.left, picture.top = Inches(x), Inches(y)

    def _render_table(self, slide, spec: M.TableSlide) -> None:
        top = self._title_block(slide, spec.title)
        geo, theme = self.geo, self.theme
        height = self._content_height(top) - (0.42 if spec.caption else 0)
        rows, cols = len(spec.rows) + 1, len(spec.headers)
        row_h = min(0.52, height / rows)
        table_h = row_h * rows

        shape = slide.shapes.add_table(rows, cols, Inches(geo.margin_x), Inches(top),
                                       Inches(geo.content_w), Inches(table_h))
        table = shape.table
        table.first_row = True

        if spec.column_widths:
            total = sum(spec.column_widths)
            for i, ratio in enumerate(spec.column_widths):
                table.columns[i].width = Emu(int(geo.content_w * 914400 * ratio / total))

        longest = max((len(str(c)) for row in spec.rows for c in row), default=8)
        size = max(9, min(theme.size_small, int(theme.size_small - max(0, cols - 4) - longest // 14)))

        for c, header in enumerate(spec.headers):
            self._style_cell(table.cell(0, c), str(header), size=size, bold=True,
                             fill=theme.accent, color="#FFFFFF")
        for r, row in enumerate(spec.rows, start=1):
            band = theme.surface if r % 2 == 0 else theme.background
            for c, value in enumerate(row):
                self._style_cell(table.cell(r, c), str(value), size=size, fill=band,
                                 color=theme.text)

        if spec.caption:
            box = add_textbox(slide, geo.margin_x, top + table_h + 0.12, geo.content_w, 0.36)
            write_paragraph(box.text_frame, spec.caption, theme, size=theme.size_caption,
                            color="muted", space_after=0, first=True)

    def _style_cell(self, cell, text: str, *, size: int, fill: str, color: str,
                    bold: bool = False) -> None:
        cell.text = ""
        cell.fill.solid()
        cell.fill.fore_color.rgb = hex_to_rgb(fill)
        cell.margin_left = cell.margin_right = Inches(0.1)
        cell.margin_top = cell.margin_bottom = Inches(0.05)
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.LEFT
        run = para.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = hex_to_rgb(color)
        from .render import apply_typeface
        apply_typeface(run, self.theme)

    def _render_chart(self, slide, spec: M.ChartSlide) -> None:
        top = self._title_block(slide, spec.title)
        geo, theme = self.geo, self.theme
        height = self._content_height(top) - (0.42 if spec.caption else 0)

        if spec.chart_type == "scatter":
            graphic = self._add_scatter(slide, spec, top, height)
        else:
            data = CategoryChartData()
            data.categories = spec.categories
            for series in spec.series:
                data.add_series(series.name, series.values)
            graphic = slide.shapes.add_chart(
                _CHART_TYPES[spec.chart_type], Inches(geo.margin_x), Inches(top),
                Inches(geo.content_w), Inches(height), data)

        chart = graphic.chart
        chart.has_title = False
        chart.font.size = Pt(theme.size_caption)
        chart.font.name = theme.font_latin
        chart.font.color.rgb = theme.rgb_muted

        if spec.show_legend and (len(spec.series) > 1 or spec.chart_type in {"pie", "doughnut"}):
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
        else:
            chart.has_legend = False

        self._color_chart(chart, spec)
        if spec.show_data_labels:
            plot = chart.plots[0]
            plot.has_data_labels = True
            labels = plot.data_labels
            labels.font.size = Pt(theme.size_caption - 1)
            labels.font.color.rgb = theme.rgb_muted
            if spec.chart_type in {"column", "bar"}:
                labels.position = XL_LABEL_POSITION.OUTSIDE_END

        if spec.caption:
            box = add_textbox(slide, geo.margin_x, top + height + 0.12, geo.content_w, 0.36)
            write_paragraph(box.text_frame, spec.caption, theme, size=theme.size_caption,
                            color="muted", space_after=0, first=True)

    def _add_scatter(self, slide, spec: M.ChartSlide, top: float, height: float):
        try:
            xs = [float(c) for c in spec.categories]
        except ValueError as exc:
            raise ValueError(
                "scatter 차트의 categories는 숫자여야 합니다. 항목 이름이라면 line이나 column을 쓰세요."
            ) from exc
        data = XyChartData()
        for series in spec.series:
            entry = data.add_series(series.name)
            for x, y in zip(xs, series.values):
                entry.add_data_point(x, y)
        geo = self.geo
        return slide.shapes.add_chart(
            XL_CHART_TYPE.XY_SCATTER_LINES, Inches(geo.margin_x), Inches(top),
            Inches(geo.content_w), Inches(height), data)

    def _color_chart(self, chart, spec: M.ChartSlide) -> None:
        theme = self.theme
        palette = [theme.accent, theme.accent_alt, theme.positive, theme.negative,
                   theme.text_muted, theme.divider]
        if spec.chart_type in {"pie", "doughnut"}:
            points = chart.plots[0].series[0].points
            for i, point in enumerate(points):
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = hex_to_rgb(palette[i % len(palette)])
            return
        for i, series in enumerate(chart.series):
            color = hex_to_rgb(palette[i % len(palette)])
            if spec.chart_type in {"line"}:
                series.format.line.color.rgb = color
                series.format.line.width = Pt(2.25)
            else:
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = color

    def _render_blank(self, slide, spec: M.BlankSlide) -> None:
        theme = self.theme
        if spec.title:
            self._title_block(slide, spec.title)
        for block in spec.blocks:
            box = add_textbox(slide, block.x, block.y, block.w, block.h)
            size = block.size or theme.size_body
            write_paragraph(box.text_frame, block.text, theme, size=size, bold=block.bold,
                            color=block.color, align=block.align, first=True)
            self._autofit(box.text_frame, [(block.text, 0)], size, block.w, block.h)


__all__ = ["DeckBuilder"]
