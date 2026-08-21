"""ppt-mcp 동작 검증."""

from __future__ import annotations

import asyncio
import collections
import json
import zipfile
from pathlib import Path

import pytest
from pptx import Presentation

from layout_audit import audit
from ppt_mcp.builder import DeckBuilder
from ppt_mcp.config import PathNotAllowed, resolve_read_path, resolve_write_path, DECK_SUFFIXES
from ppt_mcp.inspection import inspect_deck
from ppt_mcp.models import (BulletsSlide, DeckSpec, KpiItem, KpiSlide, QuoteSlide,
                            TableSlide, TitleSlide)
from ppt_mcp.render import normalize_bullets
from ppt_mcp.server import build_server
from ppt_mcp.textfit import display_width, fit_font_size
from ppt_mcp.theme import get_theme

EXAMPLE = Path(__file__).resolve().parents[1] / "examples" / "weekly_report.json"
NS_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def build(spec_dict: dict, settings) -> Presentation:
    return DeckBuilder(DeckSpec.model_validate(spec_dict), settings).build()


# --- 경로 안전장치 ---------------------------------------------------------

def test_상대경로는_작업디렉터리_안으로_해석된다(settings):
    path = resolve_write_path("보고서.pptx", settings)
    assert path.parent == settings.output_dir


def test_작업디렉터리_밖_저장은_거부된다(settings):
    with pytest.raises(PathNotAllowed):
        resolve_write_path("/etc/evil.pptx", settings)
    with pytest.raises(PathNotAllowed):
        resolve_write_path("../../탈출.pptx", settings)


def test_확장자_제한이_적용된다(settings):
    with pytest.raises(PathNotAllowed):
        resolve_write_path("보고서.exe", settings)
    with pytest.raises(PathNotAllowed):
        resolve_read_path("/etc/passwd", settings, suffixes=DECK_SUFFIXES)


def test_확장자_없으면_pptx가_붙는다(settings):
    assert resolve_write_path("보고서", settings).suffix == ".pptx"


# --- 불릿·텍스트 처리 ------------------------------------------------------

def test_불릿_문자열의_들여쓰기_기호를_단계로_읽는다():
    got = normalize_bullets(["최상위", "- 1단계", "-- 2단계", "  들여쓴 것", "• 점"])
    assert [(b.text, b.level) for b in got] == [
        ("최상위", 0), ("1단계", 1), ("2단계", 2), ("들여쓴 것", 1), ("점", 1)]


def test_한글은_전각으로_폭을_센다():
    assert display_width("안녕") == 4.0
    assert display_width("abcd") == 4.0


def test_내용이_많으면_글자크기를_줄인다():
    box = (11.5, 5.0)
    적음 = fit_font_size([("짧은 항목", 0)] * 3, *box, base_pt=18, min_pt=11)
    많음 = fit_font_size([("아주 긴 항목입니다 " * 6, 0)] * 10, *box, base_pt=18, min_pt=11)
    assert 적음 == 18
    assert 많음 < 적음 >= 11


# --- 덱 생성 ---------------------------------------------------------------

def test_예제_스펙이_그대로_생성된다(settings):
    spec = json.loads(EXAMPLE.read_text())
    prs = build(spec, settings)
    assert len(prs.slides) == len(spec["slides"])


def test_모든_슬라이드_종류가_생성된다(settings, tiny_png):
    from ppt_mcp.models import SLIDE_TYPES
    slides = [
        {"type": "title", "title": "표지"},
        {"type": "agenda", "items": ["가", "나", "다"]},
        {"type": "section", "title": "간지", "number": "01"},
        {"type": "bullets", "title": "불릿", "bullets": ["가", "- 나"]},
        {"type": "two_column", "title": "2단",
         "left": {"heading": "왼쪽", "bullets": ["가"]},
         "right": {"heading": "오른쪽", "text": "문단"}},
        {"type": "comparison", "title": "비교",
         "left": {"heading": "As-Is", "bullets": ["가"]},
         "right": {"heading": "To-Be", "bullets": ["나"]}},
        {"type": "table", "title": "표", "headers": ["A", "B"], "rows": [["1", "2"]]},
        {"type": "chart", "title": "차트", "categories": ["1월", "2월"],
         "series": [{"name": "값", "values": [1, 2]}]},
        {"type": "image", "title": "이미지", "image": {"path": str(tiny_png)}},
        {"type": "quote", "text": "인용", "attribution": "누군가"},
        {"type": "kpi", "title": "지표", "items": [{"value": "42%", "label": "비율"}]},
        {"type": "timeline", "title": "일정",
         "steps": [{"label": "1분기", "title": "준비"}, {"label": "2분기", "title": "실행"}]},
        {"type": "blank", "title": "자유",
         "blocks": [{"text": "메모", "x": 1, "y": 2, "w": 4, "h": 1}]},
    ]
    assert {s["type"] for s in slides} == set(SLIDE_TYPES), "슬라이드 종류가 빠졌다"
    prs = build({"title": "전종류", "slides": slides}, settings)
    assert len(prs.slides) == len(SLIDE_TYPES)


@pytest.mark.parametrize("theme", ["carbon_light", "carbon_dark", "minimal", "vivid"])
def test_테마별_레이아웃_불변식(settings, theme):
    spec = json.loads(EXAMPLE.read_text())
    spec["theme"] = theme
    assert audit(build(spec, settings)) == []


def test_긴_내용도_넘치거나_겹치지_않는다(settings):
    긴줄 = "아주 길고 장황해서 한 줄에 절대 들어가지 않는 항목입니다 " * 3
    prs = build({"title": "과다", "slides": [
        {"type": "bullets", "title": "매우 긴 제목을 넣어 제목 줄바꿈까지 확인하는 슬라이드입니다",
         "subtitle": "부제도 꽤 길게 붙여 봅니다", "bullets": [긴줄] * 12},
        {"type": "agenda", "items": [f"{i}. 아주 긴 목차 항목" for i in range(9)]},
        {"type": "comparison", "title": "비교",
         "left": {"heading": "왼쪽", "bullets": [긴줄] * 4},
         "right": {"heading": "오른쪽", "bullets": [긴줄] * 4}},
        {"type": "kpi", "title": "지표", "items": [
            {"value": "1,234,567건", "label": "아주 긴 지표 이름입니다", "delta": "+99%p",
             "tone": "positive"}] * 4},
        {"type": "quote", "text": 긴줄 * 2, "attribution": "긴 인용"},
    ]}, settings)
    assert audit(prs) == []


@pytest.mark.parametrize("aspect,expected_w", [("16:9", 13.333), ("4:3", 10.0)])
def test_슬라이드_비율(settings, aspect, expected_w):
    prs = build({"title": "비율", "aspect": aspect,
                 "slides": [{"type": "title", "title": "표지"}]}, settings)
    assert prs.slide_width / 914400 == pytest.approx(expected_w, abs=0.01)
    assert audit(prs) == []


# --- 한글 폰트 -------------------------------------------------------------

def test_한글_폰트가_ea까지_지정된다(settings, tmp_path):
    theme = get_theme("carbon_light")
    prs = build({"title": "폰트", "slides": [
        {"type": "bullets", "title": "한글 제목", "bullets": ["한글 본문"]}]}, settings)
    path = tmp_path / "font.pptx"
    prs.save(str(path))

    runs = [run for slide in Presentation(str(path)).slides for shape in slide.shapes
            if shape.has_text_frame for para in shape.text_frame.paragraphs for run in para.runs]
    assert runs
    for run in runs:
        rPr = run._r.find(f"{NS_A}rPr")
        assert rPr.find(f"{NS_A}ea").get("typeface") == theme.font_ea
        assert rPr.find(f"{NS_A}latin").get("typeface") == theme.font_latin
        # OOXML 스키마상 latin → ea → cs 순서를 지켜야 PowerPoint가 파일을 연다
        order = [child.tag.split("}")[1] for child in rPr]
        assert order.index("latin") < order.index("ea") < order.index("cs")


# --- 슬라이드 단위 편집 ----------------------------------------------------

def test_삭제후_추가해도_파트이름이_겹치지_않는다(settings, tmp_path):
    """python-pptx는 새 파트 이름을 슬라이드 수 + 1로 정해서, 중간을 지우면 충돌한다."""
    builder = DeckBuilder(DeckSpec(title="편집", slides=[]), settings)
    builder.add_slide(TitleSlide(title="A"))
    builder.add_slide(BulletsSlide(title="B", bullets=["x"]))
    builder.add_slide(KpiSlide(title="C", items=[KpiItem(value="1", label="l")]))
    builder.delete_slide(1)
    builder.add_slide(QuoteSlide(text="Q"), index=1)
    builder.add_slide(BulletsSlide(title="D", bullets=["y"]))
    builder.delete_slide(0)
    builder.move_slide(0, 2)

    path = tmp_path / "edit.pptx"
    builder.prs.save(str(path))

    names = zipfile.ZipFile(path).namelist()
    duplicates = [n for n, count in collections.Counter(names).items() if count > 1]
    assert duplicates == []
    assert len(Presentation(str(path)).slides) == 3


def test_슬라이드_이동_순서(settings):
    builder = DeckBuilder(DeckSpec(title="순서", slides=[]), settings)
    for name in "ABC":
        builder.add_slide(BulletsSlide(title=name, bullets=["x"]))
    builder.move_slide(0, 2)
    titles = [inspect_title(slide) for slide in builder.prs.slides]
    assert titles == ["B", "C", "A"]


def inspect_title(slide) -> str:
    return slide.shapes[0].text_frame.text


def test_없는_슬라이드_삭제는_오류(settings):
    builder = DeckBuilder(DeckSpec(title="x", slides=[]), settings)
    with pytest.raises(IndexError):
        builder.delete_slide(0)


# --- 템플릿 ---------------------------------------------------------------

def test_템플릿의_레이아웃을_쓰고_원본_슬라이드는_비운다(settings, tmp_path):
    seed = Presentation()
    seed.slides.add_slide(seed.slide_layouts[0]).shapes.title.text = "원래 있던 슬라이드"
    template = tmp_path / "corp.pptx"
    seed.save(str(template))

    prs = build({"title": "템플릿", "template": str(template), "slides": [
        {"type": "title", "title": "표지"},
        {"type": "bullets", "title": "본문", "bullets": ["가"]},
        {"type": "section", "title": "간지"},
    ]}, settings)

    assert len(prs.slides) == 3
    assert [s.slide_layout.name for s in prs.slides] == [
        "Title Slide", "Title and Content", "Section Header"]
    all_text = " ".join(sh.text_frame.text for s in prs.slides for sh in s.shapes
                        if sh.has_text_frame)
    assert "원래 있던 슬라이드" not in all_text


def test_빈_플레이스홀더는_남지_않는다(settings, tmp_path):
    seed = Presentation()
    template = tmp_path / "corp.pptx"
    seed.save(str(template))
    prs = build({"title": "t", "template": str(template),
                 "slides": [{"type": "bullets", "title": "제목", "bullets": ["가"]}]}, settings)
    for shape in prs.slides[0].shapes:
        if shape.has_text_frame:
            assert shape.text_frame.text.strip(), "빈 플레이스홀더가 남았다"


# --- 검증 ------------------------------------------------------------------

def test_표의_행_길이가_다르면_거부된다():
    with pytest.raises(ValueError):
        TableSlide(title="t", headers=["A", "B"], rows=[["1"]])


def test_차트_계열_길이가_다르면_거부된다():
    with pytest.raises(ValueError):
        DeckSpec.model_validate({"title": "t", "slides": [
            {"type": "chart", "title": "c", "categories": ["가", "나"],
             "series": [{"name": "s", "values": [1]}]}]})


def test_scatter는_숫자축을_요구한다(settings):
    with pytest.raises(ValueError, match="숫자"):
        build({"title": "t", "slides": [
            {"type": "chart", "title": "c", "chart_type": "scatter",
             "categories": ["가", "나"], "series": [{"name": "s", "values": [1, 2]}]}]}, settings)


def test_없는_테마는_거부된다(settings):
    with pytest.raises(ValueError, match="테마"):
        build({"title": "t", "theme": "없는테마", "slides": []}, settings)


def test_이미지_원격주소는_기본_차단(settings):
    with pytest.raises(PermissionError):
        build({"title": "t", "slides": [
            {"type": "image", "image": {"url": "https://example.com/a.png"}}]}, settings)


def test_슬라이드_수_상한(settings):
    from dataclasses import replace
    limited = replace(settings, max_slides=2)
    with pytest.raises(ValueError, match="너무 많"):
        build({"title": "t", "slides": [{"type": "section", "title": str(i)} for i in range(3)]},
              limited)


# --- 읽기 ------------------------------------------------------------------

def test_생성한_덱을_다시_읽으면_제목과_노트가_보인다(settings):
    spec = json.loads(EXAMPLE.read_text())
    prs = build(spec, settings)
    settings.ensure_output_dir()
    path = settings.output_dir / "r.pptx"
    prs.save(str(path))

    report = inspect_deck(str(path), settings)
    assert report.slide_count == len(spec["slides"])
    assert report.aspect == "16:9"
    assert report.author == spec["author"]
    assert report.slides[0].title == spec["slides"][0]["title"]
    assert report.slides[0].notes == spec["slides"][0]["notes"]
    assert any(s.charts for s in report.slides)
    assert any(s.tables for s in report.slides)


def test_include_text_false면_본문을_생략한다(settings):
    prs = build({"title": "t", "slides": [
        {"type": "bullets", "title": "제목", "bullets": ["본문"]}]}, settings)
    settings.ensure_output_dir()
    path = settings.output_dir / "r.pptx"
    prs.save(str(path))
    assert inspect_deck(str(path), settings, include_text=False).slides[0].texts == []


# --- MCP 도구 --------------------------------------------------------------

def call(mcp, name, args=None):
    return asyncio.run(mcp.call_tool(name, args or {}))


def test_도구가_모두_등록된다(settings):
    mcp = build_server(settings)
    names = {t.name for t in asyncio.run(mcp.list_tools())}
    assert names == {
        "describe_options", "read_deck", "list_open_decks", "create_deck", "open_deck",
        "add_slide", "update_slide", "delete_slide", "move_slide", "save_deck", "close_deck"}


def test_create_deck_이_파일을_만든다(settings):
    mcp = build_server(settings)
    spec = json.loads(EXAMPLE.read_text())
    result = call(mcp, "create_deck", {"spec": spec, "output_path": "주간.pptx"}).structured_content
    assert Path(result["path"]).exists()
    assert result["slide_count"] == len(spec["slides"])
    assert result["used_template"] is False


def test_편집_워크플로(settings):
    mcp = build_server(settings)
    opened = call(mcp, "open_deck", {"spec": {"title": "편집", "slides": [
        {"type": "title", "title": "표지"},
        {"type": "bullets", "title": "지울 것", "bullets": ["x"]}]}}).structured_content
    deck_id = opened["deck_id"]
    assert opened["slide_count"] == 2

    call(mcp, "add_slide", {"deck_id": deck_id,
                            "slide": {"type": "quote", "text": "인용"}})
    call(mcp, "update_slide", {"deck_id": deck_id, "index": 1,
                               "slide": {"type": "section", "title": "간지"}})
    call(mcp, "move_slide", {"deck_id": deck_id, "from_index": 2, "to_index": 0})
    after_delete = call(mcp, "delete_slide", {"deck_id": deck_id, "index": 0}).structured_content
    assert after_delete["slide_count"] == 2

    saved = call(mcp, "save_deck", {"deck_id": deck_id}).structured_content
    assert Path(saved["path"]).exists()
    assert audit(Presentation(saved["path"])) == []

    call(mcp, "close_deck", {"deck_id": deck_id})
    assert call(mcp, "list_open_decks").structured_content["decks"] == []


def test_없는_deck_id는_오류(settings):
    mcp = build_server(settings)
    with pytest.raises(Exception, match="열려 있는 덱이 아닙니다"):
        call(mcp, "save_deck", {"deck_id": "deck_없음"})


def test_describe_options가_선택지를_알려준다(settings):
    mcp = build_server(settings)
    options = call(mcp, "describe_options").structured_content
    assert {t["name"] for t in options["themes"]} == {
        "carbon_light", "carbon_dark", "minimal", "vivid"}
    assert len(options["slide_types"]) == 13
    assert options["remote_images_allowed"] is False
    assert options["template_layouts"] is None


def test_describe_options에_템플릿_레이아웃(settings, tmp_path):
    template = tmp_path / "corp.pptx"
    Presentation().save(str(template))
    mcp = build_server(settings)
    layouts = call(mcp, "describe_options",
                   {"template": str(template)}).structured_content["template_layouts"]
    assert any(l["name"] == "Title and Content" for l in layouts)


def test_내용이_넘치면_경고를_돌려준다(settings):
    mcp = build_server(settings)
    긴줄 = "아주 길고 장황해서 한 줄에 절대 들어가지 않는 항목입니다 " * 3
    result = call(mcp, "create_deck", {"spec": {"title": "과다", "slides": [
        {"type": "bullets", "title": "과밀", "bullets": [긴줄] * 12}]}}).structured_content
    assert result["warnings"], "글자를 줄였는데 경고가 없다"
    assert "나누" in result["warnings"][0]
    assert audit(Presentation(result["path"])) == []


def test_내용이_적으면_경고가_없다(settings):
    mcp = build_server(settings)
    result = call(mcp, "create_deck", {"spec": {"title": "보통", "slides": [
        {"type": "bullets", "title": "적당", "bullets": ["가", "나", "다"]}]}}).structured_content
    assert result["warnings"] == []
